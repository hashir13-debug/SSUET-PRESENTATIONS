"""Animation_Group3_PREMIUM.pptx — dark theme, animations, premium visuals."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ── NAMESPACE ─────────────────────────────────────────────────
PNS = "http://schemas.openxmlformats.org/presentationml/2006/main"
ANS = "http://schemas.openxmlformats.org/drawingml/2006/main"

def pq(tag): return f"{{{PNS}}}{tag}"
def aq(tag): return f"{{{ANS}}}{tag}"

# ── COLORS ────────────────────────────────────────────────────
BG    = RGBColor(6,   6,  24)
CARD  = RGBColor(14, 14,  50)
CARDA = RGBColor(18, 18,  62)
HDR   = RGBColor(10, 10,  36)
C     = RGBColor(0,  210, 255)
P     = RGBColor(150, 80, 220)
O     = RGBColor(255,145,   0)
G     = RGBColor(46, 213, 115)
PK    = RGBColor(255, 90, 160)
Y     = RGBColor(255,215,   0)
W     = RGBColor(255,255, 255)
LT    = RGBColor(195,205, 232)
MU    = RGBColor(110,125, 170)

FH = "Trebuchet MS"
FB = "Calibri"

# ── SETUP ─────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ═══════════════════════════════════════════════════════════════
# SHAPE / TEXT HELPERS
# ═══════════════════════════════════════════════════════════════
def new_slide():
    sl = prs.slides.add_slide(BLANK)
    bg = sl.background.fill
    bg.solid(); bg.fore_color.rgb = BG
    return sl

def rct(sl, l, t, w, h, fill, lc=None, lw=0.5, alpha=None):
    sh = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if lc:
        sh.line.color.rgb = lc; sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh

def oval(sl, l, t, w, h, fill, lc=None, lw=0.75):
    sh = sl.shapes.add_shape(9, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if lc: sh.line.color.rgb = lc; sh.line.width = Pt(lw)
    else:   sh.line.fill.background()
    return sh

def txb(sl, text, l, t, w, h, sz, col, bold=False,
        align=PP_ALIGN.LEFT, font=FB, italic=False):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    par = tf.paragraphs[0]; par.alignment = align
    run = par.add_run()
    run.text = text; run.font.name = font; run.font.size = Pt(sz)
    run.font.color.rgb = col; run.font.bold = bold; run.font.italic = italic
    return tb

def multiline_txb(sl, lines, l, t, w, h, sz, default_col,
                  accent=None, font=FB, leading=1):
    """Text box with smart per-line coloring."""
    ACCENT_KEYS = ("EXAMPLES:", "TOOLS:", "USED IN:", "BEST FOR:", "ERA:", "CAREER:",
                   "OUTPUT:", "DISCOVERED:", "HOW IT WORKS:", "TYPES:", "AT 24fps:",
                   "REAL EXAMPLES:", "STEP ", "NOTE:", "MODERN", "FAMOUS",
                   "KEY INSIGHT:", "CRITICAL RULE:", "WHAT YOU", "ANALOGY:",
                   "STRAIGHT AHEAD:", "POSE-TO-POSE:", "FOLLOW THROUGH:",
                   "OVERLAPPING ACTION:", "IN SOFTWARE:", "EXCEPTION:",
                   "HOW TO ACHIEVE:", "THE TEST:", "APPEAL", "2D:", "3D:",
                   "Motion:", "Games:", "CHARACTERISTICS:", "CHARACTERISTICS:")
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.space_before = Pt(leading)
        run = par.add_run()
        run.text = ln; run.font.name = font; run.font.size = Pt(sz)
        stripped = ln.strip()
        if not stripped:
            run.font.color.rgb = MU
        elif accent and any(stripped.startswith(k) for k in ACCENT_KEYS):
            run.font.color.rgb = accent; run.font.bold = True
        else:
            run.font.color.rgb = default_col
    return tb

# ═══════════════════════════════════════════════════════════════
# SLIDE STRUCTURE HELPERS
# ═══════════════════════════════════════════════════════════════
def hdr(sl, title, num="", accent=C):
    """Dark header bar with accent line and number."""
    rct(sl, 0, 0, 13.33, 0.85, HDR)
    rct(sl, 0, 0.82, 13.33, 0.03, accent)
    rct(sl, 0, 0, 0.06, 0.85, accent)
    # Subtle diagonal stripe decoration
    for i in range(6):
        rct(sl, 12.5 - i*0.28, 0, 0.14, 0.85, RGBColor(16,16,52))
    if num:
        txb(sl, num,   0.15, 0.08, 0.55, 0.70, 26, accent, bold=True, font=FH)
        txb(sl, title, 0.77, 0.13, 11.3, 0.62, 19, W,      bold=True, font=FH)
        # Big watermark
        txb(sl, num, 11.2, -0.25, 2.2, 1.2, 82, accent, bold=True, font=FH)
    else:
        txb(sl, title, 0.25, 0.13, 12.7, 0.62, 19, W, bold=True, font=FH)

def info_bar(sl, text, t=0.90, accent=C):
    rct(sl, 0.28, t, 12.77, 0.40, RGBColor(10,10,38))
    rct(sl, 0.28, t, 0.055, 0.40, accent)
    txb(sl, text, 0.40, t+0.05, 12.55, 0.32, 9, LT, italic=True)

def card(sl, l, t, w, h, accent, title, lines, title_col=None):
    """Premium card with double border + accent bar."""
    tc = title_col or accent
    # Shadow/glow layer
    rct(sl, l+0.04, t+0.04, w, h, RGBColor(4,4,18))
    # Main card
    rct(sl, l, t, w, h, CARD, lc=RGBColor(30,30,80), lw=0.4)
    # Accent bar
    rct(sl, l, t, 0.06, h, accent)
    # Inner highlight
    rct(sl, l+0.06, t, w-0.06, 0.003, accent)
    # Title
    txb(sl, title, l+0.12, t+0.1, w-0.18, 0.28, 10.5, tc, bold=True, font=FH)
    # Divider
    rct(sl, l+0.12, t+0.37, w-0.24, 0.012, RGBColor(40,40,90))
    # Body
    multiline_txb(sl, lines, l+0.12, t+0.41, w-0.18, h-0.52, 9, LT, accent=accent)
    return sl.shapes[-1]

def bg_dots(sl, color, count=40, seed=7):
    """Decorative background dots (simulate particles)."""
    import random; r = random.Random(seed)
    for _ in range(count):
        x = r.uniform(0.2, 13.0); y = r.uniform(0.9, 7.3)
        s = r.uniform(0.04, 0.1)
        sh = sl.shapes.add_shape(9, Inches(x), Inches(y), Inches(s), Inches(s))
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
        sh.line.fill.background()

# ═══════════════════════════════════════════════════════════════
# ANIMATION HELPERS
# ═══════════════════════════════════════════════════════════════
def add_fade_transition(sl, dur_ms=600):
    """Add fade slide transition."""
    tr = etree.SubElement(sl._element, pq("transition"))
    tr.set("spd", "med"); tr.set("dur", str(dur_ms))
    etree.SubElement(tr, pq("fade"))

def add_push_transition(sl, direction="l", dur_ms=500):
    """Add push slide transition (slides push in from right)."""
    tr = etree.SubElement(sl._element, pq("transition"))
    tr.set("spd", "med"); tr.set("dur", str(dur_ms))
    push = etree.SubElement(tr, pq("push"))
    push.set("dir", direction)

def animate_shapes(sl, shape_ids, start_ms=250, step_ms=180,
                   dur_ms=550, effect="fade"):
    """
    Auto-play sequential entrance animations.
    effect: "fade" | "float" (fade + upward motion)
    """
    if not shape_ids:
        return

    timing = etree.SubElement(sl._element, pq("timing"))
    tnLst  = etree.SubElement(timing, pq("tnLst"))
    par0   = etree.SubElement(tnLst, pq("par"))
    root   = etree.SubElement(par0, pq("cTn"))
    root.set("id", "1"); root.set("dur", "indefinite")
    root.set("restart", "whenNotActive"); root.set("nodeType", "tmRoot")
    ch = etree.SubElement(root, pq("childTnLst"))

    cid = 2
    for i, spid in enumerate(shape_ids):
        delay = start_ms + i * step_ms

        # Outer par — time delay
        op = etree.SubElement(ch, pq("par"))
        oc = etree.SubElement(op, pq("cTn"))
        oc.set("id", str(cid)); cid += 1; oc.set("fill", "hold")
        st = etree.SubElement(oc, pq("stCondLst"))
        co = etree.SubElement(st, pq("cond")); co.set("delay", str(delay))
        ich = etree.SubElement(oc, pq("childTnLst"))

        # Inner par — preset effect (presetID=10 = Fade)
        ip = etree.SubElement(ich, pq("par"))
        ic = etree.SubElement(ip, pq("cTn"))
        ic.set("id", str(cid)); cid += 1
        ic.set("presetID", "10"); ic.set("presetClass", "entr")
        ic.set("presetSubtype", "0"); ic.set("fill", "hold")
        ic.set("grpId", str(i)); ic.set("nodeType", "withEffect")
        ist = etree.SubElement(ic, pq("stCondLst"))
        ico = etree.SubElement(ist, pq("cond")); ico.set("delay", "0")
        ich2 = etree.SubElement(ic, pq("childTnLst"))

        # Fade effect
        ae = etree.SubElement(ich2, pq("animEffect"))
        ae.set("transition", "in"); ae.set("filter", "fade")
        ab = etree.SubElement(ae, pq("cBhvr"))
        ac = etree.SubElement(ab, pq("cTn"))
        ac.set("id", str(cid)); cid += 1; ac.set("dur", str(dur_ms))
        at = etree.SubElement(ab, pq("tgtEl"))
        asp = etree.SubElement(at, pq("spTgt")); asp.set("spid", str(spid))

        if effect == "float":
            # Motion: move from +3% Y to 0 (float up)
            am = etree.SubElement(ich2, pq("animMotion"))
            am.set("origin", "parent")
            am.set("path", "M 0 0.03 L 0 0 E")
            am.set("pathEditMode", "relative")
            amb = etree.SubElement(am, pq("cBhvr"))
            amc = etree.SubElement(amb, pq("cTn"))
            amc.set("id", str(cid)); cid += 1; amc.set("dur", str(dur_ms))
            amt = etree.SubElement(amb, pq("tgtEl"))
            amsp = etree.SubElement(amt, pq("spTgt")); amsp.set("spid", str(spid))

    etree.SubElement(timing, pq("bldLst"))

def animate_title_elements(sl, shape_ids):
    """Staggered entrance for title slide elements."""
    animate_shapes(sl, shape_ids, start_ms=100, step_ms=250, dur_ms=700, effect="float")

# ═══════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════

# ── S1 TITLE ──────────────────────────────────────────────────
s1 = new_slide()
bg_dots(s1, RGBColor(0, 40, 80),  count=30, seed=1)
bg_dots(s1, RGBColor(40, 0, 80),  count=20, seed=2)

# Big glow circles
for r, col in [(3.5, RGBColor(6,12,36)), (2.4, RGBColor(8,16,44)),
               (1.4, RGBColor(12,22,54)), (0.5, RGBColor(20,36,80))]:
    oval(s1, 12.83-r, 0.5-r, r*2, r*2, col, lc=RGBColor(0,50,90), lw=0.5)

# Film reel rings (decorative)
for r, col in [(0.75, C), (0.45, C), (0.15, C)]:
    oval(s1, 11.8-r, 0.2-r, r*2, r*2, RGBColor(8,20,50),
         lc=RGBColor(0,100,140), lw=1.0)

# Course label
e1 = txb(s1, "GRAPHICS & MULTIMEDIA  •  SMIT", 0, 1.05, 13.33, 0.45,
          13, O, bold=True, align=PP_ALIGN.CENTER, font=FH)

# ANIMATION title
e2 = txb(s1, "ANIMATION", 0, 1.5, 13.33, 2.3, 100, C,
          bold=True, align=PP_ALIGN.CENTER, font=FH)

# Subtitle
e3 = txb(s1, "The Art & Science of Bringing Images to Life",
          0, 3.6, 13.33, 0.55, 15, LT, align=PP_ALIGN.CENTER, italic=True)

# Divider line segments
e4 = rct(s1, 4.4, 4.3, 4.53, 0.025, C)

# Badge
rct(s1, 4.85, 4.48, 3.63, 0.42, RGBColor(8,22,50), lc=C, lw=0.75)
e5 = txb(s1, "Topics 1 — 9  •  Group 3", 4.85, 4.49, 3.63, 0.42, 12, C,
          bold=True, align=PP_ALIGN.CENTER, font=FH)

# Members
members = [
    ("Hashir Junaid",     "2023F-BCS-358"),
    ("Taha Haider",       "2023F-BCS-079"),
    ("Abdul Rahman Baig", "2023F-BCS-100"),
]
mem_shapes = []
for i, (name, roll) in enumerate(members):
    lx = 1.5 + i * 3.6
    rct(s1, lx, 5.06, 3.35, 0.82, RGBColor(16,16,58), lc=C, lw=0.5)
    m1 = txb(s1, name, lx+0.15, 5.11, 3.05, 0.35, 12, W, bold=True, font=FH)
    txb(s1, roll, lx+0.15, 5.46, 3.05, 0.3, 9.5, MU)
    mem_shapes.append(m1.shape_id)

# Animate title elements
animate_title_elements(s1, [e1.shape_id, e2.shape_id, e3.shape_id,
                             e4.shape_id, e5.shape_id] + mem_shapes)
add_fade_transition(s1, 800)

# ── S2 AGENDA ─────────────────────────────────────────────────
s2 = new_slide()
bg_dots(s2, RGBColor(0,40,70), count=25, seed=3)
hdr(s2, "WHAT WE'LL COVER", accent=C)

agenda = [
    (C,  "01", "Introduction to Animation",    "Definition • FPS • POV • Phi Phenomenon"),
    (P,  "02", "Usage of Animation",           "Art • Storytelling • Science • Education"),
    (O,  "03", "Animation Techniques",         "Cel • Path • Computer Animation"),
    (G,  "04", "Types of Animation",           "Keyframes • Tweening • Pencil Test"),
    (PK, "05", "Computer Animation",           "FK • IK • Morphing • Digital vs Traditional"),
    (Y,  "06", "Levels of Computer Animation", "Basic • Intermediate • Advanced"),
    (C,  "07", "12 Principles of Animation",   "Disney's universal gold standard"),
    (P,  "08", "Animation Space",              "2D  •  2½D  •  3D"),
    (O,  "09", "Animation Process",            "Organize → Tool → Build → Post"),
]
cols = [0.28, 4.53, 8.78]
rows = [0.93, 2.75, 4.57]
agenda_ids = []
for i, (acc, num, topic, sub) in enumerate(agenda):
    lx = cols[i % 3]; ty = rows[i // 3]
    rct(s2, lx+0.04, ty+0.04, 4.1, 1.72, RGBColor(4,4,18))
    sh = rct(s2, lx, ty, 4.1, 1.72, CARD, lc=RGBColor(28,28,72), lw=0.4)
    rct(s2, lx, ty, 0.06, 1.72, acc)
    rct(s2, lx+0.06, ty, 4.04, 0.003, acc)
    rct(s2, lx+0.12, ty+0.14, 0.65, 0.65, RGBColor(20,20,64))
    txb(s2, num,   lx+0.12, ty+0.17, 0.65, 0.62, 20, acc,
        bold=True, align=PP_ALIGN.CENTER, font=FH)
    txb(s2, topic, lx+0.87, ty+0.17, 3.1, 0.38, 11, LT, bold=True, font=FH)
    txb(s2, sub,   lx+0.87, ty+0.58, 3.1, 0.35, 8.5, MU)
    agenda_ids.append(sh.shape_id)

animate_shapes(s2, agenda_ids, start_ms=200, step_ms=150, effect="float")
add_push_transition(s2)

# ── S3 INTRODUCTION ───────────────────────────────────────────
s3 = new_slide()
bg_dots(s3, RGBColor(0,50,80), count=20, seed=4)
hdr(s3, "INTRODUCTION TO ANIMATION", "01", C)

# Filmstrip
rct(s3, 0.28, 0.92, 12.77, 0.60, RGBColor(8,8,30))
rct(s3, 0.28, 0.92, 0.055, 0.60, C)
for fi in range(4):
    fx = 1.0 + fi * 2.9
    rct(s3, fx, 0.96, 2.4, 0.52, RGBColor(16,16,54))
    # Moving ball (different Y each frame)
    ball_y = [1.2, 1.08, 0.99, 1.15]
    oval(s3, fx+0.95, ball_y[fi], 0.28, 0.28, C)
    txb(s3, f"F{fi+1}", fx+0.08, 0.98, 0.4, 0.18, 7, MU)
txb(s3, "24 FRAMES PER SECOND  —  ILLUSION OF MOTION",
    0.35, 0.94, 12.6, 0.22, 8, MU, italic=True, align=PP_ALIGN.CENTER)

info_bar(s3, '"Animation = sequential still images exploiting Persistence of Vision + Phi Phenomenon to create the ILLUSION OF MOVEMENT."', t=1.57)

grid_data = [
    (C, "DEFINITION",
     ["From Latin anima — 'to give soul'.",
      "• Series of still images shown rapidly",
      "• 24 fps = cinema standard",
      "• 30 fps = TV broadcast",
      "• 60 fps = modern gaming",
      "• Below 12 fps = visible flicker",
      "", "First animation: Eadweard Muybridge",
      "horse gallop sequence (1878)"]),
    (P, "PERSISTENCE OF VISION",
     ["Physical effect in the human eye.",
      "• Retina retains image ~50ms after it disappears",
      "• Afterimage overlaps the next frame",
      "• Creates seamless continuous flow",
      "", "DISCOVERED: Peter Mark Roget (1824)",
      "", "Result: brain cannot distinguish",
      "frames at 24+ fps"]),
    (O, "PHI PHENOMENON",
     ["Purely PSYCHOLOGICAL motion effect.",
      "• Brain infers motion between two static",
      "  objects in rapid sequence",
      "• Does NOT need retinal retention",
      "", "DISCOVERED: Max Wertheimer (1912)",
      "", "Examples: Flip-book, traffic light",
      "arrows, neon sign sequences"]),
    (G, "ILLUSION OF MOVEMENT",
     ["Combined result of BOTH effects.",
      "", "POV alone  ≠  enough",
      "Phi alone  ≠  enough",
      "Together  =  seamless motion",
      "", "The brain WANTS to see motion —",
      "animators exploit this biological",
      "'glitch' in human perception."]),
]
positions = [(0.28,2.03),(6.78,2.03),(0.28,4.65),(6.78,4.65)]
card_ids = []
for (acc, ttl, lns), (lx, ty) in zip(grid_data, positions):
    card(s3, lx, ty, 6.22, 2.58, acc, ttl, lns)
    card_ids.append(s3.shapes[-2].shape_id)

animate_shapes(s3, card_ids, start_ms=300, step_ms=180, effect="float")
add_push_transition(s3)

# ── S4 USAGE ──────────────────────────────────────────────────
s4 = new_slide()
bg_dots(s4, RGBColor(30,0,70), count=22, seed=5)
hdr(s4, "USAGE OF ANIMATION", "02", P)

usage_data = [
    (C,  "🎨", "ARTISTIC PURPOSES",
     ["• Fine art & experimental film",
      "• Music video production",
      "• Abstract visual compositions",
      "• Title sequence design",
      "• Gallery installations",
      "", "REAL EXAMPLES:",
      "Disney's Fantasia (1940)",
      "Yellow Submarine (1968)",
      "Spider-Man: Into the Spider-Verse"]),
    (O,  "📖", "STORYTELLING",
     ["• Narrative feature films",
      "• TV series & web animations",
      "• Video game cutscenes",
      "• Interactive web comics",
      "", "REAL EXAMPLES:",
      "Pixar Toy Story series",
      "Studio Ghibli classics",
      "Arcane — Netflix (2021)",
      "Attack on Titan"]),
    (G,  "🔬", "SCIENTIFIC VISUALIZATION",
     ["• Medical anatomy explainers",
      "• Weather & climate models",
      "• DNA & molecular modeling",
      "• Physics simulations",
      "• Engineering CAD previews",
      "", "REAL EXAMPLES:",
      "NASA mission visualizations",
      "Medical surgical training",
      "Climate change simulations"]),
    (PK, "🎓", "INSTRUCTIONAL PURPOSES",
     ["• Educational explainer videos",
      "• Military & flight simulators",
      "• Safety procedure demos",
      "• Product assembly guides",
      "• Corporate training modules",
      "", "REAL EXAMPLES:",
      "Khan Academy lessons",
      "Boeing pilot simulators",
      "IKEA assembly instructions"]),
]
positions = [(0.28,0.93),(6.82,0.93),(0.28,4.15),(6.82,4.15)]
card_ids = []
for (acc, ico, ttl, lns), (lx, ty) in zip(usage_data, positions):
    rct(s4, lx+0.04, ty+0.04, 6.26, 3.1, RGBColor(4,4,18))
    sh = rct(s4, lx, ty, 6.26, 3.1, CARD, lc=RGBColor(28,28,72), lw=0.4)
    rct(s4, lx, ty, 0.06, 3.1, acc)
    rct(s4, lx+0.06, ty, 6.20, 0.003, acc)
    txb(s4, ico, lx+0.12, ty+0.1, 0.65, 0.65, 28, W, align=PP_ALIGN.CENTER)
    txb(s4, ttl, lx+0.82, ty+0.15, 5.3, 0.35, 11, acc, bold=True, font=FH)
    rct(s4, lx+0.82, ty+0.52, 5.3, 0.012, RGBColor(40,40,90))
    multiline_txb(s4, lns, lx+0.12, ty+0.58, 6.0, 2.42, 9.5, LT, accent=acc)
    card_ids.append(sh.shape_id)

animate_shapes(s4, card_ids, start_ms=200, step_ms=200, effect="float")
add_push_transition(s4)

# ── S5 TECHNIQUES ─────────────────────────────────────────────
s5 = new_slide()
bg_dots(s5, RGBColor(40,20,0), count=20, seed=6)
hdr(s5, "ANIMATION TECHNIQUES", "03", O)

tech_data = [
    (C, "CEL ANIMATION",
     ["Traditional hand-drawn technique.",
      "• Drawings on transparent celluloid sheets",
      "• Each sheet = one frame layer",
      "• Photographed frame-by-frame on film",
      "• Allows background reuse (saves cost)",
      "• Colors applied to BACK of cel sheet",
      "• ~12–24 unique cels per second",
      "", "ERA: 1910s – 1990s",
      "TOOLS: Disney, Warner Bros, Toei",
      "", "WORKFLOW:",
      "Story → Storyboard → Layout →",
      "Key Animation → Tweening →",
      "Inking → Painting → Photography"]),
    (O, "PATH ANIMATION",
     ["Object follows a predefined trajectory.",
      "• Animator defines the path shape",
      "• Object moves along it automatically",
      "• Supports straight lines & Bezier curves",
      "• Object can resize/deform along path",
      "• Easing controls speed variation",
      "", "TOOLS: After Effects, CSS, SVG paths",
      "", "REAL EXAMPLES:",
      "Logo reveal animations",
      "Camera fly-through sequences",
      "UI element transitions",
      "Game projectile trajectories"]),
    (G, "COMPUTER ANIMATION",
     ["Digital creation of moving images.",
      "• Can simulate cel & path techniques",
      "• Plus: physics, particles, simulation",
      "• 2D digital, 3D CG, or hybrid mix",
      "• Real-time OR pre-rendered output",
      "• Procedural & AI-assisted animation",
      "", "ERA: 1980s – present",
      "TOOLS: Pixar, DreamWorks, ILM",
      "", "REAL EXAMPLES:",
      "Toy Story, Avengers VFX,",
      "Fortnite, Hogwarts Legacy"]),
]
card_ids = []
for i, (acc, ttl, lns) in enumerate(tech_data):
    lx = 0.28 + i * 4.35
    card(s5, lx, 0.93, 4.21, 6.22, acc, ttl, lns)
    card_ids.append(s5.shapes[-2].shape_id)

animate_shapes(s5, card_ids, start_ms=200, step_ms=220, effect="float")
add_push_transition(s5)

# ── S6 CEL TYPES ──────────────────────────────────────────────
s6 = new_slide()
bg_dots(s6, RGBColor(0,40,80), count=20, seed=7)
hdr(s6, "TYPES — CEL ANIMATION", "04", C)
info_bar(s6, "Named after CELLULOID — transparent plastic sheets invented in 1869. Developed ~1914 by Raoul Barre & John Bray. Dominated animation for 80 years.")

# Keyframe timeline visual
rct(s6, 0.28, 1.38, 12.77, 0.52, RGBColor(10,10,36))
rct(s6, 0.28, 1.38, 0.055, 0.52, C)
txb(s6, "KEYFRAME TIMELINE", 0.4, 1.40, 3.0, 0.22, 8, MU, italic=True)
# Timeline line
rct(s6, 0.6, 1.59, 12.2, 0.02, RGBColor(40,40,80))
# Keyframe markers (diamond shapes)
for kx, lbl in [(1.5,"K1"),(5.0,"K2"),(8.5,"K3"),(12.0,"K4")]:
    rct(s6, kx-0.07, 1.52, 0.14, 0.14, C)
    txb(s6, lbl, kx-0.14, 1.67, 0.4, 0.18, 7, C, bold=True, align=PP_ALIGN.CENTER)
# Tween dots
for tx_ in [2.1,2.6,3.1,3.6,4.1,4.6, 5.8,6.3,6.8,7.3,7.8, 9.2,9.7,10.2,10.7,11.2,11.7]:
    oval(s6, tx_-0.04, 1.545, 0.08, 0.08, RGBColor(0,100,140))
txb(s6, "-- tweens --", 3.0, 1.67, 2.0, 0.18, 7, MU, italic=True)
txb(s6, "-- tweens --", 6.3, 1.67, 2.0, 0.18, 7, MU, italic=True)
txb(s6, "-- tweens --", 9.6, 1.67, 2.0, 0.18, 7, MU, italic=True)

cel_data = [
    (C, "CEL ANIMATION",
     ["• Transparent sheet per frame",
      "• Character on cel; BG separate",
      "• Cels stacked = final composite",
      "• Colors on BACK of cel",
      "• ~12–24 unique cels per second",
      "", "USED IN:", "Tom & Jerry, Snow White,",
      "Akira, Dragon Ball Z"]),
    (O, "KEYFRAMES",
     ["The 'anchor' frames of an action.",
      "• Show most extreme/critical poses",
      "• Drawn by SENIOR animators",
      "• Define timing & structure",
      "", "EXAMPLES:",
      "Frame 1  =  standing still",
      "Frame 12 =  peak of jump  UP",
      "Frame 24 =  landed safely DOWN",
      "", "Modern software uses identical",
      "keyframe concept on timelines"]),
    (G, "TWEENING (In-Betweening)",
     ["Frames drawn BETWEEN keyframes.",
      "• Creates smooth motion transition",
      "• Done by JUNIOR animators",
      "• More tweens = smoother motion",
      "• At 24fps: ~22 tweens per action",
      "", "MODERN AUTO-TWEEN TOOLS:",
      "Adobe Animate, Flash,",
      "Toon Boom Harmony, After Effects"]),
    (P, "PENCIL TEST",
     ["Quick pre-production quality check.",
      "• Rough pencil sketches filmed first",
      "• Played back to check timing & flow",
      "• Catches errors BEFORE expensive",
      "  inking & painting stages",
      "• Saves significant time & money",
      "", "MODERN EQUIVALENT:",
      "Animatic / rough-cut preview in",
      "editing software"]),
]
positions = [(0.28,1.98),(6.78,1.98),(0.28,4.68),(6.78,4.68)]
card_ids = []
for (acc, ttl, lns), (lx, ty) in zip(cel_data, positions):
    card(s6, lx, ty, 6.22, 2.6, acc, ttl, lns)
    card_ids.append(s6.shapes[-2].shape_id)

animate_shapes(s6, card_ids, start_ms=250, step_ms=180, effect="float")
add_push_transition(s6)

# ── S7 PATH ANIMATION ─────────────────────────────────────────
s7 = new_slide()
bg_dots(s7, RGBColor(50,30,0), count=18, seed=8)
hdr(s7, "TYPES — PATH ANIMATION", "04", O)
info_bar(s7, "An object moves along a PREDEFINED TRAJECTORY. The animator draws the path; software handles movement along it — with full control over timing and easing.")

# Left: path diagram card
rct(s7, 0.28+0.04, 1.55+0.04, 6.22, 5.62, RGBColor(4,4,18))
rct(s7, 0.28, 1.55, 6.22, 5.62, CARD, lc=RGBColor(28,28,72), lw=0.4)
rct(s7, 0.28, 1.55, 0.06, 5.62, O)
txb(s7, "PATH TYPES & CONTROLS", 0.40, 1.65, 6.0, 0.28, 10.5, O, bold=True, font=FH)
rct(s7, 0.40, 1.96, 6.0, 0.012, RGBColor(40,40,90))

# Straight line
txb(s7, "X  STRAIGHT LINE  (rigid / mechanical)", 0.40, 2.04, 5.9, 0.25, 9.5, RGBColor(220,60,100))
rct(s7, 0.5, 2.35, 5.5, 0.03, RGBColor(180,60,90))
oval(s7, 0.5-0.12, 2.27, 0.24, 0.24, RGBColor(220,60,100))
oval(s7, 5.9, 2.27, 0.24, 0.24, RGBColor(220,60,100))

# Bezier path
txb(s7, "BEZIER CURVE  (smooth / organic — preferred)", 0.40, 2.58, 5.9, 0.25, 9.5, G)
# Approximate curve with segments
import math
for seg in range(10):
    t_ = seg / 10; t2 = (seg+1)/10
    # Simple parabola: y = 3.1 - 0.65*(1-(2t-1)^2)
    y1 = 3.1 - 0.65*(1-(2*t_-1)**2)
    y2 = 3.1 - 0.65*(1-(2*t2-1)**2)
    x1 = 0.5 + t_*5.4; x2 = 0.5 + t2*5.4
    rct(s7, x1, y1, x2-x1+0.02, 0.03, G)
# Control point markers
oval(s7, 0.5-0.1, 3.1-0.1, 0.2, 0.2, O)
oval(s7, 5.85, 3.1-0.1, 0.2, 0.2, O)
oval(s7, 2.0, 2.47, 0.16, 0.16, RGBColor(255,200,100), lc=O, lw=0.75)
oval(s7, 4.0, 2.47, 0.16, 0.16, RGBColor(255,200,100), lc=O, lw=0.75)
rct(s7, 0.5, 3.08, 1.62, 0.015, RGBColor(255,145,0,))
rct(s7, 4.1, 3.08, 1.83, 0.015, RGBColor(255,145,0))
txb(s7, "Control handles", 1.8, 2.36, 2.5, 0.22, 7.5, MU, italic=True, align=PP_ALIGN.CENTER)
txb(s7, "Easing: Slow → Fast → Slow (gives weight)", 0.40, 3.28, 5.9, 0.22, 8.5, MU, italic=True)

multiline_txb(s7,
    ["", "RESIZING & RESHAPING:",
     "Object can scale as it travels the path",
     "(e.g., ball squashing on bounce impact)",
     "", "TOOLS: After Effects • CSS • SVG",
     "        Unity NavMesh • Game engines"],
    0.40, 3.55, 5.9, 3.5, 9.5, LT, accent=O)

# Right: steps card
rct(s7, 6.78+0.04, 1.55+0.04, 6.22, 5.62, RGBColor(4,4,18))
rct(s7, 6.78, 1.55, 6.22, 5.62, CARD, lc=RGBColor(28,28,72), lw=0.4)
rct(s7, 6.78, 1.55, 0.06, 5.62, O)
txb(s7, "HOW PATH ANIMATION WORKS", 6.90, 1.65, 6.0, 0.28, 10.5, O, bold=True, font=FH)
rct(s7, 6.90, 1.96, 6.0, 0.012, RGBColor(40,40,90))

steps = [
    (C,  "1", "Draw the PATH — straight, Bezier, or freehand"),
    (O,  "2", "Set START and END points along the path"),
    (G,  "3", "Define DURATION — travel time for the object"),
    (P,  "4", "Apply EASING — slow in, slow out (natural feel)"),
    (PK, "5", "Set ORIENTATION — face direction of travel?"),
    (Y,  "6", "Preview, adjust keyframes, render final output"),
]
for i, (acc, num, desc) in enumerate(steps):
    ty_ = 2.05 + i * 0.83
    rct(s7, 6.92, ty_, 0.50, 0.50, RGBColor(20,20,64))
    txb(s7, num, 6.92, ty_+0.06, 0.50, 0.42, 16, acc,
        bold=True, align=PP_ALIGN.CENTER, font=FH)
    txb(s7, desc, 7.52, ty_+0.1, 5.3, 0.42, 9.5, LT)

add_push_transition(s7)

# ── S8 COMPUTER ANIMATION VS ──────────────────────────────────
s8 = new_slide()
bg_dots(s8, RGBColor(0,50,20), count=20, seed=9)
hdr(s8, "COMPUTER ANIMATION", "05", G)
info_bar(s8, "Computer animation uses software to create motion — producing effects impossible by hand. Revolutionized the industry from the 1980s onward.")

txb(s8, "TRADITIONAL CEL  vs  COMPUTER ANIMATION",
    0, 1.52, 13.33, 0.38, 15, Y, bold=True, align=PP_ALIGN.CENTER, font=FH)

# CEL column
rct(s8, 0.28+0.04, 1.98+0.04, 5.9, 5.2, RGBColor(4,4,18))
sh_cel = rct(s8, 0.28, 1.98, 5.9, 5.2, CARD, lc=RGBColor(28,28,72), lw=0.4)
rct(s8, 0.28, 1.98, 5.9, 0.42, RGBColor(0,50,80))
txb(s8, "CEL ANIMATION", 0.36, 2.04, 5.74, 0.34, 13, C,
    bold=True, align=PP_ALIGN.CENTER, font=FH)

cel_rows = [
    "Hand-drawn on physical celluloid sheets",
    "Every frame drawn manually — very slow",
    "Physical storage: film cans & warehouses",
    "Expensive: cels, paint, camera lab costs",
    "Difficult to modify finished frames",
    "Organic hand-crafted visual warmth",
    "Production: months per 1 minute of film",
]
for i, row in enumerate(cel_rows):
    rct(s8, 0.35, 2.48+i*0.38, 5.75, 0.36,
        RGBColor(16,16,54) if i%2==0 else RGBColor(20,20,62))
    txb(s8, row, 0.45, 2.5+i*0.38, 5.55, 0.32, 9.5, LT)

# VS badge
rct(s8, 6.2, 3.4, 0.93, 1.1, RGBColor(30,20,5), lc=O, lw=1)
txb(s8, "VS", 6.2, 3.55, 0.93, 0.8, 28, O, bold=True,
    align=PP_ALIGN.CENTER, font=FH)

# Computer column
rct(s8, 7.15+0.04, 1.98+0.04, 5.9, 5.2, RGBColor(4,4,18))
sh_comp = rct(s8, 7.15, 1.98, 5.9, 5.2, CARD, lc=RGBColor(28,28,72), lw=0.4)
rct(s8, 7.15, 1.98, 5.9, 0.42, RGBColor(5,50,20))
txb(s8, "COMPUTER ANIMATION", 7.23, 2.04, 5.74, 0.34, 13, G,
    bold=True, align=PP_ALIGN.CENTER, font=FH)

comp_rows = [
    "CREATED entirely within digital software",
    "Automation handles tweens & physics",
    "Digital files — backup & share instantly",
    "One-time software cost; no consumables",
    "Copy, adjust, modify any frame instantly",
    "Any style: toon to photorealistic",
    "Real-time previews; much faster pipeline",
]
for i, row in enumerate(comp_rows):
    rct(s8, 7.22, 2.48+i*0.38, 5.75, 0.36,
        RGBColor(16,16,54) if i%2==0 else RGBColor(20,20,62))
    txb(s8, "CHECK  " + row, 7.32, 2.5+i*0.38, 5.55, 0.32, 9.5, G)

animate_shapes(s8, [sh_cel.shape_id, sh_comp.shape_id],
               start_ms=200, step_ms=300, effect="float")
add_push_transition(s8)

# ── S9 FK / IK / MORPHING ─────────────────────────────────────
s9 = new_slide()
bg_dots(s9, RGBColor(0,50,20), count=20, seed=10)
hdr(s9, "KINEMATICS & MORPHING", "05", G)

fk_data = [
    (C, "FORWARD KINEMATICS (FK)",
     ["Animating from PARENT to CHILD joint.",
      "", "HOW IT WORKS:",
      "Move SHOULDER ==>",
      "Elbow follows automatically ==>",
      "Wrist follows elbow ==>",
      "Finger follows wrist",
      "", "ANALOGY: A puppet — control the",
      "parent; all children follow in chain.",
      "", "BEST FOR:",
      "• Mechanical / robotic arms",
      "• Waving & swinging motions",
      "• Tail & tentacle animation"]),
    (O, "INVERSE KINEMATICS (IK)",
     ["Reverse of FK — move END POINT;",
      "the chain auto-adjusts.",
      "", "HOW IT WORKS:",
      "Move HAND position ==>",
      "Software calculates elbow ==>",
      "And shoulder automatically",
      "", "ANALOGY: Reaching for a cup —",
      "you aim your hand; body figures",
      "out the rest automatically.",
      "", "BEST FOR:",
      "• Character walking & running",
      "• Reaching & grabbing actions"]),
    (P, "MORPHING",
     ["Smooth transformation from one shape",
      "into a completely DIFFERENT one.",
      "", "HOW IT WORKS:",
      "Map control points on source",
      "Map matching points on target",
      "Interpolate all points together",
      "Creates seamless in-between frames",
      "", "TYPES: 2D image morph • 3D mesh",
      "deformation • Vector path morphing",
      "", "FAMOUS EXAMPLES:",
      "MJ 'Black or White' (face morph)",
      "T-1000 — Terminator 2"]),
]
card_ids = []
for i, (acc, ttl, lns) in enumerate(fk_data):
    lx = 0.28 + i * 4.35
    rct(s9, lx+0.04, 0.93+0.04, 4.21, 6.22, RGBColor(4,4,18))
    sh = rct(s9, lx, 0.93, 4.21, 6.22, CARD, lc=RGBColor(28,28,72), lw=0.4)
    rct(s9, lx, 0.93, 0.06, 6.22, acc)
    txb(s9, ttl, lx+0.12, 1.03, 4.0, 0.28, 10.5, acc, bold=True, font=FH)
    rct(s9, lx+0.12, 1.34, 4.0, 0.012, RGBColor(40,40,90))

    # Joint diagram
    jx = lx + 2.0
    for j in range(3):
        jy = 1.48 + j*0.58
        oval(s9, jx-0.2, jy-0.2, 0.4, 0.4, RGBColor(20,20,64), lc=acc, lw=1.5)
        labels = ["SHOULDER","ELBOW","WRIST"]
        txb(s9, labels[j], jx-0.25, jy-0.1, 0.5, 0.2, 6, acc,
            align=PP_ALIGN.CENTER, bold=True)
        if j < 2:
            rct(s9, jx-0.03, jy+0.2, 0.06, 0.38, acc)

    if i == 1:  # IK target
        rct(s9, jx-0.35, 3.4, 0.7, 0.35, RGBColor(40,20,5), lc=acc, lw=1)
        txb(s9, "TARGET", jx-0.35, 3.42, 0.7, 0.28, 7, acc,
            bold=True, align=PP_ALIGN.CENTER)
    if i == 2:  # Morphing shapes
        oval(s9, lx+0.5, 1.48, 0.7, 0.7, RGBColor(40,10,60), lc=P, lw=1.5)
        txb(s9, "-->", lx+1.28, 1.6, 0.5, 0.4, 16, P, align=PP_ALIGN.CENTER)
        rct(s9, lx+1.85, 1.48, 0.7, 0.7, RGBColor(40,10,60), lc=P, lw=1.5)
        txb(s9, "Circle --> Star", lx+0.4, 2.3, 2.3, 0.22, 7.5, MU,
            italic=True, align=PP_ALIGN.CENTER)

    multiline_txb(s9, lns, lx+0.12, 3.22, 4.0, 3.8, 9, LT, accent=acc)
    card_ids.append(sh.shape_id)

animate_shapes(s9, card_ids, start_ms=200, step_ms=220, effect="float")
add_push_transition(s9)

# ── S10 LEVELS ────────────────────────────────────────────────
s10 = new_slide()
bg_dots(s10, RGBColor(50,45,0), count=22, seed=11)
hdr(s10, "LEVELS OF COMPUTER ANIMATION", "06", Y)

level_data = [
    (C, "BASIC  LEVEL   [  *  ]",
     ["WHAT YOU CREATE:",
      "• Animated GIFs",
      "• CSS / JS transitions",
      "• PowerPoint animations",
      "• Simple 2D sprites",
      "• Basic logo animations",
      "", "TOOLS:", "MS PowerPoint, Canva, GIMP",
      "", "Skill: Beginner",
      "Time: Days to Weeks",
      "", "CAREER:",
      "Social media creator,",
      "Marketing & design assistant"]),
    (O, "INTERMEDIATE  LEVEL  [  * *  ]",
     ["WHAT YOU CREATE:",
      "• 2D character animation",
      "• Motion graphics & infographics",
      "• Complex path animations",
      "• Basic 3D object animation",
      "• Game character movement",
      "", "TOOLS:", "After Effects, Adobe Animate,",
      "Blender, Spine, Toon Boom",
      "", "Skill: Intermediate",
      "Time: Months",
      "", "CAREER:",
      "Motion designer, 2D animator"]),
    (P, "ADVANCED  LEVEL  [  * * *  ]",
     ["WHAT YOU CREATE:",
      "• Full CG feature films",
      "• Realistic physics simulation",
      "• Procedural animation",
      "• Motion capture integration",
      "• Real-time VFX in games",
      "• AI-assisted animation",
      "", "TOOLS:", "Autodesk Maya, Houdini,",
      "Unreal Engine, Cinema 4D",
      "", "Skill: Professional",
      "Time: Years",
      "", "CAREER:",
      "VFX artist, 3D animator, TD"]),
]
card_ids = []
for i, (acc, ttl, lns) in enumerate(level_data):
    lx = 0.28 + i * 4.35
    rct(s10, lx+0.04, 0.93+0.04, 4.21, 6.22, RGBColor(4,4,18))
    sh = rct(s10, lx, 0.93, 4.21, 6.22, CARD, lc=RGBColor(28,28,72), lw=0.4)
    rct(s10, lx, 0.93, 0.06, 6.22, acc)
    txb(s10, ttl, lx+0.12, 1.02, 4.0, 0.28, 10.5, acc, bold=True, font=FH)
    # Skill bar
    rct(s10, lx+0.12, 1.36, 4.0, 0.16, RGBColor(25,25,65))
    rct(s10, lx+0.12, 1.36, [1.33,2.67,4.0][i], 0.16, acc)
    rct(s10, lx+0.12, 1.55, 4.0, 0.012, RGBColor(40,40,90))
    multiline_txb(s10, lns, lx+0.12, 1.6, 4.0, 5.45, 9, LT, accent=acc)
    card_ids.append(sh.shape_id)

animate_shapes(s10, card_ids, start_ms=200, step_ms=220, effect="float")
add_push_transition(s10)

# ── S11–S13: 12 PRINCIPLES (helper) ───────────────────────────
def principles_slide(prs_obj, accent, part, title_sfx, info_text, cards_data):
    sl = new_slide()
    bg_dots(sl, RGBColor(0,40,80), count=18, seed=11+part)
    hdr(sl, f"12 PRINCIPLES OF ANIMATION  •  Part {part} / 3", "07", accent)
    info_bar(sl, info_text, accent=accent)
    positions = [(0.28,1.55),(6.78,1.55),(0.28,4.28),(6.78,4.28)]
    cids = []
    for (acc, ttl, lns), (lx, ty) in zip(cards_data, positions):
        card(sl, lx, ty, 6.22, 2.65, acc, ttl, lns)
        cids.append(sl.shapes[-2].shape_id)
    animate_shapes(sl, cids, start_ms=250, step_ms=180, effect="float")
    add_push_transition(sl)
    return sl

s11 = principles_slide(prs, C, 1, "P1",
    "Developed by Disney masters Ollie Johnston & Frank Thomas (The Illusion of Life, 1981). Universal gold standard for ALL animators.",
    [
        (C, "01 — SQUASH & STRETCH",
         ["Exaggerate deformation for weight.",
          "• Object SQUASHES on impact",
          "• Object STRETCHES when moving fast",
          "• Golden rule: VOLUME stays constant",
          "• Shows elasticity, mass, and energy",
          "", "EXAMPLES:",
          "Rubber ball bounce — squashes flat",
          "on landing, stretches tall in mid-air",
          "", "Most fundamental principle of animation"]),
        (O, "02 — ANTICIPATION",
         ["Small opposite action BEFORE main action.",
          "• Character crouches before jumping",
          "• Arm swings back before throwing",
          "• Boxer winds up before punching",
          "• Prepares the audience mentally",
          "", "EXAMPLES:",
          "Bugs Bunny ducking low before leap",
          "", "Without anticipation, actions feel",
          "sudden, unexpected, unconvincing"]),
        (G, "03 — STAGING",
         ["Present ONE clear idea at a time.",
          "• Use clear, readable silhouettes",
          "• Camera angle serves the story",
          "• Avoid distracting backgrounds",
          "• Direct the audience's eye deliberately",
          "", "EXAMPLES:",
          "Character in profile view making",
          "a dramatic gesture — instantly readable",
          "", "Think: theater staging, not photography"]),
        (P, "04 — STRAIGHT AHEAD & POSE-TO-POSE",
         ["Two approaches to creating animation:",
          "", "STRAIGHT AHEAD:",
          "Draw frame by frame from start to end",
          "= spontaneous, fluid, surprising",
          "", "POSE-TO-POSE:",
          "Draw key poses first, fill in betweens",
          "= controlled, planned, consistent",
          "", "Best professional work combines BOTH"]),
    ])

s12 = principles_slide(prs, O, 2, "P2",
    "These principles govern how objects move AFTER the main action — the subtle details that separate amateur from truly professional animation.",
    [
        (C, "05 — FOLLOW THROUGH & OVERLAPPING",
         ["FOLLOW THROUGH:",
          "Parts keep moving after main action stops",
          "• Hair swings after head stops",
          "• Coat settles after character stops",
          "", "OVERLAPPING ACTION:",
          "Different parts move at different rates",
          "• Arms lag slightly behind torso",
          "", "Makes motion feel ORGANIC & real"]),
        (O, "06 — EASE IN & EASE OUT",
         ["Objects ACCELERATE at start and",
          "DECELERATE at end of every movement.",
          "• Nothing in nature starts at full speed",
          "• Car: slow start = full speed = brake",
          "", "IN SOFTWARE:",
          "Bezier easing curves control this.",
          "Linear = robotic feel",
          "Ease = natural, believable weight",
          "", "Most impactful principle for realism"]),
        (G, "07 — ARCS",
         ["Natural motion follows ARC paths,",
          "NOT straight lines.",
          "• Arms swing in arcs",
          "• Heads nod in arcs",
          "• Thrown objects travel in parabolas",
          "", "EXCEPTION:",
          "Mechanical / robot movement IS",
          "intentionally straight — use linear",
          "paths deliberately for machines",
          "", "Breaking arcs = looks robotic"]),
        (P, "08 — SECONDARY ACTION",
         ["Supporting actions that ENRICH",
          "the primary movement.",
          "", "EXAMPLES:",
          "Walking + arms swing (secondary)",
          "Running + hair flows (secondary)",
          "Talking + eyebrow (secondary)",
          "", "CRITICAL RULE:",
          "Secondary must NEVER overpower main",
          "action — supports, never steals focus"]),
    ])

s13 = principles_slide(prs, P, 3, "P3",
    'Walt Disney: "The secret of animation is the exaggeration of the truth." These final 4 principles deliver exactly that truth.',
    [
        (C, "09 — TIMING",
         ["Frames = perceived speed, weight & personality.",
          "MORE frames = slower / heavier",
          "FEWER frames = faster / lighter",
          "", "AT 24fps:",
          "• Eye blink: 6–8 frames",
          "• Head turn: 10–15 frames",
          "• Heavy punch: 3–4 frames",
          "• Gentle wave: 30+ frames",
          "", "TIMING is the HEARTBEAT of animation"]),
        (O, "10 — EXAGGERATION",
         ["Push actions BEYOND realistic limits",
          "for emotional impact.",
          "• Giant cartoon eyes & reactions",
          "• Impossible stretching",
          "• Over-the-top expressions",
          "", "KEY INSIGHT:",
          "NOT about being unrealistic —",
          "it's making the TRUTH more VISIBLE",
          "and ENTERTAINING to the audience"]),
        (G, "11 — SOLID DRAWING",
         ["Characters must feel like 3D objects",
          "with real weight and form.",
          "• Understand 3D basic shapes",
          "• 'Draw through' the form",
          "• Consistent proportions all angles",
          "• Avoid 'twinning' (mirror poses)",
          "• Strong line of action through body",
          "", "Even in 2D — audience must FEEL",
          "the three-dimensional volume"]),
        (P, "12 — APPEAL",
         ["Characters must be visually interesting.",
          "", "APPEAL is NOT CUTENESS:",
          "• A villain has appeal",
          "• An anti-hero has appeal",
          "• Even monsters can have appeal",
          "", "HOW TO ACHIEVE:",
          "• Clear, readable design",
          "• Dynamic, asymmetric poses",
          "• Charisma & personality",
          "• Strong distinctive silhouette"]),
    ])

# ── S14 ANIMATION SPACE ───────────────────────────────────────
s14 = new_slide()
bg_dots(s14, RGBColor(50,0,30), count=22, seed=14)
hdr(s14, "ANIMATION SPACE", "08", PK)
info_bar(s14, "Animation exists in three spatial dimensions — each with different tools, techniques, visual results, and hardware requirements.", accent=PK)

space_data = [
    (C, "2D SPACE  (X + Y only)",
     ["Width + Height only. No depth.",
      "", "CHARACTERISTICS:",
      "• Flat, like drawing on paper",
      "• No real depth perception",
      "• Simple to create & render",
      "• Low hardware requirements",
      "• Classic cartoon aesthetic",
      "• Fast production pipeline",
      "", "EXAMPLES:",
      "Tom & Jerry, Classic Disney,",
      "Original Mario & Pac-Man,",
      "Mobile game animations",
      "", "TOOLS:",
      "Adobe Animate, Toon Boom,",
      "Spine, TVPaint, Aseprite"]),
    (O, "2.5D SPACE  (pseudo-3D)",
     ["2D art + ILLUSION of depth via",
      "layers, parallax & shadows.",
      "", "CHARACTERISTICS:",
      "• Parallax scrolling creates depth",
      "• Objects appear on Z-axis visually",
      "• Still technically 2D sprites",
      "• 'Pseudo-3D' / 'Layered 2D'",
      "• Great performance/visual tradeoff",
      "", "EXAMPLES:",
      "Paper Mario, South Park,",
      "Ori and the Blind Forest",
      "Disney Multiplane Camera",
      "", "TOOLS:",
      "After Effects Z-depth,",
      "Unity 2D with parallax"]),
    (G, "3D SPACE  (X + Y + Z)",
     ["Full volumetric environment.",
      "True depth, volume, perspective.",
      "", "CHARACTERISTICS:",
      "• Viewable from any camera angle",
      "• Realistic lighting & shadows",
      "• Requires powerful hardware",
      "• Longer render times",
      "• Maximum visual versatility",
      "", "EXAMPLES:",
      "Toy Story, Frozen, Moana,",
      "Fortnite, Call of Duty,",
      "Architectural visualizations",
      "", "TOOLS:",
      "Blender, Autodesk Maya,",
      "Cinema 4D, Unreal Engine"]),
]
card_ids = []
for i, (acc, ttl, lns) in enumerate(space_data):
    lx = 0.28 + i * 4.35
    rct(s14, lx+0.04, 1.55+0.04, 4.21, 5.7, RGBColor(4,4,18))
    sh = rct(s14, lx, 1.55, 4.21, 5.7, CARD, lc=RGBColor(28,28,72), lw=0.4)
    rct(s14, lx, 1.55, 0.06, 5.7, acc)

    # Visual diagram per space type
    if i == 0:  # 2D flat box with face
        rct(s14, lx+0.5, 1.65, 3.2, 0.8, RGBColor(8,24,50), lc=acc, lw=1.5)
        for ex in [lx+1.1, lx+2.5]:
            oval(s14, ex-0.12, 1.82, 0.24, 0.24, acc)
        txb(s14, "2D FLAT — X  Y", lx+0.5, 1.72, 3.2, 0.22, 8, MU,
            align=PP_ALIGN.CENTER)
    elif i == 1:  # 2.5D layers
        for j, (ly_, opa) in enumerate([(1.65, 0.15),(1.83, 0.25),(2.01, 0.4)]):
            rct(s14, lx+0.3+j*0.2, ly_, 3.2-j*0.4, 0.32,
                RGBColor(int(8+j*10), int(24+j*8), int(50+j*10)), lc=acc, lw=0.75)
            txb(s14, ["BG","MID","FG"][j], lx+0.4+j*0.2, ly_+0.07, 0.5, 0.2, 7, acc,
                align=PP_ALIGN.CENTER)
    else:  # 3D cube
        rct(s14, lx+0.5, 1.68, 1.8, 0.36, RGBColor(15,50,25), lc=acc, lw=1)
        rct(s14, lx+0.5, 2.04, 0.9, 0.5,  RGBColor(10,35,18), lc=acc, lw=1)
        rct(s14, lx+1.4, 2.04, 0.9, 0.5,  RGBColor(12,42,22), lc=acc, lw=1)
        txb(s14, "X    Y    Z", lx+0.5, 2.56, 1.8, 0.22, 8, acc,
            bold=True, align=PP_ALIGN.CENTER)

    txb(s14, ttl, lx+0.12, 2.65, 4.0, 0.28, 10.5, acc, bold=True, font=FH)
    rct(s14, lx+0.12, 2.97, 4.0, 0.012, RGBColor(40,40,90))
    multiline_txb(s14, lns, lx+0.12, 3.02, 4.0, 4.15, 9, LT, accent=acc)
    card_ids.append(sh.shape_id)

animate_shapes(s14, card_ids, start_ms=200, step_ms=220, effect="float")
add_push_transition(s14)

# ── S15 ANIMATION PROCESS ─────────────────────────────────────
s15 = new_slide()
bg_dots(s15, RGBColor(40,25,0), count=18, seed=15)
hdr(s15, "ANIMATION PROCESS", "09", O)
info_bar(s15, "Every animation project — from a 10-second ad to a 2-hour feature film — follows this same 4-stage pipeline.")

steps_data = [
    (C,  "1", "ORGANIZE EXECUTION",
     ["• Define project goals & scope",
      "• Assign team roles",
      "• Set milestones & deadlines",
      "• Gather reference material",
      "• Create style guide",
      "• Budget & risk planning",
      "", "OUTPUT:",
      "Project bible, schedule,",
      "mood boards, storyboard"]),
    (O,  "2", "CHOOSE TOOL",
     ["• Match tool to project type",
      "• Consider team skill level",
      "• Account for budget",
      "", "2D: Adobe Animate, Toon Boom",
      "3D: Maya, Blender, Cinema 4D",
      "Motion: After Effects",
      "Games: Unity, Unreal Engine"]),
    (G,  "3", "BUILD & TWEAK",
     ["• Create rough blocking first",
      "• Refine with keyframes",
      "• Apply the 12 Principles",
      "• Sync with audio / music",
      "• Team review cycles",
      "• Iterate until approved",
      "", "NOTE:",
      "This step is ITERATIVE —",
      "cycle through many times"]),
    (P,  "4", "POST-PROCESS",
     ["• Render final frames",
      "• Color grading & correction",
      "• Add VFX & overlays",
      "• Sound design & music mix",
      "• Quality assurance testing",
      "• Format for delivery",
      "• Archive project files",
      "", "OUTPUT:",
      "Final MP4, MOV, WebM"]),
]
icons = ["PLAN","TOOL","BUILD","POST"]
card_w = 2.88
card_ids = []
for i, (acc, num, ttl, lns) in enumerate(steps_data):
    lx = 0.28 + i*(card_w+0.38)
    if i > 0:
        txb(s15, ">", lx-0.37, 3.4, 0.34, 0.55, 22, O, bold=True,
            align=PP_ALIGN.CENTER, font=FH)
    rct(s15, lx+0.04, 1.55+0.04, card_w, 5.62, RGBColor(4,4,18))
    sh = rct(s15, lx, 1.55, card_w, 5.62, CARD, lc=RGBColor(28,28,72), lw=0.4)
    rct(s15, lx, 1.55, 0.06, 5.62, acc)
    rct(s15, lx, 1.55, card_w, 0.36, acc)
    txb(s15, f"STEP {num}", lx, 1.57, card_w, 0.32, 11, BG,
        bold=True, align=PP_ALIGN.CENTER, font=FH)
    # Circle number
    oval(s15, lx+(card_w/2)-0.38, 1.99, 0.76, 0.76, RGBColor(20,20,64), lc=acc, lw=2)
    txb(s15, num, lx+(card_w/2)-0.38, 2.05, 0.76, 0.64, 24, acc,
        bold=True, align=PP_ALIGN.CENTER, font=FH)
    txb(s15, icons[i], lx+(card_w/2)-1.0, 2.85, 2.0, 0.28, 9, acc,
        bold=True, align=PP_ALIGN.CENTER, font=FH)
    txb(s15, ttl, lx+0.12, 3.18, card_w-0.18, 0.28, 10.5, acc, bold=True, font=FH)
    rct(s15, lx+0.12, 3.49, card_w-0.18, 0.012, RGBColor(40,40,90))
    multiline_txb(s15, lns, lx+0.12, 3.54, card_w-0.18, 3.55, 9, LT, accent=acc)
    card_ids.append(sh.shape_id)

animate_shapes(s15, card_ids, start_ms=200, step_ms=200, effect="float")
add_push_transition(s15)

# ── S16 SUMMARY ───────────────────────────────────────────────
s16 = new_slide()
bg_dots(s16, RGBColor(0,30,60), count=25, seed=16)
hdr(s16, "SUMMARY — KEY TAKEAWAYS", accent=C)

summary = [
    (C,  "01", "Introduction",     "Animation = sequential images + Persistence of Vision + Phi Phenomenon = seamless illusion of motion."),
    (P,  "02", "Usage",            "Art, storytelling, scientific visualization, instructional training. Used across every major industry."),
    (O,  "03", "Techniques",       "Three core: Cel (hand-drawn), Path (trajectory-based), Computer (digital — most versatile today)."),
    (G,  "04", "Cel Animation",    "Keyframes + tweening = smooth motion. Pencil test validates timing before expensive final production."),
    (PK, "05", "Path Animation",   "Object follows predefined path (straight or Bezier) with easing for natural, weighted motion feel."),
    (C,  "06", "Computer Anim.",   "FK/IK chains for rigs. Morphing for shape transforms. Far faster & more flexible than traditional cel."),
    (P,  "07", "Levels",           "Basic (GIFs/CSS) — Intermediate (motion graphics) — Advanced (feature films, VFX, AI animation)."),
    (O,  "08", "12 Principles",    "Disney's Squash & Stretch, Anticipation, Timing, Ease, Arcs, and 8 more — universal standard."),
    (G,  "09", "Animation Space",  "2D (flat) — 2.5D (pseudo-depth via parallax) — 3D (full volumetric, any camera angle)."),
    (Y,  "10", "Process",          "Organize -- Choose Tool -- Build & Tweak (iterative) -- Post-Process. Every project, every studio."),
]
row_ids = []
for i, (acc, num, ttl, desc) in enumerate(summary):
    col = i // 5; row = i % 5
    lx = 0.28 if col == 0 else 6.95
    ty = 0.94 + row * 1.22
    rct(s16, lx+0.03, ty+0.03, 6.2, 1.15, RGBColor(4,4,18))
    sh = rct(s16, lx, ty, 6.2, 1.15, CARD, lc=RGBColor(28,28,72), lw=0.4)
    rct(s16, lx, ty, 0.06, 1.15, acc)
    rct(s16, lx+0.06, ty, 6.14, 0.003, acc)
    rct(s16, lx+0.12, ty+0.08, 0.62, 0.62, RGBColor(20,20,64))
    txb(s16, num, lx+0.12, ty+0.12, 0.62, 0.55, 18, acc,
        bold=True, align=PP_ALIGN.CENTER, font=FH)
    txb(s16, ttl,  lx+0.84, ty+0.08, 5.25, 0.3, 10.5, acc, bold=True, font=FH)
    txb(s16, desc, lx+0.84, ty+0.4,  5.25, 0.68, 9, LT)
    row_ids.append(sh.shape_id)

rct(s16, 6.72, 0.94, 0.055, 6.1, RGBColor(35,35,80))
animate_shapes(s16, row_ids, start_ms=200, step_ms=120, effect="float")
add_push_transition(s16)

# ── S17 THANK YOU ─────────────────────────────────────────────
s17 = new_slide()
bg_dots(s17, RGBColor(0,40,80), count=35, seed=17)
bg_dots(s17, RGBColor(40,0,80), count=25, seed=18)

# Big glow circles
for r, col in [(3.8, RGBColor(6,12,36)),(2.6, RGBColor(8,16,44)),
               (1.6, RGBColor(12,22,54)),(0.6, RGBColor(20,36,80))]:
    oval(s17, 12.83-r, 0.5-r, r*2, r*2, col, lc=RGBColor(0,50,90), lw=0.5)

# Film reel
for r_, col in [(0.75, C),(0.45, C),(0.15, C)]:
    oval(s17, 11.8-r_, 0.2-r_, r_*2, r_*2, RGBColor(8,20,50),
         lc=RGBColor(0,100,140), lw=1.0)

e1 = txb(s17, "GRAPHICS & MULTIMEDIA  •  SMIT", 0, 1.0, 13.33, 0.45,
          13, O, bold=True, align=PP_ALIGN.CENTER, font=FH)
e2 = txb(s17, "THANK YOU", 0, 1.45, 13.33, 2.3, 100, C,
          bold=True, align=PP_ALIGN.CENTER, font=FH)
e3 = txb(s17, "Animation: Where Imagination Meets Technology",
          0, 3.55, 13.33, 0.5, 15, LT, align=PP_ALIGN.CENTER, italic=True)
e4 = rct(s17, 4.4, 4.2, 4.53, 0.025, C)
e5 = txb(s17, "GROUP 3", 0, 4.36, 13.33, 0.38, 13, O,
          bold=True, align=PP_ALIGN.CENTER, font=FH)

mem_ids2 = []
for i, (name, roll) in enumerate(members):
    lx = 1.5 + i * 3.6
    rct(s17, lx, 4.9, 3.35, 0.82, RGBColor(16,16,58), lc=C, lw=0.5)
    m = txb(s17, name, lx+0.15, 4.95, 3.05, 0.35, 12, W, bold=True, font=FH)
    txb(s17, roll, lx+0.15, 5.3, 3.05, 0.3, 9.5, MU)
    mem_ids2.append(m.shape_id)

txb(s17, "Introduction  •  Usage  •  Techniques  •  Types  •  Computer Animation  •  Levels  •  12 Principles  •  Space  •  Process",
    0, 5.88, 13.33, 0.35, 8.5, MU, align=PP_ALIGN.CENTER, italic=True)

animate_title_elements(s17, [e1.shape_id, e2.shape_id, e3.shape_id,
                              e4.shape_id, e5.shape_id] + mem_ids2)
add_fade_transition(s17, 800)

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
out = r"C:\Users\Noman Traders\Desktop\SMIT -AI-CLAUDE\Graphics n Multimedia\Animation_Group3_PREMIUM.pptx"
prs.save(out)
print("SAVED OK ->", out)
