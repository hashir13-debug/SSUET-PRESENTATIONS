"""
Animation Presentation — ULTRA v3
Group 3 | SMIT Graphics & Multimedia
Film-strip theme + visual diagrams
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

BG    = RGBColor(4,   4,  18)
CARD  = RGBColor(13, 13,  48)
CARD2 = RGBColor(19, 19,  64)
FILM  = RGBColor(7,   7,  28)
CYAN  = RGBColor(0,  215, 255)
PURP  = RGBColor(155, 80, 225)
ORG   = RGBColor(255, 148,  0)
GRN   = RGBColor(46,  220, 115)
PINK  = RGBColor(255,  90, 162)
YEL   = RGBColor(255, 215,   0)
WHT   = RGBColor(255, 255, 255)
LT    = RGBColor(192, 206, 238)
DARK  = RGBColor(4,    4,  18)
MUTED = RGBColor(105, 120, 158)

def dk(c, f=0.35):
    return RGBColor(min(255,int(c[0]*f)), min(255,int(c[1]*f)), min(255,int(c[2]*f)))

def hx(c):
    return f'{c[0]:02X}{c[1]:02X}{c[2]:02X}'

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ─── CORE ────────────────────────────────────────────
def ns():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BG
    return sl

def R(sl, l, t, w, h, fill=CARD, line=None):
    sh = sl.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = fill if line is None else line
    return sh

def Ov(sl, l, t, w, h, fill=CARD, line=None):
    sh = sl.shapes.add_shape(9, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = fill if line is None else line
    return sh

def GR(sl, l, t, w, h, c1, c2, ang=135, lc=None):
    sh = sl.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c1
    sh.line.color.rgb = c1 if lc is None else lc
    spPr = sh._element.find(qn('p:spPr'))
    gx = (f'<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
          f'<a:gsLst>'
          f'<a:gs pos="0"><a:srgbClr val="{hx(c1)}"/></a:gs>'
          f'<a:gs pos="100000"><a:srgbClr val="{hx(c2)}"/></a:gs>'
          f'</a:gsLst>'
          f'<a:lin ang="{ang*60000}" scaled="0"/>'
          f'</a:gradFill>')
    gf = etree.fromstring(gx)
    sol = spPr.find(qn('a:solidFill'))
    if sol is not None: spPr.remove(sol)
    spPr.append(gf)
    return sh

def GR3(sl, l, t, w, h, c1, c2, c3, ang=135):
    sh = sl.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c1; sh.line.color.rgb = c1
    spPr = sh._element.find(qn('p:spPr'))
    gx = (f'<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
          f'<a:gsLst>'
          f'<a:gs pos="0"><a:srgbClr val="{hx(c1)}"/></a:gs>'
          f'<a:gs pos="50000"><a:srgbClr val="{hx(c2)}"/></a:gs>'
          f'<a:gs pos="100000"><a:srgbClr val="{hx(c3)}"/></a:gs>'
          f'</a:gsLst>'
          f'<a:lin ang="{ang*60000}" scaled="0"/>'
          f'</a:gradFill>')
    gf = etree.fromstring(gx)
    sol = spPr.find(qn('a:solidFill'))
    if sol is not None: spPr.remove(sol)
    spPr.append(gf)
    return sh

def T(sl, text, l, t, w, h, sz=14, col=WHT, bold=False,
      italic=False, align=PP_ALIGN.LEFT, wrap=True):
    bx = sl.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.color.rgb = col
    r.font.bold = bold; r.font.italic = italic
    return bx

# ─── FILM STRIP ──────────────────────────────────────
def film_strip(sl, y, h=Inches(0.3)):
    R(sl, Inches(0), y, Inches(13.33), h, FILM)
    hw, hh = Inches(0.14), Inches(0.14)
    hy = y + (h - hh) / 2
    for i in range(44):
        hx_p = Inches(0.12) + i * Inches(0.3)
        sh = sl.shapes.add_shape(9, hx_p, hy, hw, hh)
        sh.fill.solid(); sh.fill.fore_color.rgb = BG
        sh.line.color.rgb = RGBColor(20, 20, 55)

# ─── VISUAL: ANIMATION FRAME SEQUENCE ────────────────
def frame_seq(sl, x, y, n=5, fw=Inches(0.82), fh=Inches(1.05), accent=CYAN):
    """Visual showing still frames becoming animation"""
    gap = Inches(0.08)
    stickfigs = ["—", "\\", "|", "/", "—"]   # simple motion hint
    for i in range(n):
        fx = x + i * (fw + gap)
        alpha_f = 0.25 + (i / max(n-1,1)) * 0.75
        border_c = RGBColor(int(accent[0]*alpha_f), int(accent[1]*alpha_f), int(accent[2]*alpha_f))
        # Frame box
        R(sl, fx, y, fw, fh, dk(accent, 0.12), border_c)
        # Inner darker area (the "image")
        R(sl, fx+Inches(0.05), y+Inches(0.05), fw-Inches(0.1), fh-Inches(0.28),
          dk(accent, 0.07))
        # Simple figure hint text
        T(sl, "●", fx, y+Inches(0.15), fw, Inches(0.38),
          sz=14, col=border_c, align=PP_ALIGN.CENTER, bold=True)
        # Frame label
        T(sl, f"f{i+1}", fx, y+fh-Inches(0.25), fw, Inches(0.22),
          sz=8, col=border_c, align=PP_ALIGN.CENTER, bold=True)
        if i < n-1:
            T(sl, "▶", fx+fw, y+fh/2-Inches(0.13), gap+Inches(0.02), Inches(0.26),
              sz=7, col=accent, align=PP_ALIGN.CENTER)
    T(sl, "Still frames  →  Illusion of Motion  →  ANIMATION",
      x, y+fh+Inches(0.06), n*(fw+gap), Inches(0.28),
      sz=9, col=MUTED, italic=True, align=PP_ALIGN.CENTER)

# ─── VISUAL: PLAY BUTTON ─────────────────────────────
def play_btn(sl, cx, cy, sz, color):
    Ov(sl, cx-sz/2, cy-sz/2, sz, sz, dk(color,0.2), color)
    tri = sl.shapes.add_shape(5, cx-sz*0.12, cy-sz*0.23, sz*0.32, sz*0.46)
    tri.fill.solid(); tri.fill.fore_color.rgb = color; tri.line.color.rgb = color

# ─── VISUAL: HEX / CIRCLE BADGE ──────────────────────
def badge(sl, cx, cy, sz, color, label, sublabel=""):
    Ov(sl, cx-sz/2, cy-sz/2, sz, sz, dk(color,0.22), color)
    Ov(sl, cx-sz*0.35, cy-sz*0.35, sz*0.7, sz*0.7, dk(color,0.12), dk(color,0.5))
    T(sl, label, cx-sz/2, cy-sz*0.32, sz, sz*0.62,
      sz=int(sz/Inches(1)*10), col=color, bold=True, align=PP_ALIGN.CENTER)
    if sublabel:
        T(sl, sublabel, cx-sz/2, cy+sz*0.15, sz, sz*0.3,
          sz=int(sz/Inches(1)*6), col=LT, align=PP_ALIGN.CENTER)

# ─── HEADER ──────────────────────────────────────────
def hdr(sl, num, title, accent):
    film_strip(sl, Inches(0))
    GR(sl, Inches(0), Inches(0.3), Inches(13.33), Inches(1.05), accent, dk(accent,0.42), 0)
    T(sl, num, Inches(10.6), Inches(-0.08), Inches(2.7), Inches(1.55),
      sz=92, col=dk(accent,0.3), bold=True, align=PP_ALIGN.RIGHT)
    T(sl, title, Inches(0.45), Inches(0.36), Inches(11.0), Inches(0.9),
      sz=27, col=WHT, bold=True)

def ibar(sl, text, accent, y=None):
    y = y or Inches(1.38)
    GR(sl, Inches(0.35), y, Inches(12.63), Inches(0.72), dk(accent,0.24), dk(accent,0.38), 0)
    GR(sl, Inches(0.35), y, Inches(0.2), Inches(0.72), accent, dk(accent))
    T(sl, text, Inches(0.66), y+Inches(0.07), Inches(12.1), Inches(0.6),
      sz=11.5, col=LT, italic=True)

# ─── CARD LAYOUTS ────────────────────────────────────
def card4(sl, cards, sy=Inches(1.38)):
    cw, ch = Inches(6.22), Inches(2.88)
    pos = [(Inches(0.35), sy), (Inches(6.68), sy),
           (Inches(0.35), sy+Inches(3.02)), (Inches(6.68), sy+Inches(3.02))]
    for i, (ac, ti, bo) in enumerate(cards[:4]):
        l, top = pos[i]
        R(sl, l, top, cw, ch, CARD)
        GR(sl, l, top, cw, Inches(0.58), ac, dk(ac,0.48), 0)
        GR(sl, l, top, Inches(0.2), ch, ac, dk(ac))
        # Accent dot
        Ov(sl, l+Inches(0.26), top+Inches(0.14), Inches(0.28), Inches(0.28), ac)
        T(sl, ti, l+Inches(0.65), top+Inches(0.07), cw-Inches(0.75), Inches(0.5),
          sz=13, col=ac, bold=True)
        T(sl, bo, l+Inches(0.28), top+Inches(0.62), cw-Inches(0.38), ch-Inches(0.72),
          sz=10.5, col=LT)

def card3(sl, cards, sy=Inches(1.38)):
    cw, ch = Inches(4.28), Inches(6.0)
    for i, (ac, ti, bo) in enumerate(cards[:3]):
        l = Inches(0.32) + i*Inches(4.45)
        R(sl, l, sy, cw, ch, CARD)
        GR(sl, l, sy, cw, Inches(0.68), ac, dk(ac,0.48), 0)
        GR(sl, l, sy, Inches(0.2), ch, ac, dk(ac))
        T(sl, ti, l+Inches(0.3), sy+Inches(0.08), cw-Inches(0.4), Inches(0.62),
          sz=14, col=DARK, bold=True, align=PP_ALIGN.CENTER)
        T(sl, bo, l+Inches(0.28), sy+Inches(0.78), cw-Inches(0.38), ch-Inches(0.9),
          sz=11, col=LT)

def card4h(sl, cards, sy=Inches(2.15), ch=Inches(5.0)):
    cw = Inches(3.08)
    for i, (ac, ti, bo) in enumerate(cards[:4]):
        l = Inches(0.38) + i*Inches(3.24)
        R(sl, l, sy, cw, ch, CARD)
        GR(sl, l, sy, cw, Inches(0.62), ac, dk(ac,0.48), 0)
        GR(sl, l, sy, Inches(0.18), ch, ac, dk(ac))
        T(sl, ti, l+Inches(0.28), sy+Inches(0.08), cw-Inches(0.36), Inches(0.56),
          sz=12, col=DARK, bold=True, align=PP_ALIGN.CENTER)
        T(sl, bo, l+Inches(0.26), sy+Inches(0.72), cw-Inches(0.36), ch-Inches(0.84),
          sz=10.5, col=LT)

# ════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════
s1 = ns()
GR3(s1, Inches(0), Inches(0), Inches(13.33), Inches(7.5),
    RGBColor(3,3,16), RGBColor(8,5,32), RGBColor(4,4,18), 135)

film_strip(s1, Inches(0))
film_strip(s1, Inches(7.2))

# Decorative glow circles
for x,y,sz,c in [
    (Inches(10.2), Inches(-1.0), Inches(4.5), dk(CYAN,0.14)),
    (Inches(11.8), Inches(5.0),  Inches(3.0), dk(PURP,0.14)),
    (Inches(-0.8), Inches(4.5),  Inches(3.2), dk(ORG,0.11)),
    (Inches(0.2),  Inches(-0.6), Inches(2.2), dk(CYAN,0.09)),
]:
    sh = s1.shapes.add_shape(9, x, y, sz, sz)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.color.rgb = c

# Accent bars
GR(s1, Inches(0), Inches(0.3), Inches(13.33), Inches(0.18), CYAN, PURP, 0)
GR(s1, Inches(0), Inches(7.02), Inches(13.33), Inches(0.18), PURP, ORG, 0)

# Visual: frame sequence top-right corner
frame_seq(s1, Inches(8.5), Inches(0.6), n=4, fw=Inches(0.92), fh=Inches(1.2), accent=CYAN)

# Watermark
T(s1, "ANIM", Inches(-0.4), Inches(0.2), Inches(9), Inches(5.2),
  sz=195, col=RGBColor(9,9,38), bold=True)

# Main title
T(s1, "ANIMATION", Inches(0.5), Inches(0.95), Inches(12.33), Inches(2.3),
  sz=88, col=CYAN, bold=True, align=PP_ALIGN.CENTER)

# Subtitle
GR(s1, Inches(1.8), Inches(3.12), Inches(9.73), Inches(0.8), dk(CYAN,0.26), dk(PURP,0.26), 0)
T(s1, "The Art & Science of Bringing Images to Life",
  Inches(1.8), Inches(3.14), Inches(9.73), Inches(0.76),
  sz=19, col=LT, italic=True, align=PP_ALIGN.CENTER)

GR(s1, Inches(3.2), Inches(4.05), Inches(6.93), Inches(0.07), CYAN, PURP, 0)

T(s1, "Graphics & Multimedia  •  SMIT  •  Group 3",
  Inches(1), Inches(4.18), Inches(11.33), Inches(0.55),
  sz=15, col=ORG, bold=True, align=PP_ALIGN.CENTER)

# Play button + topics badge
play_btn(s1, Inches(5.5), Inches(4.88), Inches(0.62), CYAN)
GR(s1, Inches(5.82), Inches(4.58), Inches(3.85), Inches(0.72), dk(CYAN,0.38), dk(PURP,0.38), 0)
T(s1, "Topics 1 — 9", Inches(5.82), Inches(4.58), Inches(3.85), Inches(0.72),
  sz=17, col=WHT, bold=True, align=PP_ALIGN.CENTER)

# Member cards
members = [
    ("Hashir Junaid",      "2023F-BCS-358"),
    ("Taha Haider",        "2023F-BCS-079"),
    ("Abdul Rahman Baig",  "2023F-BCS-100"),
]
for i,(name,roll) in enumerate(members):
    x = Inches(0.85) + i*Inches(3.9)
    GR(s1, x, Inches(5.62), Inches(3.65), Inches(1.0), CARD, CARD2, 0)
    Ov(s1, x+Inches(0.15), Inches(5.78), Inches(0.28), Inches(0.28), CYAN)
    T(s1, name, x+Inches(0.55), Inches(5.68), Inches(3.05), Inches(0.46),
      sz=13, col=WHT, bold=True)
    T(s1, roll, x+Inches(0.55), Inches(6.1), Inches(3.05), Inches(0.42),
      sz=11, col=MUTED)

# ════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════
s2 = ns()
film_strip(s2, Inches(0))
GR(s2, Inches(0), Inches(0.3), Inches(13.33), Inches(1.1), CYAN, PURP, 0)
T(s2, "WHAT WE'LL COVER TODAY", Inches(0.45), Inches(0.36), Inches(12), Inches(0.88),
  sz=33, col=WHT, bold=True)
T(s2, "9 Topics  •  Animation: Complete Foundation  •  Group 3",
  Inches(0.45), Inches(1.1), Inches(12), Inches(0.32), sz=11, col=dk(CYAN,1.8), italic=True)

agenda = [
    ("01","Introduction to Animation",   CYAN),
    ("02","Usage of Animation",          PURP),
    ("03","Animation Techniques",        ORG),
    ("04","Types of Animation",          GRN),
    ("05","Computer Animation",          PINK),
    ("06","Levels of Computer Animation",YEL),
    ("07","12 Basic Principles",         CYAN),
    ("08","Animation Space",             PURP),
    ("09","Animation Process",           ORG),
]
col_x = [Inches(0.32), Inches(4.62), Inches(8.95)]
row_y = [Inches(1.52), Inches(2.85), Inches(4.18)]
for i,(num,topic,ac) in enumerate(agenda):
    row, col = divmod(i, 3)
    x, y = col_x[col], row_y[row]
    R(s2, x, y, Inches(3.9), Inches(1.15), CARD)
    GR(s2, x, y, Inches(0.72), Inches(1.15), ac, dk(ac,0.5), 90)
    Ov(s2, x+Inches(0.22), y+Inches(0.4), Inches(0.28), Inches(0.28), dk(ac,0.4))
    T(s2, num, x, y, Inches(0.72), Inches(1.15),
      sz=21, col=DARK, bold=True, align=PP_ALIGN.CENTER)
    T(s2, topic, x+Inches(0.82), y+Inches(0.27), Inches(2.98), Inches(0.68),
      sz=13.5, col=WHT)

GR(s2, Inches(0), Inches(5.38), Inches(13.33), Inches(2.12), BG, RGBColor(8,8,38), 90)
T(s2, "From hand-drawn cels to AI-powered 3D — 9 complete topics await.",
  Inches(1), Inches(5.7), Inches(11.33), Inches(0.7),
  sz=20, col=MUTED, italic=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# SLIDE 3 — INTRODUCTION TO ANIMATION
# ════════════════════════════════════════════════════
s3 = ns()
hdr(s3, "01", "INTRODUCTION TO ANIMATION", CYAN)

# Visual: frame sequence illustration
frame_seq(s3, Inches(7.55), Inches(1.42), n=5, fw=Inches(0.98), fh=Inches(1.2), accent=CYAN)

# Definition box (left side)
GR(s3, Inches(0.35), Inches(1.42), Inches(6.85), Inches(1.05), dk(CYAN,0.26), dk(CYAN,0.14), 0)
GR(s3, Inches(0.35), Inches(1.42), Inches(0.2), Inches(1.05), CYAN, dk(CYAN))
T(s3,
  '"Animation: Rapidly displaying STATIC images to create the ILLUSION of\n'
  ' MOVEMENT — exploiting the brain\'s own perception against itself."',
  Inches(0.65), Inches(1.48), Inches(6.52), Inches(0.98),
  sz=12.5, col=WHT, italic=True)

cards = [
    (CYAN,  "DEFINITION",
     "From Latin 'anima' — to give soul\n\n"
     "• Series of still frames shown rapidly\n"
     "• 24 fps = cinema gold standard\n"
     "• 30 fps = TV broadcast\n"
     "• 60 fps = modern gaming\n"
     "• Below 12 fps = visible flicker\n\n"
     "HISTORY:\n"
     "First animation: Eadweard Muybridge\nhorse gallop sequence (1878)"),
    (PURP,  "PERSISTENCE OF VISION",
     "Physical effect in the human eye\n\n"
     "• Retina keeps image ~50ms after\n  stimulus disappears\n"
     "• Afterimage overlaps next frame\n"
     "• Creates a seamless visual flow\n"
     "• Discovered: Peter Roget (1824)\n\n"
     "WHY: Brain cannot distinguish\nindividual frames at 24+ fps\n"
     "— sees continuous motion instead"),
    (ORG,   "PHI PHENOMENON",
     "Purely PSYCHOLOGICAL motion effect\n\n"
     "• Brain infers motion between two\n  stationary objects in sequence\n"
     "• Does NOT need retinal retention\n"
     "• Discovered: Wertheimer (1912)\n\n"
     "CLASSIC EXAMPLES:\n"
     "• Flip-book animation\n"
     "• Neon sign arrows 'flowing'\n"
     "• Traffic light direction arrows"),
    (GRN,   "ILLUSION OF MOVEMENT",
     "Combined result of BOTH effects\n\n"
     "• Persistence of Vision alone = NOT enough\n"
     "• Phi Phenomenon alone = NOT enough\n"
     "• BOTH together = seamless motion!\n\n"
     "The brain WANTS to see motion.\nAnimators exploit this biological\n'glitch' in human perception.\n\n"
     "This is WHY animation feels REAL\neven though it is 100% artificial"),
]
card4(s3, cards, sy=Inches(2.52))

# ════════════════════════════════════════════════════
# SLIDE 4 — USAGE OF ANIMATION
# ════════════════════════════════════════════════════
s4 = ns()
hdr(s4, "02", "USAGE OF ANIMATION", PURP)

# Category badges at top
categories = [
    (CYAN, "ARTISTIC",      "Fine Art / Music Videos"),
    (ORG,  "STORYTELLING",  "Films / Games / Series"),
    (GRN,  "SCIENTIFIC",    "Research / Medicine / NASA"),
    (PINK, "INSTRUCTIONAL", "Education / Training / Sims"),
]
for i,(ac,lbl,sub) in enumerate(categories):
    x = Inches(0.35) + i*Inches(3.25)
    GR(s4, x, Inches(1.35), Inches(3.15), Inches(0.72), dk(ac,0.25), dk(ac,0.42), 0)
    GR(s4, x, Inches(1.35), Inches(0.2), Inches(0.72), ac, dk(ac))
    T(s4, lbl, x+Inches(0.28), Inches(1.4), Inches(2.82), Inches(0.36),
      sz=12, col=ac, bold=True)
    T(s4, sub, x+Inches(0.28), Inches(1.74), Inches(2.82), Inches(0.28),
      sz=9, col=MUTED)

cards = [
    (CYAN, "ARTISTIC PURPOSES",
     "• Fine art & experimental films\n"
     "• Music video production\n"
     "• Abstract visual compositions\n"
     "• Title sequence design\n"
     "• Gallery art installations\n\n"
     "ICONIC EXAMPLES:\n"
     "  Disney Fantasia (1940)\n"
     "  Yellow Submarine (1968)\n"
     "  Spider-Man: Into the Spider-Verse\n"
     "  (2018) — Oscar winner"),
    (ORG, "STORYTELLING",
     "• Narrative feature films\n"
     "• TV series & streaming shows\n"
     "• Video game cutscenes\n"
     "• Interactive web comics\n"
     "• Character-driven dramas\n\n"
     "ICONIC EXAMPLES:\n"
     "  Pixar Toy Story series\n"
     "  Studio Ghibli classics\n"
     "  Arcane — Netflix (2021)\n"
     "  Attack on Titan"),
    (GRN, "SCIENTIFIC VISUALIZATION",
     "• Human anatomy explainers\n"
     "• Weather & climate models\n"
     "• DNA & molecular modeling\n"
     "• Physics fluid simulations\n"
     "• Engineering CAD previews\n\n"
     "ICONIC EXAMPLES:\n"
     "  NASA mission visualizations\n"
     "  Medical surgical training\n"
     "  Climate change simulations\n"
     "  Molecular biology videos"),
    (PINK, "INSTRUCTIONAL",
     "• Educational explainer videos\n"
     "• Military & flight simulators\n"
     "• Safety procedure demos\n"
     "• Product assembly guides\n"
     "• Corporate onboarding modules\n\n"
     "ICONIC EXAMPLES:\n"
     "  Khan Academy lessons\n"
     "  Boeing pilot simulators\n"
     "  IKEA assembly guides\n"
     "  YouTube tutorial channels"),
]
card4(s4, cards, sy=Inches(2.12))

# ════════════════════════════════════════════════════
# SLIDE 5 — ANIMATION TECHNIQUES
# ════════════════════════════════════════════════════
s5 = ns()
hdr(s5, "03", "ANIMATION TECHNIQUES", ORG)

# Visual: technique evolution arrow
GR(s5, Inches(0.35), Inches(1.38), Inches(12.63), Inches(0.48), dk(ORG,0.2), dk(ORG,0.12), 0)
years = [("1910s – 1990s", "CEL ANIMATION", CYAN),
         ("1970s – Now",   "PATH ANIMATION", ORG),
         ("1980s – Now",   "COMPUTER ANIM", GRN)]
for i,(era,lbl,ac) in enumerate(years):
    x = Inches(0.55) + i*Inches(4.35)
    GR(s5, x, Inches(1.38), Inches(4.1), Inches(0.48), dk(ac,0.32), dk(ac,0.22), 0)
    T(s5, f"{lbl}  •  {era}", x+Inches(0.1), Inches(1.42), Inches(3.9), Inches(0.4),
      sz=10, col=ac, bold=True)

cards = [
    (CYAN, "CEL ANIMATION",
     "Traditional hand-drawn technique\n\n"
     "• Drawings on transparent celluloid sheets\n"
     "• Each sheet = one frame of animation\n"
     "• Cels layered over static backgrounds\n"
     "• Allows background reuse (cost saving)\n"
     "• Photographed frame-by-frame on film\n\n"
     "WORKFLOW:\n"
     "  Story → Storyboard → Key Animation\n"
     "  → In-Betweening → Ink → Paint → Photo\n\n"
     "STUDIOS: Disney, Hanna-Barbera,\n"
     "          Warner Bros, Toei Animation"),
    (ORG, "PATH ANIMATION",
     "Object follows a predefined trajectory\n\n"
     "• Animator defines the path shape first\n"
     "• Object moves along it automatically\n"
     "• Supports straight lines & Bezier curves\n"
     "• Object can resize/reshape along path\n"
     "• Easing controls speed variation\n\n"
     "COMMON USES:\n"
     "  Logo reveals, camera fly-throughs,\n"
     "  UI transitions, game projectiles\n\n"
     "TOOLS: After Effects, CSS, SVG paths,\n"
     "        Unity, Unreal Engine"),
    (GRN, "COMPUTER ANIMATION",
     "Digital creation of moving images\n\n"
     "• Can simulate both Cel & Path techniques\n"
     "• PLUS: physics, particles, simulation\n"
     "• 2D digital, 3D CG, or hybrid mix\n"
     "• Real-time OR pre-rendered output\n"
     "• Procedural & AI-assisted animation\n\n"
     "ADVANTAGES:\n"
     "  Unlimited camera angles\n"
     "  Easy to iterate & fix errors\n"
     "  Physics simulation built-in\n"
     "  Automatic in-betweening\n\n"
     "STUDIOS: Pixar, DreamWorks, ILM"),
]
card3(s5, cards, sy=Inches(1.92))

# Evolution arrows
for i in range(2):
    ax = Inches(4.46) + i*Inches(4.45)
    GR(s5, ax, Inches(4.0), Inches(0.22), Inches(0.7), ORG, dk(ORG), 0)
    T(s5, "▶", ax-Inches(0.06), Inches(3.9), Inches(0.38), Inches(0.9),
      sz=20, col=ORG, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# SLIDE 6 — CEL ANIMATION DEEP DIVE
# ════════════════════════════════════════════════════
s6 = ns()
hdr(s6, "04", "TYPES OF ANIMATION — CEL ANIMATION", CYAN)
ibar(s6,
     "Named after CELLULOID — transparent plastic sheets invented in 1869. "
     "Raoul Barré & John Bray developed the cel technique ~1914. Dominated animation for 80 years.", CYAN)

# Visual: frame/cel stack illustration
for i in range(4):
    off = i * Inches(0.12)
    alpha = 0.15 + i*0.2
    c = RGBColor(int(CYAN[0]*alpha), int(CYAN[1]*alpha), int(CYAN[2]*alpha))
    R(s6, Inches(11.6)-off, Inches(2.2)+off, Inches(1.45), Inches(1.8), CARD, c)
T(s6, "CEL\nSTACK", Inches(11.5), Inches(2.1), Inches(1.6), Inches(0.55),
  sz=9, col=MUTED, align=PP_ALIGN.CENTER)

cards = [
    (CYAN, "CEL ANIMATION",
     "• Transparent sheet per frame of film\n"
     "• Character on cel; background separate\n"
     "• Cels stacked = final composite frame\n"
     "• Background reuse saves huge cost\n"
     "• Colors applied to BACK of cel sheet\n"
     "• 12–24 unique cels per second of film\n\n"
     "USED IN:\n"
     "Tom & Jerry, Looney Tunes,\n"
     "Snow White, Bambi, Akira, Pinocchio"),
    (ORG, "KEYFRAMES",
     "• The 'anchor' frames of an action\n"
     "• Show most extreme or critical poses\n"
     "• Drawn by SENIOR animators\n"
     "• Define the timing & structure\n"
     "• Think of them as checkpoints\n\n"
     "EXAMPLE:\n"
     "Frame 1 = character standing still\n"
     "Frame 12 = peak of jump (keyframe)\n"
     "Frame 24 = character landed\n\n"
     "Modern software uses identical\nkeyframe concept in timelines"),
    (GRN, "TWEENING (In-Betweening)",
     "• Frames drawn BETWEEN keyframes\n"
     "• Creates smooth motion transition\n"
     "• Done by JUNIOR animators\n"
     "• More tweens = smoother motion\n"
     "• At 24fps: ~22 tweens per action\n\n"
     "MODERN EQUIVALENT:\n"
     "Software AUTO-tweens between\nkeyframes — no manual drawing!\n\n"
     "Tools: Adobe Animate, Flash,\n"
     "Toon Boom Harmony, After Effects"),
    (PURP, "PENCIL TEST",
     "• Quick pre-production validation test\n"
     "• Rough pencil sketches filmed first\n"
     "• Played back to check timing & flow\n"
     "• Catches errors BEFORE expensive\n  inking & painting stages\n"
     "• Saves significant time and money\n\n"
     "MODERN EQUIVALENT:\n"
     "Animatic / rough cut preview\nplayed in editing software\n\n"
     "Standard practice at EVERY\nmajor animation studio worldwide"),
]
card4(s6, cards, sy=Inches(2.15))

# ════════════════════════════════════════════════════
# SLIDE 7 — PATH ANIMATION DEEP DIVE
# ════════════════════════════════════════════════════
s7 = ns()
hdr(s7, "04", "TYPES OF ANIMATION — PATH ANIMATION", ORG)
ibar(s7,
     "PATH ANIMATION: Object follows a predefined trajectory. "
     "Animator draws the path; software handles movement along it automatically.", ORG)

# Visual: path diagram (left visual panel)
R(s7, Inches(0.35), Inches(2.12), Inches(6.1), Inches(5.1), CARD)
GR(s7, Inches(0.35), Inches(2.12), Inches(6.1), Inches(0.62), ORG, dk(ORG,0.5), 0)
GR(s7, Inches(0.35), Inches(2.12), Inches(0.2), Inches(5.1), ORG, dk(ORG))
T(s7, "KEY CONCEPTS", Inches(0.65), Inches(2.17), Inches(5.7), Inches(0.52),
  sz=15, col=DARK, bold=True)

# Visual: curved path using circles as waypoints
path_pts_visual = [
    (Inches(1.2), Inches(3.0)),
    (Inches(2.8), Inches(2.5)),
    (Inches(4.0), Inches(3.4)),
    (Inches(5.5), Inches(2.8)),
]
for i,(px,py) in enumerate(path_pts_visual):
    sz = Inches(0.22) if i not in [0, len(path_pts_visual)-1] else Inches(0.3)
    c = ORG if i in [0, len(path_pts_visual)-1] else CYAN
    Ov(s7, px-sz/2, py-sz/2, sz, sz, c)
    if i < len(path_pts_visual)-1:
        nx, ny = path_pts_visual[i+1]
        # Dashed line hint (thin rect between points)
        mid_x = (px+nx)/2 - Inches(0.02)
        mid_y = (py+ny)/2 - Inches(0.01)
        R(s7, mid_x, mid_y, Inches(0.04), Inches(0.04), MUTED)
T(s7, "START", path_pts_visual[0][0]-Inches(0.4), path_pts_visual[0][1]+Inches(0.1),
  Inches(0.7), Inches(0.22), sz=8, col=ORG, bold=True)
T(s7, "END", path_pts_visual[-1][0]-Inches(0.1), path_pts_visual[-1][1]+Inches(0.1),
  Inches(0.5), Inches(0.22), sz=8, col=ORG, bold=True)
T(s7, "Bezier curve path\n(object follows this route)",
  Inches(1.0), Inches(3.65), Inches(4.8), Inches(0.45),
  sz=9.5, col=MUTED, italic=True, align=PP_ALIGN.CENTER)

pts = [
    (ORG,  "PREDETERMINED PATH",
     "Route defined BEFORE animation. Like a roller-coaster — object follows exactly what's laid out."),
    (CYAN, "STRAIGHT LINES",
     "Linear movement for mechanical objects, text reveals, UI animations."),
    (GRN,  "BEZIER CURVES",
     "Smooth organic paths via control points. Natural feel for balls, flying objects, camera."),
    (PINK, "RESIZING & RESHAPING",
     "Object scales or deforms as it travels the path — adds dynamism (e.g. ball squash on bounce)."),
]
yo = Inches(4.22)
for ac,ti,de in pts:
    GR(s7, Inches(0.58), yo, Inches(0.12), Inches(0.8), ac, dk(ac))
    T(s7, ti, Inches(0.82), yo, Inches(5.5), Inches(0.35), sz=11.5, col=ac, bold=True)
    T(s7, de, Inches(0.82), yo+Inches(0.34), Inches(5.5), Inches(0.5), sz=10.5, col=LT)
    yo += Inches(0.92)

# Right panel: How it works
R(s7, Inches(6.78), Inches(2.12), Inches(6.17), Inches(5.1), CARD)
GR(s7, Inches(6.78), Inches(2.12), Inches(6.17), Inches(0.62), ORG, dk(ORG,0.5), 0)
GR(s7, Inches(6.78), Inches(2.12), Inches(0.2), Inches(5.1), ORG, dk(ORG))
T(s7, "HOW PATH ANIMATION WORKS", Inches(7.08), Inches(2.17), Inches(5.77), Inches(0.52),
  sz=15, col=DARK, bold=True)

steps6 = [
    (CYAN,  "1","Draw the PATH — straight line, Bezier curve, or freehand"),
    (ORG,   "2","Set START and END points along the path"),
    (GRN,   "3","Define DURATION — how long the journey takes"),
    (PURP,  "4","Apply EASING — slow in / slow out for natural feel"),
    (PINK,  "5","Set ORIENTATION — does object face direction of travel?"),
    (YEL,   "6","Preview, adjust, repeat — render when satisfied"),
]
yo = Inches(2.86)
for ac,n,desc in steps6:
    GR(s7, Inches(6.92), yo, Inches(0.55), Inches(0.62), ac, dk(ac))
    T(s7, n, Inches(6.92), yo, Inches(0.55), Inches(0.62),
      sz=17, col=DARK, bold=True, align=PP_ALIGN.CENTER)
    T(s7, desc, Inches(7.57), yo+Inches(0.1), Inches(5.28), Inches(0.55), sz=11.5, col=LT)
    yo += Inches(0.76)

GR(s7, Inches(0.35), Inches(7.1), Inches(12.63), Inches(0.36), dk(GRN,0.25), dk(GRN,0.38), 0)
T(s7, "Used in: CSS animations  •  After Effects motion paths  •  Unity NavMesh  •  SVG paths  •  Game projectiles",
  Inches(0.55), Inches(7.12), Inches(12.2), Inches(0.32), sz=11, col=GRN)

# ════════════════════════════════════════════════════
# SLIDE 8 — COMPUTER ANIMATION (CEL vs COMPUTER)
# ════════════════════════════════════════════════════
s8 = ns()
hdr(s8, "05", "COMPUTER ANIMATION", GRN)
ibar(s8,
     "Computer animation uses software to create motion — simulating traditional techniques "
     "or producing entirely new effects impossible by hand.", GRN)

T(s8, "CEL  vs  COMPUTER  ANIMATION",
  Inches(0.5), Inches(2.1), Inches(12.33), Inches(0.6),
  sz=20, col=YEL, bold=True, align=PP_ALIGN.CENTER)

GR(s8, Inches(0.35), Inches(2.75), Inches(5.65), Inches(0.52), CYAN, dk(CYAN,0.5), 0)
T(s8, "CEL ANIMATION", Inches(0.35), Inches(2.75), Inches(5.65), Inches(0.52),
  sz=14, col=DARK, bold=True, align=PP_ALIGN.CENTER)
GR(s8, Inches(7.35), Inches(2.75), Inches(5.65), Inches(0.52), GRN, dk(GRN,0.5), 0)
T(s8, "COMPUTER ANIMATION", Inches(7.35), Inches(2.75), Inches(5.65), Inches(0.52),
  sz=14, col=DARK, bold=True, align=PP_ALIGN.CENTER)

# VS badge
GR(s8, Inches(6.03), Inches(2.62), Inches(1.27), Inches(1.05), ORG, dk(ORG,0.5), 90)
T(s8, "VS", Inches(6.03), Inches(2.62), Inches(1.27), Inches(1.05),
  sz=26, col=WHT, bold=True, align=PP_ALIGN.CENTER)

rows = [
    ("Hand-drawn on physical celluloid sheets", "Created entirely within digital software"),
    ("Every frame drawn manually — very slow", "Automation handles tweens, physics, rendering"),
    ("Physical storage: film cans, warehouses", "Digital files — backup & share instantly"),
    ("Expensive: cels, paint, camera, lab fees", "One-time software cost; no consumables"),
    ("Difficult to modify any finished frame", "Copy, adjust, modify any frame instantly"),
    ("Organic, hand-crafted visual warmth", "Any style: toon to photorealistic"),
    ("Months per 1 minute of finished film", "Real-time previews; faster iterations"),
]
for i,(cel,comp) in enumerate(rows):
    y = Inches(3.32) + i*Inches(0.56)
    bg = CARD if i%2==0 else CARD2
    R(s8, Inches(0.35), y, Inches(5.65), Inches(0.52), bg)
    R(s8, Inches(7.35), y, Inches(5.65), Inches(0.52), bg)
    T(s8, cel,  Inches(0.48),  y+Inches(0.07), Inches(5.38), Inches(0.44), sz=11, col=LT)
    T(s8, comp, Inches(7.48), y+Inches(0.07), Inches(5.38), Inches(0.44), sz=11, col=GRN)

# ════════════════════════════════════════════════════
# SLIDE 9 — FK / IK / MORPHING
# ════════════════════════════════════════════════════
s9 = ns()
hdr(s9, "05", "KINEMATICS, INVERSE KINEMATICS & MORPHING", GRN)

# Visual: joint chain diagram for FK
ov_positions_fk = [(Inches(1.25), Inches(2.0)), (Inches(1.72), Inches(2.68)),
                   (Inches(2.18), Inches(3.35)), (Inches(2.55), Inches(4.0))]
labels_fk = ["SHOULDER", "ELBOW", "WRIST", "FINGER"]
for i,(px,py) in enumerate(ov_positions_fk):
    Ov(s9, px, py, Inches(0.32), Inches(0.32), CYAN if i==0 else dk(CYAN, 0.4+i*0.18), CYAN)
    if i < len(ov_positions_fk)-1:
        # Line to next joint (as thin rect)
        nx, ny = ov_positions_fk[i+1]
        mid_x = (px+nx)/2
        mid_y = (py+ny)/2
        R(s9, mid_x, mid_y, Inches(0.06), Inches(0.06), MUTED)
    T(s9, labels_fk[i], px+Inches(0.38), py+Inches(0.04),
      Inches(1.0), Inches(0.25), sz=8, col=CYAN)

T(s9, "FK: Parent → Child", Inches(0.9), Inches(4.45), Inches(2.5), Inches(0.28),
  sz=9, col=MUTED, italic=True)

cards = [
    (CYAN, "FORWARD KINEMATICS (FK)",
     "Animating from PARENT → CHILD joint\n\n"
     "HOW IT WORKS:\n"
     "  Move SHOULDER →\n"
     "  Elbow follows automatically →\n"
     "  Wrist follows →\n"
     "  Finger follows\n\n"
     "ANALOGY: A puppet — control\nthe parent; all children follow.\n\n"
     "BEST FOR:\n"
     "  Mechanical/robotic arms\n"
     "  Waving & swinging motions\n"
     "  Tail & tentacle animation\n"
     "  Flag ripples in wind\n\n"
     "Tools: Maya rigs, Blender"),
    (ORG, "INVERSE KINEMATICS (IK)",
     "Move END POINT; chain auto-adjusts\n\n"
     "HOW IT WORKS:\n"
     "  Move HAND position →\n"
     "  Software calculates elbow →\n"
     "  And shoulder automatically\n\n"
     "ANALOGY: Reaching for a cup —\nyou aim the hand; your body\nfigures out the joints itself.\n\n"
     "BEST FOR:\n"
     "  Character walking & running\n"
     "  Reaching & grabbing actions\n"
     "  Real-time game characters\n"
     "  Interactive VR animation"),
    (PURP, "MORPHING",
     "Smooth shape/image transformation\n\n"
     "HOW IT WORKS:\n"
     "  Map points on source image\n"
     "  Map matching points on target\n"
     "  Interpolate all points together\n"
     "  Creates seamless in-betweens\n\n"
     "TYPES:\n"
     "  2D: image-to-image morph\n"
     "  3D: mesh deformation morph\n"
     "  Shape: vector path morphing\n\n"
     "FAMOUS EXAMPLES:\n"
     "  MJ 'Black or White' (1991)\n"
     "  T-1000 — Terminator 2\n"
     "  Logo reveal animations"),
]
card3(s9, cards, sy=Inches(1.38))

# ════════════════════════════════════════════════════
# SLIDE 10 — LEVELS OF COMPUTER ANIMATION
# ════════════════════════════════════════════════════
s10 = ns()
hdr(s10, "06", "LEVELS OF COMPUTER ANIMATION", YEL)

# Visual: level progression bar
for i,ac in enumerate([CYAN, ORG, PURP]):
    x = Inches(0.35) + i*Inches(4.35)
    GR(s10, x, Inches(1.38), Inches(4.22), Inches(0.42), dk(ac,0.25), dk(ac,0.38), 0)
    stars = "★" * (i+1) + "☆" * (2-i)
    T(s10, stars, x, Inches(1.38), Inches(4.22), Inches(0.42),
      sz=15, col=ac, bold=True, align=PP_ALIGN.CENTER)

cards = [
    (CYAN, "BASIC  ★☆☆",
     "Entry-level animation tools\n\n"
     "WHAT YOU CREATE:\n"
     "  • Animated GIFs\n"
     "  • CSS / JS transitions\n"
     "  • PowerPoint animations\n"
     "  • Simple 2D sprites\n"
     "  • Basic logo animations\n\n"
     "TOOLS:\n"
     "  MS PowerPoint, Canva,\n"
     "  GIMP, basic CSS/JS\n\n"
     "SKILL: Beginner\n"
     "TIME TO LEARN: Days to weeks\n\n"
     "CAREER: Social media content,\nmarketing materials"),
    (ORG, "INTERMEDIATE  ★★☆",
     "Deeper control over animation\n\n"
     "WHAT YOU CREATE:\n"
     "  • 2D character animation\n"
     "  • Motion graphics\n"
     "  • Complex path animations\n"
     "  • Basic 3D object animation\n"
     "  • Game character movement\n\n"
     "TOOLS:\n"
     "  Adobe After Effects,\n"
     "  Adobe Animate,\n"
     "  Blender (basics), Spine\n\n"
     "SKILL: Intermediate\n"
     "TIME TO LEARN: Months\n\n"
     "CAREER: Motion designer,\n2D animator, game developer"),
    (PURP, "ADVANCED  ★★★",
     "Studio-level production quality\n\n"
     "WHAT YOU CREATE:\n"
     "  • Full CG feature films\n"
     "  • Realistic physics simulations\n"
     "  • Procedural animation\n"
     "  • Motion capture integration\n"
     "  • Real-time VFX in games\n"
     "  • AI-assisted animation\n\n"
     "TOOLS:\n"
     "  Autodesk Maya, Houdini,\n"
     "  Unreal Engine, Cinema 4D\n\n"
     "SKILL: Professional\n"
     "TIME TO LEARN: Years\n\n"
     "CAREER: VFX artist, 3D animator,\ntechnical director, TD"),
]
card3(s10, cards, sy=Inches(1.85))

# ════════════════════════════════════════════════════
# SLIDE 11 — 12 PRINCIPLES PART 1 (1-4)
# ════════════════════════════════════════════════════
s11 = ns()
hdr(s11, "07", "12 BASIC PRINCIPLES OF ANIMATION  •  Part 1 of 3", CYAN)
ibar(s11,
     "Developed by Disney masters Ollie Johnston & Frank Thomas (The Illusion of Life, 1981). "
     "Still the universal GOLD STANDARD — used by every professional animator on Earth.", CYAN)

# Visual: principle counter
for i in range(12):
    x = Inches(0.35) + i*Inches(0.98)
    c = CYAN if i < 4 else MUTED
    Ov(s11, x, Inches(2.04), Inches(0.5), Inches(0.5), dk(c,0.3) if i >= 4 else dk(c,0.5), c)
    T(s11, str(i+1), x, Inches(2.04), Inches(0.5), Inches(0.5),
      sz=9, col=c if i < 4 else MUTED, bold=(i<4), align=PP_ALIGN.CENTER)

cards = [
    (CYAN, "01  SQUASH & STRETCH",
     "Exaggerate deformation for weight & life\n\n"
     "• Object SQUASHES on impact\n"
     "• Object STRETCHES when moving fast\n"
     "• Golden rule: VOLUME stays constant!\n"
     "• Shows elasticity, mass, and life\n\n"
     "EXAMPLE:\n"
     "A rubber ball bounce —\nsquash on landing,\nstretch in the air.\n\n"
     "Most fundamental principle\nused in almost every shot"),
    (ORG, "02  ANTICIPATION",
     "Small opposite action BEFORE main action\n\n"
     "• Character crouches before jumping\n"
     "• Arm swings back before throwing\n"
     "• Boxer winds up before punching\n"
     "• Prepares the audience mentally\n\n"
     "EXAMPLE:\n"
     "Bugs Bunny ducking low before\na massive comedic leap forward.\n\n"
     "Without anticipation, actions\nfeel sudden and unconvincing"),
    (GRN, "03  STAGING",
     "Present ONE clear idea at a time\n\n"
     "• Use clear, readable silhouettes\n"
     "• Camera angle serves the story\n"
     "• Avoid distracting backgrounds\n"
     "• Direct the audience's eye\n\n"
     "EXAMPLE:\n"
     "Character in profile making a\ndramatic gesture — instantly\nreadable from any distance.\n\n"
     "Think: theater staging, not\ncrowded photography"),
    (PURP, "04  STRAIGHT AHEAD\n& POSE-TO-POSE",
     "Two approaches to animation:\n\n"
     "STRAIGHT AHEAD:\n"
     "• Draw frame by frame start → end\n"
     "• Spontaneous, fluid, surprising\n"
     "• Great for fire, water, chaos\n\n"
     "POSE-TO-POSE:\n"
     "• Draw key poses first\n"
     "• Fill in betweens after\n"
     "• Controlled, planned, consistent\n"
     "• Great for character acting\n\n"
     "Best work combines BOTH"),
]
card4(s11, cards, sy=Inches(2.58))

# ════════════════════════════════════════════════════
# SLIDE 12 — 12 PRINCIPLES PART 2 (5-8)
# ════════════════════════════════════════════════════
s12 = ns()
hdr(s12, "07", "12 BASIC PRINCIPLES OF ANIMATION  •  Part 2 of 3", ORG)
ibar(s12,
     "These principles govern how objects move AFTER the main action — "
     "the details that separate amateur animation from professional studio work.", ORG)

for i in range(12):
    x = Inches(0.35) + i*Inches(0.98)
    c = ORG if 4 <= i < 8 else MUTED
    Ov(s12, x, Inches(2.04), Inches(0.5), Inches(0.5), dk(c,0.3) if i < 4 or i >= 8 else dk(c,0.5), c)
    T(s12, str(i+1), x, Inches(2.04), Inches(0.5), Inches(0.5),
      sz=9, col=c if 4 <= i < 8 else MUTED, bold=(4<=i<8), align=PP_ALIGN.CENTER)

cards = [
    (CYAN, "05  FOLLOW THROUGH\n& OVERLAPPING ACTION",
     "FOLLOW THROUGH:\n"
     "Parts keep moving after\nmain action stops.\n"
     "• Hair swings after head stops\n"
     "• Coat settles after character stops\n\n"
     "OVERLAPPING ACTION:\n"
     "Different parts move at\ndifferent rates and with delay.\n"
     "• Arms lag behind the torso\n"
     "• Secondary elements trail\n\n"
     "Makes motion feel ORGANIC\nand physically believable"),
    (ORG, "06  EASE IN & EASE OUT\n(Slow In / Slow Out)",
     "Acceleration at start, deceleration\nat end of every movement.\n\n"
     "• Nothing in nature starts at full\n  speed or stops instantly\n"
     "• A car: slow → full speed → brake\n"
     "• Gives physical weight & real feel\n\n"
     "In software:\n"
     "  Bezier easing curves\n"
     "  Linear = robotic / artificial\n"
     "  Ease = natural / believable\n\n"
     "Most impactful principle\nfor 'weight' and feel in motion"),
    (GRN, "07  ARCS",
     "Natural motion follows ARC paths,\nnot straight lines.\n\n"
     "• Arms swing in arcs\n"
     "• Heads nod in arcs\n"
     "• Thrown objects: parabola\n"
     "• Even eyes dart in arcs\n"
     "• Running legs follow arcs\n\n"
     "EXCEPTION:\n"
     "Mechanical/robot movement is\nintentionally straight — use\nlinear paths for machines.\n\n"
     "Breaking arcs = instant 'robot'"),
    (PURP, "08  SECONDARY ACTION",
     "Supporting actions that ENRICH\nthe primary movement.\n\n"
     "EXAMPLES:\n"
     "  Walking (main) + arms swing\n"
     "  Running (main) + hair flows\n"
     "  Talking (main) + eyebrow moves\n\n"
     "CRITICAL RULE:\n"
     "Secondary must NEVER overpower\nthe main action — it supports,\nit does not steal focus.\n\n"
     "Adds richness and believability\nto every scene"),
]
card4(s12, cards, sy=Inches(2.58))

# ════════════════════════════════════════════════════
# SLIDE 13 — 12 PRINCIPLES PART 3 (9-12)
# ════════════════════════════════════════════════════
s13 = ns()
hdr(s13, "07", "12 BASIC PRINCIPLES OF ANIMATION  •  Part 3 of 3", PURP)
ibar(s13,
     "Walt Disney: 'The secret of animation is the exaggeration of the truth.' "
     "These final 4 principles deliver exactly that truth.", PURP)

for i in range(12):
    x = Inches(0.35) + i*Inches(0.98)
    c = PURP if i >= 8 else MUTED
    Ov(s13, x, Inches(2.04), Inches(0.5), Inches(0.5), dk(c,0.3) if i < 8 else dk(c,0.5), c)
    T(s13, str(i+1), x, Inches(2.04), Inches(0.5), Inches(0.5),
      sz=9, col=c if i >= 8 else MUTED, bold=(i>=8), align=PP_ALIGN.CENTER)

cards = [
    (CYAN, "09  TIMING",
     "Number of frames = perceived\nspeed, weight, and personality.\n\n"
     "MORE frames = slower, heavier\nFEWER frames = faster, lighter\n\n"
     "AT 24fps:\n"
     "  Eye blink: 6–8 frames\n"
     "  Head turn: 10–15 frames\n"
     "  Heavy punch: 3–4 frames\n"
     "  Gentle wave: 30+ frames\n\n"
     "TIMING is the HEARTBEAT\nof animation. Wrong timing =\nnothing else matters."),
    (ORG, "10  EXAGGERATION",
     "Push actions BEYOND realistic\nlimits for emotional impact.\n\n"
     "• Giant cartoon eyes\n"
     "• Impossible stretching & squash\n"
     "• Over-the-top expressions\n"
     "• Massive emotional reactions\n\n"
     "KEY INSIGHT:\n"
     "NOT about being unrealistic —\nit's making the TRUTH more\nVISIBLE and ENTERTAINING.\n\n"
     "Exaggerate the ESSENCE.\nA sad character cries buckets,\nnot small tears."),
    (GRN, "11  SOLID DRAWING",
     "Characters must feel like 3D\nobjects with real weight & form.\n\n"
     "PRINCIPLES:\n"
     "  • Understand 3D basic shapes\n"
     "  • 'Draw through' the form\n"
     "  • Consistent character proportions\n"
     "  • Avoid 'twinning' (mirror poses)\n"
     "  • Strong line of action\n\n"
     "Even in 2D — the audience must\nFEEL the three-dimensional volume.\n\n"
     "Modern 3D: rigs must deform\nbelievably under all conditions"),
    (PURP, "12  APPEAL",
     "Characters must be visually\ninteresting to WATCH.\n\n"
     "APPEAL is NOT cuteness:\n"
     "  • A powerful villain has appeal\n"
     "  • A complex anti-hero has appeal\n"
     "  • Even monsters can have appeal\n\n"
     "HOW TO ACHIEVE:\n"
     "  • Clear, readable design\n"
     "  • Dynamic, asymmetric poses\n"
     "  • Charisma & personality\n"
     "  • Strong recognizable silhouette\n\n"
     "THE TEST: Would an audience\nwant to keep watching this character?"),
]
card4(s13, cards, sy=Inches(2.58))

# ════════════════════════════════════════════════════
# SLIDE 14 — ANIMATION SPACE
# ════════════════════════════════════════════════════
s14 = ns()
hdr(s14, "08", "ANIMATION SPACE", PINK)
ibar(s14,
     "Animation can exist in THREE types of spatial dimensions — "
     "each with different tools, techniques, rendering requirements, and visual results.", PINK)

# Visual: space dimension diagram
# 2D axes visual
R(s14, Inches(0.38), Inches(2.08), Inches(3.95), Inches(0.38), dk(CYAN,0.3), CYAN)
T(s14, "X", Inches(0.45), Inches(2.08), Inches(0.5), Inches(0.38),
  sz=12, col=CYAN, bold=True, align=PP_ALIGN.CENTER)
T(s14, "—————————→", Inches(0.82), Inches(2.08), Inches(3.45), Inches(0.38),
  sz=11, col=CYAN, align=PP_ALIGN.CENTER)

R(s14, Inches(0.38), Inches(2.1), Inches(0.38), Inches(3.55), dk(CYAN,0.3), CYAN)
T(s14, "Y\n|\n|\n↓", Inches(0.38), Inches(2.1), Inches(0.38), Inches(3.55),
  sz=10, col=CYAN, bold=True, align=PP_ALIGN.CENTER)

cards = [
    (CYAN, "2D SPACE",
     "X-axis (width) + Y-axis (height)\nNo depth — completely flat.\n\n"
     "CHARACTERISTICS:\n"
     "• Flat, like drawing on paper\n"
     "• No real depth perception\n"
     "• Simple to create & render\n"
     "• Low hardware requirements\n"
     "• Classic cartoon aesthetic\n\n"
     "EXAMPLES:\n"
     "Tom & Jerry, classic Disney,\noriginal Mario & Pac-Man,\nmobile game animations,\nFlash web animations\n\n"
     "TOOLS:\n"
     "Adobe Animate, Toon Boom,\nSpine, TVPaint, Aseprite"),
    (ORG, "2½D SPACE\n(Two-and-a-Half D)",
     "2D art + ILLUSION of depth via\nlayers, parallax & shadows.\n\n"
     "CHARACTERISTICS:\n"
     "• Parallax scrolling = depth effect\n"
     "• Objects appear on Z-axis visually\n"
     "• Still technically 2D sprites\n"
     "• 'Pseudo-3D' / 'Layered 2D'\n"
     "• Great performance/visual ratio\n\n"
     "EXAMPLES:\n"
     "Paper Mario, South Park,\nOri and the Blind Forest,\nDisney multiplane camera\n(Sleeping Beauty, 1959)\n\n"
     "TOOLS:\n"
     "After Effects (Z-depth layers),\nUnity 2D with parallax"),
    (GRN, "3D SPACE",
     "X-axis + Y-axis + Z-axis (depth)\nFull volumetric environment.\n\n"
     "CHARACTERISTICS:\n"
     "• True depth, volume, perspective\n"
     "• Viewable from ANY camera angle\n"
     "• Realistic lighting & shadows\n"
     "• Requires powerful hardware\n"
     "• Maximum visual versatility\n\n"
     "EXAMPLES:\n"
     "Toy Story, Frozen, Moana,\nFortnite, Call of Duty,\narchitectural visualizations,\nmedical & surgical simulations\n\n"
     "TOOLS:\n"
     "Blender, Maya, Cinema 4D,\nHoudini, Unreal Engine"),
]
card3(s14, cards, sy=Inches(2.12))

# ════════════════════════════════════════════════════
# SLIDE 15 — ANIMATION PROCESS
# ════════════════════════════════════════════════════
s15 = ns()
hdr(s15, "09", "ANIMATION PROCESS", ORG)
ibar(s15,
     "Every animation project — from a 10-second ad to a 2-hour feature film — follows this same 4-stage pipeline.", ORG)

# Visual: pipeline connector bar
GR(s15, Inches(0.38), Inches(4.35), Inches(12.57), Inches(0.12), ORG, PURP, 0)

steps = [
    (CYAN,  "1","ORGANIZE\nEXECUTION",
     "• Define project goals & scope\n"
     "• Assign team roles\n"
     "• Set milestones & deadlines\n"
     "• Gather reference material\n"
     "• Create style guide\n"
     "• Budget & risk planning\n\n"
     "OUTPUT:\n"
     "Project bible, schedule,\nmood boards, storyboard"),
    (ORG,   "2","CHOOSE\nANIMATION TOOL",
     "• Match tool to project type\n"
     "• Consider team skill level\n"
     "• Account for budget\n\n"
     "2D: Adobe Animate,\n     Toon Boom Harmony\n"
     "3D: Maya, Blender,\n     Cinema 4D\n"
     "MG: After Effects\n"
     "Games: Unity, Unreal"),
    (GRN,   "3","BUILD & TWEAK\nSEQUENCES",
     "• Create rough blocking\n"
     "• Refine with keyframes\n"
     "• Apply the 12 Principles\n"
     "• Sync with audio/music\n"
     "• Team review cycles\n"
     "• Iterate until approved\n\n"
     "⟳ ITERATIVE STEP —\nrepeated many times\nbefore moving to Step 4"),
    (PURP,  "4","POST-PROCESS\nANIMATION",
     "• Render final frames\n"
     "• Color grading & correction\n"
     "• Add VFX & overlays\n"
     "• Sound design & music mix\n"
     "• QA testing & review\n"
     "• Format for delivery\n"
     "• Archive project files\n\n"
     "OUTPUT:\n"
     "Final MP4, MOV, WebM"),
]

for i,(ac,n,ti,bo) in enumerate(steps):
    l = Inches(0.38) + i*Inches(3.24)
    t = Inches(2.12)
    h = Inches(5.0)
    R(s15, l, t, Inches(3.08), h, CARD)
    GR(s15, l, t, Inches(3.08), Inches(0.65), ac, dk(ac,0.5), 0)
    GR(s15, l, t, Inches(0.2), h, ac, dk(ac))
    # Step circle
    circ = s15.shapes.add_shape(9, l+Inches(1.24), Inches(4.2), Inches(0.62), Inches(0.62))
    circ.fill.solid(); circ.fill.fore_color.rgb = ac; circ.line.color.rgb = ac
    T(s15, n, l+Inches(1.24), Inches(4.2), Inches(0.62), Inches(0.62),
      sz=19, col=DARK, bold=True, align=PP_ALIGN.CENTER)
    T(s15, ti, l+Inches(0.28), t+Inches(0.1), Inches(2.85), Inches(0.6),
      sz=12, col=DARK, bold=True, align=PP_ALIGN.CENTER)
    T(s15, bo, l+Inches(0.28), t+Inches(0.78), Inches(2.85), h-Inches(0.9),
      sz=10.5, col=LT)
    if i < 3:
        T(s15, "▶", l+Inches(3.12), Inches(4.05), Inches(0.22), Inches(0.58),
          sz=14, col=ORG, bold=True, align=PP_ALIGN.CENTER)

GR(s15, Inches(0.38), Inches(7.1), Inches(12.57), Inches(0.36), dk(GRN,0.25), dk(GRN,0.38), 0)
T(s15, "KEY INSIGHT: Step 3 is ITERATIVE — you cycle through Build & Tweak many times before Post-Process.",
  Inches(0.55), Inches(7.12), Inches(12.2), Inches(0.32), sz=11.5, col=GRN, bold=True)

# ════════════════════════════════════════════════════
# SLIDE 16 — SUMMARY
# ════════════════════════════════════════════════════
s16 = ns()
film_strip(s16, Inches(0))
GR(s16, Inches(0), Inches(0.3), Inches(13.33), Inches(1.05), CYAN, PURP, 0)
T(s16, "SUMMARY — KEY TAKEAWAYS",
  Inches(0.45), Inches(0.36), Inches(12.5), Inches(0.9),
  sz=33, col=WHT, bold=True, align=PP_ALIGN.CENTER)

left = [
    ("01",CYAN,  "Introduction",   "Animation = sequences exploiting Persistence of Vision + Phi Phenomenon to create seamless motion illusion."),
    ("02",PURP,  "Usage",          "Art, storytelling, scientific visualization, instructional. Animation is used across every major industry."),
    ("03",ORG,   "Techniques",     "3 core techniques: Cel (hand-drawn), Path (trajectory), Computer (digital — most versatile today)."),
    ("04",GRN,   "Cel Animation",  "Keyframes + tweening = smooth motion. Pencil test validates timing before expensive final production."),
    ("05",PINK,  "Path Animation", "Object follows predefined path (straight or Bezier) with easing for natural, weighted motion feel."),
]
right = [
    ("06",CYAN,  "Computer Anim",  "FK/IK chains for character rigs. Morphing for shape transforms. Faster & more flexible than traditional."),
    ("07",PURP,  "Levels",         "Basic (GIFs/CSS) → Intermediate (motion graphics) → Advanced (feature films, real-time VFX, AI animation)."),
    ("08",ORG,   "12 Principles",  "Disney's Squash & Stretch, Anticipation, Timing, Ease, Arcs, and 8 more — still the universal standard."),
    ("09",GRN,   "Anim Space",     "2D (flat) → 2½D (pseudo-depth via parallax) → 3D (full volumetric, any camera angle, photorealism)."),
    ("  ",YEL,   "Process",        "Organize → Choose Tool → Build & Tweak (iterative) → Post-Process. Every project, every studio."),
]
for i,(n,c,ti,de) in enumerate(left):
    y = Inches(1.42)+i*Inches(1.18)
    GR(s16, Inches(0.35), y, Inches(0.72), Inches(1.0), c, dk(c,0.5), 90)
    T(s16, n, Inches(0.35), y, Inches(0.72), Inches(1.0), sz=17, col=DARK, bold=True, align=PP_ALIGN.CENTER)
    R(s16, Inches(1.12), y, Inches(5.5), Inches(1.0), CARD)
    T(s16, ti, Inches(1.2), y+Inches(0.04), Inches(5.35), Inches(0.38), sz=12, col=c, bold=True)
    T(s16, de, Inches(1.2), y+Inches(0.42), Inches(5.35), Inches(0.55), sz=10.5, col=LT)
for i,(n,c,ti,de) in enumerate(right):
    y = Inches(1.42)+i*Inches(1.18)
    GR(s16, Inches(7.05), y, Inches(0.72), Inches(1.0), c, dk(c,0.5), 90)
    T(s16, n, Inches(7.05), y, Inches(0.72), Inches(1.0), sz=17, col=DARK, bold=True, align=PP_ALIGN.CENTER)
    R(s16, Inches(7.82), y, Inches(5.13), Inches(1.0), CARD)
    T(s16, ti, Inches(7.9), y+Inches(0.04), Inches(5.0), Inches(0.38), sz=12, col=c, bold=True)
    T(s16, de, Inches(7.9), y+Inches(0.42), Inches(5.0), Inches(0.55), sz=10.5, col=LT)
R(s16, Inches(6.6), Inches(1.32), Inches(0.06), Inches(6.0), RGBColor(30,30,80))

# ════════════════════════════════════════════════════
# SLIDE 17 — THANK YOU
# ════════════════════════════════════════════════════
s17 = ns()
GR3(s17, Inches(0), Inches(0), Inches(13.33), Inches(7.5),
    RGBColor(3,3,16), RGBColor(8,5,32), RGBColor(4,4,18), 135)

film_strip(s17, Inches(0))
film_strip(s17, Inches(7.2))

for x,y,sz,c in [
    (Inches(10.2), Inches(-1.0), Inches(4.5), dk(CYAN,0.14)),
    (Inches(11.8), Inches(5.0),  Inches(3.0), dk(PURP,0.14)),
    (Inches(-0.8), Inches(4.5),  Inches(3.2), dk(ORG,0.11)),
]:
    sh = s17.shapes.add_shape(9, x, y, sz, sz)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.color.rgb = c

GR(s17, Inches(0), Inches(0.3), Inches(13.33), Inches(0.18), CYAN, PURP, 0)
GR(s17, Inches(0), Inches(7.02), Inches(13.33), Inches(0.18), PURP, ORG, 0)

T(s17, "ANIM", Inches(-0.4), Inches(0.2), Inches(9), Inches(5.2),
  sz=195, col=RGBColor(9,9,38), bold=True)

T(s17, "THANK YOU", Inches(0.5), Inches(0.88), Inches(12.33), Inches(2.4),
  sz=88, col=CYAN, bold=True, align=PP_ALIGN.CENTER)

GR(s17, Inches(2.2), Inches(3.05), Inches(8.93), Inches(0.78), dk(CYAN,0.28), dk(PURP,0.28), 0)
T(s17, "Animation: Where Imagination Meets Technology",
  Inches(2.2), Inches(3.07), Inches(8.93), Inches(0.75),
  sz=18, col=LT, italic=True, align=PP_ALIGN.CENTER)

GR(s17, Inches(3.5), Inches(3.98), Inches(6.33), Inches(0.07), CYAN, PURP, 0)

T(s17, "Group 3  •  Graphics & Multimedia  •  SMIT",
  Inches(1), Inches(4.12), Inches(11.33), Inches(0.55),
  sz=15, col=ORG, bold=True, align=PP_ALIGN.CENTER)

members = [
    ("Hashir Junaid",     "2023F-BCS-358"),
    ("Taha Haider",       "2023F-BCS-079"),
    ("Abdul Rahman Baig", "2023F-BCS-100"),
]
for i,(name,roll) in enumerate(members):
    x = Inches(0.85)+i*Inches(3.9)
    GR(s17, x, Inches(4.82), Inches(3.65), Inches(1.05), CARD, CARD2, 0)
    Ov(s17, x+Inches(0.15), Inches(4.96), Inches(0.3), Inches(0.3), CYAN)
    T(s17, name, x+Inches(0.58), Inches(4.86), Inches(3.0), Inches(0.5),
      sz=14, col=WHT, bold=True)
    T(s17, roll, x+Inches(0.58), Inches(5.3), Inches(3.0), Inches(0.45),
      sz=11, col=MUTED)

T(s17, "Topics Covered: Introduction  •  Usage  •  Techniques  •  Types  •  Computer Animation  •  Levels  •  12 Principles  •  Space  •  Process",
  Inches(0.5), Inches(6.58), Inches(12.33), Inches(0.5),
  sz=10.5, col=RGBColor(75,88,135), italic=True, align=PP_ALIGN.CENTER)

# ─── SAVE ───────────────────────────────────────────
out = r"C:\Users\Noman Traders\Desktop\SMIT -AI-CLAUDE\Graphics n Multimedia\Animation_Group3_ULTRA.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
