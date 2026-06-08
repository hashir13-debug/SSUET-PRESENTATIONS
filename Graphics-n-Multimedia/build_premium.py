"""
Animation Presentation — PREMIUM VERSION
Group 3 | SMIT Graphics & Multimedia
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── PALETTE ─────────────────────────────────────────
BG    = RGBColor(6,   6,  28)
CARD  = RGBColor(16, 16,  58)
CARD2 = RGBColor(22, 22,  75)
CYAN  = RGBColor(0,  210, 255)
PURP  = RGBColor(150, 80, 220)
ORG   = RGBColor(255, 145,  0)
GRN   = RGBColor(46,  213, 115)
PINK  = RGBColor(255,  90, 160)
YEL   = RGBColor(255, 215,  0)
WHT   = RGBColor(255, 255, 255)
LT    = RGBColor(190, 205, 235)
DARK  = RGBColor(6,    6,  28)
MUTED = RGBColor(110, 125, 165)

def dk(c, f=0.35):
    return RGBColor(min(255,int(c[0]*f)), min(255,int(c[1]*f)), min(255,int(c[2]*f)))

def lk(c, f=1.6):
    return RGBColor(min(255,int(c[0]*f)), min(255,int(c[1]*f)), min(255,int(c[2]*f)))

def hx(c):
    return f'{c[0]:02X}{c[1]:02X}{c[2]:02X}'

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── CORE HELPERS ────────────────────────────────────

def ns():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide

def R(slide, l, t, w, h, fill=CARD, line=None):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = fill if line is None else line
    return sh

def GR(slide, l, t, w, h, c1, c2, ang=135, lc=None):
    """Gradient rectangle."""
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c1
    sh.line.color.rgb = c1 if lc is None else lc
    spPr = sh._element.find(qn('p:spPr'))
    gx = (f'<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
          f'<a:gsLst>'
          f'<a:gs pos="0"><a:srgbClr val="{hx(c1)}"/></a:gs>'
          f'<a:gs pos="100000"><a:srgbClr val="{hx(c2)}"/></a:gs>'
          f'</a:gsLst>'
          f'<a:lin ang="{ang*60000}" scaled="0"/>'
          f'</a:gradFill>')
    gf = etree.fromstring(gx)
    sol = spPr.find(qn('a:solidFill'))
    if sol is not None: spPr.remove(sol)
    spPr.append(gf)
    return sh

def GR3(slide, l, t, w, h, c1, c2, c3, ang=135):
    """3-stop gradient rectangle."""
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c1
    sh.line.color.rgb = c1
    spPr = sh._element.find(qn('p:spPr'))
    gx = (f'<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
          f'<a:gsLst>'
          f'<a:gs pos="0"><a:srgbClr val="{hx(c1)}"/></a:gs>'
          f'<a:gs pos="50000"><a:srgbClr val="{hx(c2)}"/></a:gs>'
          f'<a:gs pos="100000"><a:srgbClr val="{hx(c3)}"/></a:gs>'
          f'</a:gsLst>'
          f'<a:lin ang="{ang*60000}" scaled="0"/>'
          f'</a:gradFill>')
    gf = etree.fromstring(gx)
    sol = spPr.find(qn('a:solidFill'))
    if sol is not None: spPr.remove(sol)
    spPr.append(gf)
    return sh

def T(slide, text, l, t, w, h, sz=14, col=WHT, bold=False,
      italic=False, align=PP_ALIGN.LEFT, wrap=True):
    bx = slide.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.color.rgb = col
    r.font.bold = bold; r.font.italic = italic
    return bx

def hdr(slide, num, title, accent):
    """Gradient header + big watermark number."""
    GR(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.15), accent, dk(accent,0.5), 0)
    # Watermark number
    T(slide, num, Inches(11.3), Inches(-0.2), Inches(2.0), Inches(1.5),
      sz=90, col=dk(accent,0.45), bold=True, align=PP_ALIGN.RIGHT)
    T(slide, title, Inches(0.45), Inches(0.1), Inches(11.0), Inches(0.98),
      sz=30, col=WHT, bold=True)

def ibar(slide, text, accent, y=None):
    y = y or Inches(1.18)
    GR(slide, Inches(0.38), y, Inches(12.57), Inches(0.82), dk(accent,0.28), dk(accent,0.45), 0)
    GR(slide, Inches(0.38), y, Inches(0.2), Inches(0.82), accent, dk(accent))
    T(slide, text, Inches(0.7), y+Inches(0.1), Inches(12.0), Inches(0.68),
      sz=12, col=LT, italic=True)

def card4(slide, cards, sy=Inches(1.18)):
    """4 cards in 2×2 grid."""
    cw, ch = Inches(6.25), Inches(2.9)
    pos = [(Inches(0.38), sy), (Inches(6.7), sy),
           (Inches(0.38), sy+Inches(3.05)), (Inches(6.7), sy+Inches(3.05))]
    for i,(ac,ti,bo) in enumerate(cards[:4]):
        l,top = pos[i]
        R(slide, l, top, cw, ch, CARD)
        GR(slide, l, top, cw, Inches(0.6), ac, dk(ac,0.5), 0)
        # Left glow bar
        GR(slide, l, top, Inches(0.18), ch, ac, dk(ac))
        T(slide, ti, l+Inches(0.28), top+Inches(0.06), cw-Inches(0.38), Inches(0.54),
          sz=14, col=ac, bold=True)
        T(slide, bo, l+Inches(0.28), top+Inches(0.65), cw-Inches(0.38), ch-Inches(0.76),
          sz=11, col=LT)

def card3(slide, cards, sy=Inches(1.18)):
    """3 cards in a row."""
    cw, ch = Inches(4.3), Inches(6.12)
    for i,(ac,ti,bo) in enumerate(cards[:3]):
        l = Inches(0.35) + i*Inches(4.45)
        R(slide, l, sy, cw, ch, CARD)
        GR(slide, l, sy, cw, Inches(0.7), ac, dk(ac,0.5), 0)
        GR(slide, l, sy, Inches(0.18), ch, ac, dk(ac))
        T(slide, ti, l+Inches(0.28), sy+Inches(0.08), cw-Inches(0.38), Inches(0.62),
          sz=14, col=DARK, bold=True, align=PP_ALIGN.CENTER)
        T(slide, bo, l+Inches(0.26), sy+Inches(0.8), cw-Inches(0.36), ch-Inches(0.92),
          sz=11, col=LT)

def card4h(slide, cards, sy=Inches(2.15), ch=Inches(5.1)):
    """4 cards horizontal row (narrow)."""
    cw = Inches(3.08)
    for i,(ac,ti,bo) in enumerate(cards[:4]):
        l = Inches(0.38) + i*Inches(3.24)
        R(slide, l, sy, cw, ch, CARD)
        GR(slide, l, sy, cw, Inches(0.62), ac, dk(ac,0.5), 0)
        GR(slide, l, sy, Inches(0.18), ch, ac, dk(ac))
        T(slide, ti, l+Inches(0.28), sy+Inches(0.08), cw-Inches(0.36), Inches(0.58),
          sz=12, col=DARK, bold=True, align=PP_ALIGN.CENTER)
        T(slide, bo, l+Inches(0.26), sy+Inches(0.72), cw-Inches(0.36), ch-Inches(0.84),
          sz=10.5, col=LT)

# ════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════
s1 = ns()
# Full gradient background
GR3(s1, Inches(0), Inches(0), Inches(13.33), Inches(7.5),
    RGBColor(4,4,22), RGBColor(10,6,40), RGBColor(6,6,28), 135)

# Decorative circles
for x,y,sz,c in [
    (Inches(11.5), Inches(-0.5), Inches(3.5), dk(CYAN,0.2)),
    (Inches(12.0), Inches(5.5),  Inches(2.0), dk(PURP,0.2)),
    (Inches(-0.3), Inches(5.0),  Inches(2.5), dk(ORG,0.18)),
    (Inches(0.5),  Inches(-0.3), Inches(1.5), dk(CYAN,0.15)),
]:
    sh = s1.shapes.add_shape(9, x, y, sz, sz)
    sh.fill.solid(); sh.fill.fore_color.rgb = c
    sh.line.color.rgb = c

# Top accent line
GR(s1, Inches(0), Inches(0), Inches(13.33), Inches(0.22), CYAN, PURP, 0)
# Bottom accent line
GR(s1, Inches(0), Inches(7.28), Inches(13.33), Inches(0.22), PURP, ORG, 0)

# Big watermark "ANIM"
T(s1, "ANIM", Inches(-0.2), Inches(0.3), Inches(8), Inches(4.5),
  sz=185, col=RGBColor(12,12,50), bold=True)

# Main title
T(s1, "ANIMATION", Inches(0.5), Inches(1.1), Inches(12.33), Inches(2.2),
  sz=85, col=CYAN, bold=True, align=PP_ALIGN.CENTER)

# Subtitle bar
GR(s1, Inches(2.0), Inches(3.2), Inches(9.33), Inches(0.75), dk(CYAN,0.3), dk(PURP,0.3), 0)
T(s1, "The Art & Science of Bringing Images to Life",
  Inches(2.0), Inches(3.22), Inches(9.33), Inches(0.72),
  sz=19, col=LT, italic=True, align=PP_ALIGN.CENTER)

# Divider line
GR(s1, Inches(3.5), Inches(4.1), Inches(6.33), Inches(0.07), CYAN, PURP, 0)

# Course tag
T(s1, "Graphics & Multimedia  •  SMIT  •  Group 3",
  Inches(1), Inches(4.25), Inches(11.33), Inches(0.55),
  sz=15, col=ORG, bold=True, align=PP_ALIGN.CENTER)

# Topics badge
GR(s1, Inches(4.5), Inches(4.95), Inches(4.33), Inches(0.7), dk(CYAN,0.4), dk(PURP,0.4), 0)
T(s1, "Topics 1 — 9", Inches(4.5), Inches(4.95), Inches(4.33), Inches(0.7),
  sz=17, col=WHT, bold=True, align=PP_ALIGN.CENTER)

# Member cards
members = [
    ("Hashir Junaid", "2023F-BCS-358"),
    ("Taha Haider",   "2023F-BCS-079"),
    ("Abdul Rahman Baig", "2023F-BCS-100"),
]
for i,(name,roll) in enumerate(members):
    x = Inches(1.0) + i*Inches(3.8)
    GR(s1, x, Inches(5.85), Inches(3.55), Inches(0.95), CARD, CARD2, 0)
    T(s1, name, x, Inches(5.88), Inches(3.55), Inches(0.48), sz=13, col=WHT, bold=True, align=PP_ALIGN.CENTER)
    T(s1, roll, x, Inches(6.3), Inches(3.55), Inches(0.45), sz=11, col=MUTED, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════
s2 = ns()
GR(s2, Inches(0), Inches(0), Inches(13.33), Inches(1.18), CYAN, PURP, 0)
T(s2, "WHAT WE'LL COVER TODAY", Inches(0.45), Inches(0.1), Inches(12), Inches(1.0),
  sz=34, col=WHT, bold=True)
T(s2, "9 Topics  •  Animation: Complete Foundation",
  Inches(0.45), Inches(0.72), Inches(12), Inches(0.42), sz=13, col=dk(CYAN,1.8), italic=True)

agenda = [
    ("01","Introduction to Animation",   CYAN),
    ("02","Usage of Animation",          PURP),
    ("03","Animation Techniques",        ORG),
    ("04","Types of Animation",          GRN),
    ("05","Computer Animation",          PINK),
    ("06","Levels of Computer Animation",YEL),
    ("07","12 Basic Principles",         CYAN),
    ("08","Animation Space",             PURP),
    ("09","Animation Process",           ORG),
]
# 3 columns × 3 rows
col_x = [Inches(0.35), Inches(4.68), Inches(9.0)]
row_y = [Inches(1.32), Inches(2.66), Inches(4.0)]
for i,(num,topic,ac) in enumerate(agenda):
    row,col = divmod(i,3)
    x,y = col_x[col], row_y[row]
    GR(s2, x, y, Inches(4.0), Inches(1.18), CARD, CARD2, 0)
    GR(s2, x, y, Inches(0.75), Inches(1.18), ac, dk(ac), 90)
    T(s2, num, x, y, Inches(0.75), Inches(1.18), sz=22, col=DARK, bold=True, align=PP_ALIGN.CENTER)
    T(s2, topic, x+Inches(0.85), y+Inches(0.25), Inches(3.05), Inches(0.72), sz=14, col=WHT)

# Bottom bar
GR(s2, Inches(0), Inches(5.28), Inches(13.33), Inches(2.22), BG, RGBColor(10,10,45), 90)
T(s2, "From hand-drawn cels to AI-generated 3D — we cover it all.",
  Inches(1), Inches(5.6), Inches(11.33), Inches(0.8),
  sz=20, col=MUTED, italic=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# SLIDE 3 — INTRODUCTION TO ANIMATION
# ════════════════════════════════════════════════════
s3 = ns()
hdr(s3, "01", "INTRODUCTION TO ANIMATION", CYAN)

# Big definition box
GR(s3, Inches(0.38), Inches(1.22), Inches(12.57), Inches(1.12), dk(CYAN,0.28), dk(CYAN,0.15), 0)
GR(s3, Inches(0.38), Inches(1.22), Inches(0.2), Inches(1.12), CYAN, dk(CYAN))
T(s3,
  '"Animation: The technique of rapidly displaying a SEQUENCE of static images to create the\n'
  ' ILLUSION OF MOVEMENT — exploiting the brain\'s own perception quirks against itself."',
  Inches(0.68), Inches(1.28), Inches(12.15), Inches(1.05),
  sz=14, col=WHT, italic=True)

cards = [
    (CYAN,  "DEFINITION",
     "From Latin 'anima' — to give soul.\n\n"
     "• Series of still frames shown rapidly\n"
     "• 24 fps = cinema standard\n"
     "• 30 fps = TV broadcast\n"
     "• 60 fps = modern gaming\n"
     "• < 12 fps → visible flicker\n\n"
     "First animation: Eadweard Muybridge's\nhorse gallop sequence (1878)"),
    (PURP,  "PERSISTENCE OF VISION",
     "Physical effect in the human eye.\n\n"
     "• Retina keeps an image ~50ms after\n  stimulus disappears\n"
     "• This 'afterimage' overlaps next frame\n"
     "• Creates a seamless visual flow\n"
     "• Discovered: Peter Mark Roget (1824)\n\n"
     "Result: our brains cannot perceive\nindividual frames at 24+ fps"),
    (ORG,   "PHI PHENOMENON",
     "Purely PSYCHOLOGICAL motion effect.\n\n"
     "• Brain infers motion between two\n  stationary objects shown in sequence\n"
     "• Does NOT need retinal retention\n"
     "• Discovered: Max Wertheimer (1912)\n\n"
     "Classic examples:\n"
     "• Flip-book animation\n"
     "• Neon sign arrows 'flowing'\n"
     "• Traffic light direction arrows"),
    (GRN,   "ILLUSION OF MOVEMENT",
     "The combined result of both effects.\n\n"
     "• POV alone ≠ enough\n"
     "• Phi alone ≠ enough\n"
     "• Together → seamless motion\n\n"
     "The brain WANTS to see motion —\nanimators exploit this biological\n'glitch' in human perception.\n\n"
     "This is why animation feels REAL\neven though it's completely artificial"),
]
card4(s3, cards, sy=Inches(2.42))

# ════════════════════════════════════════════════════
# SLIDE 4 — USAGE OF ANIMATION
# ════════════════════════════════════════════════════
s4 = ns()
hdr(s4, "02", "USAGE OF ANIMATION", PURP)

cards = [
    (CYAN, "ARTISTIC PURPOSES",
     "• Fine art & experimental film\n"
     "• Music video production\n"
     "• Abstract visual art\n"
     "• Title sequence design\n"
     "• Gallery installations\n\n"
     "ICONIC EXAMPLES:\n"
     "  Disney's Fantasia (1940)\n"
     "  Yellow Submarine (1968)\n"
     "  Spider-Man: Into the Spider-Verse"),
    (ORG, "STORYTELLING",
     "• Narrative feature films\n"
     "• TV series & web series\n"
     "• Video game cutscenes\n"
     "• Interactive web comics\n"
     "• Character-driven drama\n\n"
     "ICONIC EXAMPLES:\n"
     "  Pixar Toy Story series\n"
     "  Studio Ghibli classics\n"
     "  Arcane — Netflix (2021)"),
    (GRN, "SCIENTIFIC VISUALIZATION",
     "• Human anatomy explainers\n"
     "• Weather & climate models\n"
     "• DNA & molecular modeling\n"
     "• Physics simulations\n"
     "• Engineering CAD previews\n\n"
     "ICONIC EXAMPLES:\n"
     "  NASA mission visualizations\n"
     "  Medical surgical training\n"
     "  Climate change simulations"),
    (PINK, "INSTRUCTIONAL",
     "• Educational explainer videos\n"
     "• Military & flight simulators\n"
     "• Safety procedure demos\n"
     "• Product assembly guides\n"
     "• Corporate training modules\n\n"
     "ICONIC EXAMPLES:\n"
     "  Khan Academy explainers\n"
     "  Boeing pilot simulators\n"
     "  YouTube tutorial channels"),
]
card4(s4, cards, sy=Inches(1.22))

# ════════════════════════════════════════════════════
# SLIDE 5 — ANIMATION TECHNIQUES OVERVIEW
# ════════════════════════════════════════════════════
s5 = ns()
hdr(s5, "03", "ANIMATION TECHNIQUES", ORG)

cards = [
    (CYAN, "CEL ANIMATION",
     "Traditional hand-drawn technique.\n\n"
     "• Drawings on transparent celluloid sheets\n"
     "• Each sheet = one frame of animation\n"
     "• Cels layered over static backgrounds\n"
     "• Allows background reuse\n"
     "• Photographed frame-by-frame on film\n\n"
     "WORKFLOW:\n"
     "  Story → Storyboard → Layout →\n"
     "  Key animation → In-betweening →\n"
     "  Inking → Painting → Photography\n\n"
     "ERA: 1910s – 1990s\n"
     "STUDIOS: Disney, Hanna-Barbera,\n"
     "          Warner Bros, Toei Animation"),
    (ORG, "PATH ANIMATION",
     "Object follows a predefined trajectory.\n\n"
     "• Animator defines the path shape\n"
     "• Object moves along it automatically\n"
     "• Supports straight lines & Bezier curves\n"
     "• Object can resize/deform along path\n"
     "• Easing controls speed variation\n\n"
     "COMMON USES:\n"
     "  Logo reveal animations\n"
     "  Camera fly-through sequences\n"
     "  UI element transitions\n"
     "  Game projectile trajectories\n\n"
     "TOOLS: After Effects, CSS, SVG paths"),
    (GRN, "COMPUTER ANIMATION",
     "Digital creation of moving images.\n\n"
     "• Can simulate cel & path techniques\n"
     "• Plus: physics, particles, simulation\n"
     "• 2D digital, 3D CG, or hybrid mix\n"
     "• Real-time OR pre-rendered output\n"
     "• Procedural & AI-assisted animation\n\n"
     "ADVANTAGES:\n"
     "  Unlimited camera angles\n"
     "  Easy to iterate & fix mistakes\n"
     "  Physics simulation built-in\n"
     "  Automatic in-betweening\n\n"
     "ERA: 1980s – present\n"
     "STUDIOS: Pixar, DreamWorks, ILM"),
]
card3(s5, cards, sy=Inches(1.22))

# Arrow connectors
for i in range(2):
    ax = Inches(4.52)+i*Inches(4.45)
    GR(s5, ax, Inches(4.1), Inches(0.25), Inches(0.7), ORG, dk(ORG), 0)
    T(s5, "▶", ax-Inches(0.05), Inches(4.0), Inches(0.4), Inches(0.9),
      sz=20, col=ORG, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# SLIDE 6 — CEL ANIMATION DEEP DIVE
# ════════════════════════════════════════════════════
s6 = ns()
hdr(s6, "04", "TYPES OF ANIMATION — CEL ANIMATION", CYAN)
ibar(s6,
     "Named after CELLULOID — transparent plastic sheets invented in 1869. "
     "Raoul Barré & John Bray developed the cel technique ~1914. It dominated animation for 80 years.", CYAN)

cards = [
    (CYAN, "CEL ANIMATION",
     "• Transparent sheet per frame of film\n"
     "• Character drawn on cel; BG is separate\n"
     "• Cels stacked → final composite frame\n"
     "• Allows background reuse (saves cost)\n"
     "• Colors applied to BACK of cel sheet\n"
     "• ~12–24 unique cels per second\n\n"
     "USED IN:\n"
     "Tom & Jerry, Looney Tunes,\n"
     "Snow White, Bambi, Akira"),
    (ORG, "KEYFRAMES",
     "• The 'anchor' frames of an action\n"
     "• Show the most extreme or critical poses\n"
     "• Drawn by SENIOR animators\n"
     "• Define the timing structure\n"
     "• Think of them as checkpoints\n\n"
     "EXAMPLE:\n"
     "Frame 1 = character standing\n"
     "Frame 12 = peak of jump\n"
     "Frame 24 = landed\n\n"
     "Modern software uses identical\nkeyframe concept in timelines"),
    (GRN, "TWEENING (In-Betweening)",
     "• Frames drawn BETWEEN keyframes\n"
     "• Creates smooth motion transition\n"
     "• Done by junior animators\n"
     "• More tweens = smoother motion\n"
     "• At 24fps: ~22 tweens per action\n\n"
     "MODERN EQUIVALENT:\n"
     "Software auto-tweens between\nkeyframes — no manual drawing\n\n"
     "Tools: Adobe Animate, Flash,\n"
     "Toon Boom Harmony, After Effects"),
    (PURP, "PENCIL TEST",
     "• Quick pre-production test\n"
     "• Rough pencil sketches filmed first\n"
     "• Played back to check timing & flow\n"
     "• Catches errors BEFORE expensive\n  inking & painting stages\n"
     "• Saves significant time and money\n\n"
     "MODERN EQUIVALENT:\n"
     "Animatic / rough cut preview\nplayed back in editing software\n\n"
     "Standard at every major studio\nbefore committing to final art"),
]
card4(s6, cards, sy=Inches(2.1))

# ════════════════════════════════════════════════════
# SLIDE 7 — PATH ANIMATION DEEP DIVE
# ════════════════════════════════════════════════════
s7 = ns()
hdr(s7, "04", "TYPES OF ANIMATION — PATH ANIMATION", ORG)
ibar(s7,
     "PATH ANIMATION: An object follows a predefined trajectory. "
     "The animator draws the path; the software handles the movement along it.", ORG)

# Left panel: Concepts
R(s7, Inches(0.38), Inches(2.08), Inches(6.08), Inches(5.2), CARD)
GR(s7, Inches(0.38), Inches(2.08), Inches(6.08), Inches(0.62), ORG, dk(ORG,0.5), 0)
GR(s7, Inches(0.38), Inches(2.08), Inches(0.2), Inches(5.2), ORG, dk(ORG))
T(s7, "KEY CONCEPTS", Inches(0.68), Inches(2.13), Inches(5.6), Inches(0.55),
  sz=15, col=DARK, bold=True)

pts = [
    (ORG,  "PREDETERMINED PATH",
     "Route defined BEFORE animation. Like a roller-coaster track\n— the object follows exactly what's laid out."),
    (CYAN, "STRAIGHT LINES",
     "Linear movement. Used for mechanical objects, text reveals, UI animations — natural for rigid elements."),
    (GRN,  "BEZIER CURVES",
     "Smooth organic paths via control points. Used for balls, flying objects, organic motion. Far more natural."),
    (PINK, "RESIZING & RESHAPING",
     "Object can scale or deform as it travels the path — adds dynamism (e.g. ball squashing on bounce)."),
]
yo = Inches(2.78)
for ac,ti,de in pts:
    GR(s7, Inches(0.6), yo, Inches(0.12), Inches(0.88), ac, dk(ac))
    T(s7, ti, Inches(0.82), yo, Inches(5.5), Inches(0.38), sz=12, col=ac, bold=True)
    T(s7, de, Inches(0.82), yo+Inches(0.38), Inches(5.5), Inches(0.52), sz=10.5, col=LT)
    yo += Inches(1.12)

# Right panel: How it works
R(s7, Inches(6.8), Inches(2.08), Inches(6.15), Inches(5.2), CARD)
GR(s7, Inches(6.8), Inches(2.08), Inches(6.15), Inches(0.62), ORG, dk(ORG,0.5), 0)
GR(s7, Inches(6.8), Inches(2.08), Inches(0.2), Inches(5.2), ORG, dk(ORG))
T(s7, "HOW PATH ANIMATION WORKS", Inches(7.1), Inches(2.13), Inches(5.75), Inches(0.55),
  sz=15, col=DARK, bold=True)

steps6 = [
    (CYAN,  "1","Draw the PATH — straight, curve, or freehand shape"),
    (ORG,   "2","Set START and END points along the path"),
    (GRN,   "3","Define DURATION — how long the object travels"),
    (PURP,  "4","Apply EASING — slow in, slow out for natural feel"),
    (PINK,  "5","Set OBJECT ORIENTATION — face direction of travel?"),
    (YEL,   "6","Preview, adjust, repeat — render when satisfied"),
]
yo = Inches(2.82)
for ac,n,desc in steps6:
    GR(s7, Inches(6.95), yo, Inches(0.55), Inches(0.62), ac, dk(ac))
    T(s7, n, Inches(6.95), yo, Inches(0.55), Inches(0.62), sz=16, col=DARK, bold=True, align=PP_ALIGN.CENTER)
    T(s7, desc, Inches(7.6), yo+Inches(0.1), Inches(5.25), Inches(0.55), sz=12, col=LT)
    yo += Inches(0.78)

GR(s7, Inches(0.38), Inches(7.08), Inches(12.57), Inches(0.38), dk(GRN,0.28), dk(GRN,0.4), 0)
T(s7, "Used in: CSS animations  •  After Effects motion paths  •  Unity NavMesh  •  Game projectile systems",
  Inches(0.55), Inches(7.1), Inches(12.2), Inches(0.34), sz=11, col=GRN)

# ════════════════════════════════════════════════════
# SLIDE 8 — COMPUTER ANIMATION (CEL vs COMPUTER)
# ════════════════════════════════════════════════════
s8 = ns()
hdr(s8, "05", "COMPUTER ANIMATION", GRN)
ibar(s8,
     "Computer animation uses software to create motion — simulating traditional techniques "
     "or producing entirely new effects impossible by hand.", GRN)

T(s8, "CEL  vs  COMPUTER ANIMATION",
  Inches(0.5), Inches(2.08), Inches(12.33), Inches(0.6),
  sz=20, col=YEL, bold=True, align=PP_ALIGN.CENTER)

# Header row
GR(s8, Inches(0.38), Inches(2.72), Inches(5.6), Inches(0.52), CYAN, dk(CYAN,0.5), 0)
T(s8, "CEL ANIMATION", Inches(0.38), Inches(2.72), Inches(5.6), Inches(0.52),
  sz=14, col=DARK, bold=True, align=PP_ALIGN.CENTER)
GR(s8, Inches(7.35), Inches(2.72), Inches(5.6), Inches(0.52), GRN, dk(GRN,0.5), 0)
T(s8, "COMPUTER ANIMATION", Inches(7.35), Inches(2.72), Inches(5.6), Inches(0.52),
  sz=14, col=DARK, bold=True, align=PP_ALIGN.CENTER)

# VS badge
GR(s8, Inches(6.05), Inches(2.6), Inches(1.22), Inches(1.05), ORG, dk(ORG,0.5), 90)
T(s8, "VS", Inches(6.05), Inches(2.6), Inches(1.22), Inches(1.05),
  sz=26, col=WHT, bold=True, align=PP_ALIGN.CENTER)

rows = [
    ("Hand-drawn on physical celluloid sheets", "Created entirely within digital software"),
    ("Every frame drawn manually — very slow", "Automation handles tweens, physics, rendering"),
    ("Physical storage: film cans, warehouses", "Digital files — backup & share instantly"),
    ("Expensive: cels, paint, camera, lab", "One-time software cost; no consumables"),
    ("Difficult to modify finished frames", "Copy, adjust, modify any frame instantly"),
    ("Organic hand-crafted visual warmth", "Any style from toon to photorealistic"),
    ("Production: months per 1 minute of film", "Real-time previews; faster iterations"),
]
for i,(cel,comp) in enumerate(rows):
    y = Inches(3.3) + i*Inches(0.56)
    bg = CARD if i%2==0 else CARD2
    R(s8, Inches(0.38), y, Inches(5.6), Inches(0.52), bg)
    R(s8, Inches(7.35), y, Inches(5.6), Inches(0.52), bg)
    T(s8, cel,  Inches(0.5),  y+Inches(0.07), Inches(5.35), Inches(0.44), sz=11, col=LT)
    T(s8, comp, Inches(7.47), y+Inches(0.07), Inches(5.35), Inches(0.44), sz=11, col=GRN)

# ════════════════════════════════════════════════════
# SLIDE 9 — FK / IK / MORPHING
# ════════════════════════════════════════════════════
s9 = ns()
hdr(s9, "05", "KINEMATICS, INVERSE KINEMATICS & MORPHING", GRN)

cards = [
    (CYAN, "FORWARD KINEMATICS (FK)",
     "Animating from PARENT → CHILD joint.\n\n"
     "HOW IT WORKS:\n"
     "  Move SHOULDER →\n"
     "  Elbow follows →\n"
     "  Wrist follows →\n"
     "  Finger follows\n\n"
     "ANALOGY: A puppet — control\nthe parent; all children follow.\n\n"
     "BEST FOR:\n"
     "  Mechanical/robotic arms\n"
     "  Waving & swinging motions\n"
     "  Tail & tentacle animation\n"
     "  Flag ripples in wind\n\n"
     "Used in: Maya, Blender rigs"),
    (ORG, "INVERSE KINEMATICS (IK)",
     "Reverse of FK — move the END\nPOINT; chain auto-adjusts.\n\n"
     "HOW IT WORKS:\n"
     "  Move HAND position →\n"
     "  Software calculates elbow →\n"
     "  And shoulder automatically\n\n"
     "ANALOGY: Reaching for a cup —\nyou aim the hand; your body\nfigures out the rest.\n\n"
     "BEST FOR:\n"
     "  Character walking & running\n"
     "  Reaching & grabbing actions\n"
     "  Real-time game characters\n"
     "  Interactive VR animation"),
    (PURP, "MORPHING",
     "Smooth shape/image transformation.\n\n"
     "HOW IT WORKS:\n"
     "  Map points on source image\n"
     "  Map matching target points\n"
     "  Interpolate all points together\n"
     "  Creates seamless in-betweens\n\n"
     "TYPES:\n"
     "  2D: image-to-image morph\n"
     "  3D: mesh deformation\n"
     "  Shape: vector path morphing\n\n"
     "FAMOUS EXAMPLES:\n"
     "  MJ 'Black or White' (faces)\n"
     "  T-1000 — Terminator 2\n"
     "  Werewolf transformations\n"
     "  Logo reveal animations"),
]
card3(s9, cards, sy=Inches(1.22))

# ════════════════════════════════════════════════════
# SLIDE 10 — LEVELS OF COMPUTER ANIMATION
# ════════════════════════════════════════════════════
s10 = ns()
hdr(s10, "06", "LEVELS OF COMPUTER ANIMATION", YEL)

cards = [
    (CYAN, "BASIC  ★☆☆",
     "Entry-level tools, simple output.\n\n"
     "WHAT YOU CREATE:\n"
     "  • Animated GIFs\n"
     "  • CSS / JS transitions\n"
     "  • PowerPoint animations\n"
     "  • Simple 2D sprites\n"
     "  • Basic logo animations\n\n"
     "TOOLS:\n"
     "  MS PowerPoint, Canva,\n"
     "  GIMP, basic CSS/JS\n\n"
     "SKILL: Beginner\n"
     "TIME: Days to weeks\n\n"
     "CAREER: Social media content,\n"
     "marketing materials"),
    (ORG, "INTERMEDIATE  ★★☆",
     "Deeper control over animation.\n\n"
     "WHAT YOU CREATE:\n"
     "  • 2D character animation\n"
     "  • Motion graphics\n"
     "  • Complex path animations\n"
     "  • Basic 3D object animation\n"
     "  • Game character movement\n\n"
     "TOOLS:\n"
     "  Adobe After Effects,\n"
     "  Adobe Animate,\n"
     "  Blender (basics), Spine\n\n"
     "SKILL: Intermediate\n"
     "TIME: Months\n\n"
     "CAREER: Motion designer,\n"
     "2D animator, game dev"),
    (PURP, "ADVANCED  ★★★",
     "Studio-level production quality.\n\n"
     "WHAT YOU CREATE:\n"
     "  • Full CG feature films\n"
     "  • Realistic physics sim\n"
     "  • Procedural animation\n"
     "  • Motion capture integration\n"
     "  • Real-time VFX in games\n"
     "  • AI-assisted animation\n\n"
     "TOOLS:\n"
     "  Autodesk Maya, Houdini,\n"
     "  Unreal Engine, Cinema 4D\n\n"
     "SKILL: Professional\n"
     "TIME: Years\n\n"
     "CAREER: VFX artist, 3D\n"
     "animator, technical director"),
]
card3(s10, cards, sy=Inches(1.22))

# ════════════════════════════════════════════════════
# SLIDE 11 — 12 PRINCIPLES PART 1 (1-4)
# ════════════════════════════════════════════════════
s11 = ns()
hdr(s11, "07", "12 BASIC PRINCIPLES OF ANIMATION  •  Part 1 of 3", CYAN)
ibar(s11,
     "Developed by Disney masters Ollie Johnston & Frank Thomas (The Illusion of Life, 1981). "
     "Still the universal gold standard — used by every professional animator on Earth.", CYAN)

cards = [
    (CYAN, "01  SQUASH & STRETCH",
     "Exaggerate deformation for weight.\n\n"
     "• Object SQUASHES on impact\n"
     "• Object STRETCHES when moving fast\n"
     "• Golden rule: VOLUME stays constant\n"
     "• Shows elasticity, mass, and life\n\n"
     "EXAMPLE:\n"
     "A rubber ball bounce — squash\non landing, stretch in the air.\n\n"
     "Most fundamental principle —\nused in almost every shot"),
    (ORG, "02  ANTICIPATION",
     "Small opposite action BEFORE main.\n\n"
     "• Character crouches before jumping\n"
     "• Arm swings back before throwing\n"
     "• Boxer winds up before punching\n"
     "• Prepares audience mentally\n\n"
     "EXAMPLE:\n"
     "Bugs Bunny ducking low before\na massive comedic leap forward.\n\n"
     "Without anticipation, actions\nfeel sudden and unconvincing"),
    (GRN, "03  STAGING",
     "Present ONE clear idea at a time.\n\n"
     "• Use clear, readable silhouettes\n"
     "• Camera angle serves the story\n"
     "• Avoid distracting background\n"
     "• Direct the audience's eye\n\n"
     "EXAMPLE:\n"
     "Character in profile view making\na dramatic gesture — instantly\nreadable from any distance.\n\n"
     "Think: theater staging, not\nphotography"),
    (PURP, "04  STRAIGHT AHEAD\n& POSE-TO-POSE",
     "Two approaches to animation:\n\n"
     "STRAIGHT AHEAD:\n"
     "• Draw frame by frame start→end\n"
     "• Spontaneous, fluid, surprising\n"
     "• Great for fire, water, chaos\n\n"
     "POSE-TO-POSE:\n"
     "• Draw key poses first\n"
     "• Fill in betweens after\n"
     "• Controlled, planned, consistent\n"
     "• Great for character acting\n\n"
     "Best work combines BOTH"),
]
card4(s11, cards, sy=Inches(2.08))

# ════════════════════════════════════════════════════
# SLIDE 12 — 12 PRINCIPLES PART 2 (5-8)
# ════════════════════════════════════════════════════
s12 = ns()
hdr(s12, "07", "12 BASIC PRINCIPLES OF ANIMATION  •  Part 2 of 3", ORG)
ibar(s12,
     "These principles govern how objects move AFTER the main action — "
     "the details that separate amateur animation from professional work.", ORG)

cards = [
    (CYAN, "05  FOLLOW THROUGH\n& OVERLAPPING ACTION",
     "FOLLOW THROUGH:\n"
     "Parts keep moving after main\naction stops.\n"
     "• Hair swings after head stops\n"
     "• Coat settles after character stops\n\n"
     "OVERLAPPING ACTION:\n"
     "Different parts move at different\nrates and with delay.\n"
     "• Arms lag behind torso\n"
     "• Secondary elements trail\n\n"
     "Makes motion feel ORGANIC\nand physically believable"),
    (ORG, "06  EASE IN & EASE OUT\n(Slow In / Slow Out)",
     "Acceleration at start, deceleration\nat end of every movement.\n\n"
     "• Nothing in nature starts at full\n  speed or stops instantly\n"
     "• A car: slow → full speed → brake\n"
     "• Gives physical weight & feel\n\n"
     "In software:\n"
     "• Bézier easing curves\n"
     "• Linear = robotic/artificial\n"
     "• Ease = natural/believable\n\n"
     "The single most impactful\nprinciple for 'weight' in motion"),
    (GRN, "07  ARCS",
     "Natural motion follows ARC paths,\nnot straight lines.\n\n"
     "• Arms swing in arcs\n"
     "• Heads nod in arcs\n"
     "• Thrown objects: parabola\n"
     "• Even eyes dart in arcs\n"
     "• Running legs follow arcs\n\n"
     "EXCEPTION:\n"
     "Mechanical/robot movement is\nintentionally straight — use linear\npaths for machines.\n\n"
     "Breaking arcs = instant 'robot'"),
    (PURP, "08  SECONDARY ACTION",
     "Supporting actions that ENRICH\nthe primary movement.\n\n"
     "EXAMPLES:\n"
     "  Walking (main) + arms swing\n"
     "  Running (main) + hair flows\n"
     "  Talking (main) + eyebrow moves\n\n"
     "CRITICAL RULE:\n"
     "Secondary must NEVER overpower\nthe main action — it supports,\nnot steals focus.\n\n"
     "Adds richness and believability\nto every scene"),
]
card4(s12, cards, sy=Inches(2.08))

# ════════════════════════════════════════════════════
# SLIDE 13 — 12 PRINCIPLES PART 3 (9-12)
# ════════════════════════════════════════════════════
s13 = ns()
hdr(s13, "07", "12 BASIC PRINCIPLES OF ANIMATION  •  Part 3 of 3", PURP)
ibar(s13,
     "Walt Disney: 'The secret of animation is the exaggeration of the truth.' "
     "These final 4 principles deliver that truth.", PURP)

cards = [
    (CYAN, "09  TIMING",
     "Number of frames = perceived\nspeed, weight, and personality.\n\n"
     "MORE frames = slower, heavier\nFEWER frames = faster, lighter\n\n"
     "AT 24fps:\n"
     "• Eye blink: 6–8 frames\n"
     "• Head turn: 10–15 frames\n"
     "• Heavy punch: 3–4 frames\n"
     "• Gentle wave: 30+ frames\n\n"
     "TIMING is the HEARTBEAT of\nanimation. Wrong timing = nothing\nelse matters."),
    (ORG, "10  EXAGGERATION",
     "Push actions BEYOND realistic\nlimits for emotional impact.\n\n"
     "• Giant cartoon eyes\n"
     "• Impossible stretching\n"
     "• Over-the-top expressions\n"
     "• Massive emotional reactions\n\n"
     "KEY INSIGHT:\n"
     "NOT about being unrealistic —\nit's making the TRUTH more\nVISIBLE and ENTERTAINING.\n\n"
     "Exaggerate the ESSENCE, not\nrandom features. A sad character\ncries buckets, not leaks."),
    (GRN, "11  SOLID DRAWING",
     "Characters must feel like 3D\nobjects with real weight & form.\n\n"
     "PRINCIPLES:\n"
     "• Understand 3D basic shapes\n"
     "• 'Draw through' the form\n"
     "• Consistent proportions\n"
     "• Avoid 'twinning' (mirror poses)\n"
     "• Strong line of action\n\n"
     "Even in 2D — the audience must\nFEEL the three-dimensional volume.\n\n"
     "Modern 3D: rigs must deform\nbelievably under all conditions"),
    (PURP, "12  APPEAL",
     "Characters must be visually\ninteresting to WATCH.\n\n"
     "APPEAL ≠ CUTENESS:\n"
     "• A powerful villain has appeal\n"
     "• A complex anti-hero has appeal\n"
     "• Even monsters can have appeal\n\n"
     "HOW TO ACHIEVE:\n"
     "• Clear, readable design\n"
     "• Dynamic, asymmetric poses\n"
     "• Charisma & personality\n"
     "• Strong silhouette\n\n"
     "THE TEST: Would an audience\nwant to keep watching this\ncharacter? If yes → appeal."),
]
card4(s13, cards, sy=Inches(2.08))

# ════════════════════════════════════════════════════
# SLIDE 14 — ANIMATION SPACE
# ════════════════════════════════════════════════════
s14 = ns()
hdr(s14, "08", "ANIMATION SPACE", PINK)
ibar(s14,
     "Animation can exist in three types of spatial dimensions — "
     "each with different tools, techniques, and visual results.", PINK)

cards = [
    (CYAN, "2D SPACE",
     "X-axis (width) + Y-axis (height)\nonly. No depth.\n\n"
     "CHARACTERISTICS:\n"
     "• Flat, like drawing on paper\n"
     "• No real depth perception\n"
     "• Simple to create & render\n"
     "• Low hardware requirements\n"
     "• Classic cartoon aesthetic\n\n"
     "EXAMPLES:\n"
     "Tom & Jerry, classic Disney,\noriginal Mario & Pac-Man,\nmobile game animations,\nFlash web animations\n\n"
     "TOOLS:\n"
     "Adobe Animate, Toon Boom,\nSpine, TVPaint, Aseprite"),
    (ORG, "2½D SPACE\n(Two-and-a-Half D)",
     "2D art + ILLUSION of depth via\nlayers, parallax & shadows.\n\n"
     "CHARACTERISTICS:\n"
     "• Parallax scrolling = depth\n"
     "• Objects on Z-axis visually\n"
     "• Still technically 2D sprites\n"
     "• 'Pseudo-3D' / 'Layered 2D'\n"
     "• Great performance/visual ratio\n\n"
     "EXAMPLES:\n"
     "Paper Mario, South Park,\nOri and the Blind Forest,\nDisney multiplane camera\n(Sleeping Beauty, 1959)\n\n"
     "TOOLS:\n"
     "After Effects (Z-depth layers),\nUnity 2D with parallax"),
    (GRN, "3D SPACE",
     "X-axis + Y-axis + Z-axis (depth)\nFull volumetric environment.\n\n"
     "CHARACTERISTICS:\n"
     "• True depth, volume, perspective\n"
     "• Viewable from any camera angle\n"
     "• Realistic lighting & shadows\n"
     "• Requires powerful hardware\n"
     "• Maximum visual versatility\n\n"
     "EXAMPLES:\n"
     "Toy Story, Frozen, Moana,\nFortnite, Call of Duty,\narchitectural visualizations,\nmedical & surgical simulations\n\n"
     "TOOLS:\n"
     "Blender, Maya, Cinema 4D,\nHoudini, Unreal Engine"),
]
card3(s14, cards, sy=Inches(2.08))

# ════════════════════════════════════════════════════
# SLIDE 15 — ANIMATION PROCESS
# ════════════════════════════════════════════════════
s15 = ns()
hdr(s15, "09", "ANIMATION PROCESS", ORG)
ibar(s15,
     "Every animation project — from a 10-second ad to a 2-hour feature film — follows this same 4-stage pipeline.", ORG)

# Pipeline connector bar
GR(s15, Inches(0.38), Inches(4.28), Inches(12.57), Inches(0.14), ORG, PURP, 0)

steps = [
    (CYAN,  "1","ORGANIZE\nEXECUTION",
     "• Define project goals & scope\n"
     "• Assign team roles\n"
     "• Set milestones & deadlines\n"
     "• Gather reference material\n"
     "• Create style guide\n"
     "• Budget & risk planning\n\n"
     "OUTPUT:\n"
     "Project bible, schedule,\nmood boards, storyboard"),
    (ORG,   "2","CHOOSE\nANIMATION TOOL",
     "• Match tool to project type\n"
     "• Consider team skill level\n"
     "• Account for budget\n\n"
     "2D: Adobe Animate,\n     Toon Boom\n"
     "3D: Maya, Blender,\n     Cinema 4D\n"
     "MG: After Effects\n"
     "Games: Unity, Unreal"),
    (GRN,   "3","BUILD & TWEAK\nSEQUENCES",
     "• Create rough blocking\n"
     "• Refine with keyframes\n"
     "• Apply the 12 Principles\n"
     "• Sync with audio/music\n"
     "• Team review cycles\n"
     "• Iterate until approved\n\n"
     "NOTE: This step repeats\nmultiple times before\nmoving to Step 4"),
    (PURP,  "4","POST-PROCESS\nANIMATION",
     "• Render final frames\n"
     "• Color grading & correction\n"
     "• Add VFX & overlays\n"
     "• Sound design & music mix\n"
     "• QA testing\n"
     "• Format for delivery\n"
     "• Archive project files\n\n"
     "OUTPUT:\n"
     "Final MP4, MOV, WebM"),
]

for i,(ac,n,ti,bo) in enumerate(steps):
    l = Inches(0.38) + i*Inches(3.24)
    t = Inches(2.08)
    h = Inches(5.08)
    R(s15, l, t, Inches(3.08), h, CARD)
    GR(s15, l, t, Inches(3.08), Inches(0.65), ac, dk(ac,0.5), 0)
    GR(s15, l, t, Inches(0.18), h, ac, dk(ac))
    # Step number circle
    circ = s15.shapes.add_shape(9, l+Inches(1.27), Inches(4.2), Inches(0.55), Inches(0.55))
    circ.fill.solid(); circ.fill.fore_color.rgb = ac; circ.line.color.rgb = ac
    T(s15, n, l+Inches(1.27), Inches(4.2), Inches(0.55), Inches(0.55),
      sz=18, col=DARK, bold=True, align=PP_ALIGN.CENTER)
    T(s15, ti, l+Inches(0.26), t+Inches(0.1), Inches(2.85), Inches(0.6),
      sz=12, col=DARK, bold=True, align=PP_ALIGN.CENTER)
    T(s15, bo, l+Inches(0.26), t+Inches(0.78), Inches(2.85), h-Inches(0.9),
      sz=10.5, col=LT)
    # Arrow
    if i<3:
        T(s15, "▶", l+Inches(3.12), Inches(4.0), Inches(0.22), Inches(0.6),
          sz=14, col=ORG, bold=True, align=PP_ALIGN.CENTER)

GR(s15, Inches(0.38), Inches(7.08), Inches(12.57), Inches(0.38), dk(GRN,0.28), dk(GRN,0.4), 0)
T(s15, "KEY INSIGHT: Step 3 is ITERATIVE — you cycle through Build & Tweak many times before reaching Post-Process.",
  Inches(0.55), Inches(7.1), Inches(12.2), Inches(0.34), sz=11.5, col=GRN, bold=True)

# ════════════════════════════════════════════════════
# SLIDE 16 — SUMMARY
# ════════════════════════════════════════════════════
s16 = ns()
GR(s16, Inches(0), Inches(0), Inches(13.33), Inches(1.28), CYAN, PURP, 0)
T(s16, "SUMMARY — KEY TAKEAWAYS",
  Inches(0.45), Inches(0.1), Inches(12.5), Inches(1.1),
  sz=34, col=WHT, bold=True, align=PP_ALIGN.CENTER)

left = [
    ("01",CYAN,  "Introduction",   "Animation = sequences exploiting Persistence of Vision + Phi Phenomenon to create seamless motion illusion."),
    ("02",PURP,  "Usage",          "Art, storytelling, scientific visualization, instructional. Animation is used in every major industry."),
    ("03",ORG,   "Techniques",     "Three core techniques: Cel (hand-drawn), Path (trajectory), Computer (digital — most versatile today)."),
    ("04",GRN,   "Cel Animation",  "Keyframes + tweening = smooth motion. Pencil test validates timing before expensive final production."),
    ("05",PINK,  "Path Animation", "Object follows predefined path (straight or Bezier curve) with easing for natural, weighted motion."),
]
right = [
    ("06",CYAN,  "Computer Anim",  "FK/IK chains for character rigs. Morphing for shape transforms. Far faster & more flexible than traditional."),
    ("07",PURP,  "Levels",         "Basic (GIFs/CSS) → Intermediate (motion graphics) → Advanced (feature films, real-time VFX, AI animation)."),
    ("08",ORG,   "12 Principles",  "Disney's Squash & Stretch, Anticipation, Timing, Ease, Arcs, and 8 more — still the universal standard."),
    ("09",GRN,   "Anim Space",     "2D (flat) → 2½D (pseudo-depth via parallax) → 3D (full volumetric, any camera angle, photorealism)."),
    ("  ",YEL,   "Process",        "Organize → Choose Tool → Build & Tweak (iterative) → Post-Process. Every project, every studio."),
]
for i,(n,c,ti,de) in enumerate(left):
    y = Inches(1.42)+i*Inches(1.18)
    GR(s16, Inches(0.38), y, Inches(0.72), Inches(1.0), c, dk(c,0.5), 90)
    T(s16, n, Inches(0.38), y, Inches(0.72), Inches(1.0), sz=17, col=DARK, bold=True, align=PP_ALIGN.CENTER)
    R(s16, Inches(1.15), y, Inches(5.5), Inches(1.0), CARD)
    T(s16, ti, Inches(1.22), y+Inches(0.04), Inches(5.38), Inches(0.4), sz=12, col=c, bold=True)
    T(s16, de, Inches(1.22), y+Inches(0.42), Inches(5.38), Inches(0.55), sz=10.5, col=LT)
for i,(n,c,ti,de) in enumerate(right):
    y = Inches(1.42)+i*Inches(1.18)
    GR(s16, Inches(7.05), y, Inches(0.72), Inches(1.0), c, dk(c,0.5), 90)
    T(s16, n, Inches(7.05), y, Inches(0.72), Inches(1.0), sz=17, col=DARK, bold=True, align=PP_ALIGN.CENTER)
    R(s16, Inches(7.82), y, Inches(5.13), Inches(1.0), CARD)
    T(s16, ti, Inches(7.9), y+Inches(0.04), Inches(5.0), Inches(0.4), sz=12, col=c, bold=True)
    T(s16, de, Inches(7.9), y+Inches(0.42), Inches(5.0), Inches(0.55), sz=10.5, col=LT)
R(s16, Inches(6.6), Inches(1.3), Inches(0.07), Inches(6.1), RGBColor(35,35,85))

# ════════════════════════════════════════════════════
# SLIDE 17 — THANK YOU
# ════════════════════════════════════════════════════
s17 = ns()
GR3(s17, Inches(0), Inches(0), Inches(13.33), Inches(7.5),
    RGBColor(4,4,22), RGBColor(10,5,40), RGBColor(6,6,28), 135)

# Decorative circles
for x,y,sz,c in [
    (Inches(10.5), Inches(-0.8), Inches(4.0), dk(CYAN,0.18)),
    (Inches(11.5), Inches(5.0),  Inches(2.5), dk(PURP,0.18)),
    (Inches(-0.5), Inches(4.5),  Inches(3.0), dk(ORG,0.15)),
]:
    sh = s17.shapes.add_shape(9, x, y, sz, sz)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.color.rgb = c

GR(s17, Inches(0), Inches(0), Inches(13.33), Inches(0.22), CYAN, PURP, 0)
GR(s17, Inches(0), Inches(7.28), Inches(13.33), Inches(0.22), PURP, ORG, 0)

# Watermark
T(s17, "ANIM", Inches(-0.2), Inches(0.3), Inches(8), Inches(4.5),
  sz=185, col=RGBColor(12,12,50), bold=True)

T(s17, "THANK YOU", Inches(0.5), Inches(0.9), Inches(12.33), Inches(2.3),
  sz=85, col=CYAN, bold=True, align=PP_ALIGN.CENTER)

GR(s17, Inches(2.5), Inches(3.1), Inches(8.33), Inches(0.75), dk(CYAN,0.3), dk(PURP,0.3), 0)
T(s17, "Animation: Where Imagination Meets Technology",
  Inches(2.5), Inches(3.12), Inches(8.33), Inches(0.72),
  sz=18, col=LT, italic=True, align=PP_ALIGN.CENTER)

GR(s17, Inches(3.5), Inches(4.0), Inches(6.33), Inches(0.07), CYAN, PURP, 0)

T(s17, "Group 3  •  Graphics & Multimedia  •  SMIT",
  Inches(1), Inches(4.18), Inches(11.33), Inches(0.55),
  sz=15, col=ORG, bold=True, align=PP_ALIGN.CENTER)

members = [
    ("Hashir Junaid",      "2023F-BCS-358"),
    ("Taha Haider",        "2023F-BCS-079"),
    ("Abdul Rahman Baig",  "2023F-BCS-100"),
]
for i,(name,roll) in enumerate(members):
    x = Inches(1.0)+i*Inches(3.8)
    GR(s17, x, Inches(4.88), Inches(3.55), Inches(1.0), CARD, CARD2, 0)
    T(s17, name, x, Inches(4.9), Inches(3.55), Inches(0.52),
      sz=14, col=WHT, bold=True, align=PP_ALIGN.CENTER)
    T(s17, roll, x, Inches(5.38), Inches(3.55), Inches(0.45),
      sz=11, col=MUTED, align=PP_ALIGN.CENTER)

T(s17, "Topics Covered: Introduction  •  Usage  •  Techniques  •  Types  •  Computer Animation  •  Levels  •  12 Principles  •  Space  •  Process",
  Inches(0.5), Inches(6.65), Inches(12.33), Inches(0.5),
  sz=10.5, col=RGBColor(80,90,140), italic=True, align=PP_ALIGN.CENTER)

# ── SAVE ────────────────────────────────────────────
out = r"C:\Users\Noman Traders\Desktop\SMIT -AI-CLAUDE\Graphics n Multimedia\Animation_Group3_PREMIUM.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
