"""
Animation Presentation — CLEAN version
Short text + Image placeholders + Simple diagrams
Group 3 | SMIT Graphics & Multimedia
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn as OXMLqn
from lxml import etree

# ═══ COLORS ═══════════════════════════════════════════════════════════════
BG    = RGBColor(4,   4,  18)
CARD  = RGBColor(13, 13,  48)
CARD2 = RGBColor(20, 20,  65)
CYAN  = RGBColor(0,  215, 255)
PURP  = RGBColor(155, 80, 225)
ORG   = RGBColor(255, 148,  0)
GRN   = RGBColor(46,  220, 115)
PINK  = RGBColor(255,  90, 162)
YEL   = RGBColor(255, 215,   0)
WHT   = RGBColor(255, 255, 255)
MUTED = RGBColor(105, 120, 158)
IMG   = RGBColor( 8,   8,  35)
RED   = RGBColor(255,  60,  60)

def hx(c): return f'{c[0]:02X}{c[1]:02X}{c[2]:02X}'

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ═══ HELPER FUNCTIONS ═════════════════════════════════════════════════════
def new_slide():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BG
    return sl

def box(sl, l, t, w, h, fill=CARD, lc=None, lw=None):
    sh = sl.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if lc:
        sh.line.color.rgb = lc
        if lw: sh.line.width = Pt(lw)
    else: sh.line.fill.background()
    return sh

def grad(sl, l, t, w, h, c1, c2, ang=135):
    sh = sl.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c1; sh.line.fill.background()
    spPr = sh._element.find(OXMLqn('p:spPr'))
    gx = (f'<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
          f'<a:gsLst>'
          f'<a:gs pos="0"><a:srgbClr val="{hx(c1)}"/></a:gs>'
          f'<a:gs pos="100000"><a:srgbClr val="{hx(c2)}"/></a:gs>'
          f'</a:gsLst>'
          f'<a:lin ang="{ang*60000}" scaled="0"/>'
          f'</a:gradFill>')
    gf = etree.fromstring(gx)
    sol = spPr.find(OXMLqn('a:solidFill'))
    if sol is not None: spPr.remove(sol)
    spPr.append(gf); return sh

def T(sl, text, l, t, w, h, size=14, bold=False, color=WHT, align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def bullets(sl, title, items, l, t, w, h, tc=CYAN, bc=WHT, ts=15, bs=12, card_fill=CARD, pad=0.18):
    """Card with title + bullet list."""
    box(sl, l, t, w, h, fill=card_fill)
    box(sl, l, t, Inches(0.05), h, fill=tc)  # left accent strip
    yoff = Inches(pad)
    T(sl, title, l+Inches(0.12), t+yoff, w-Inches(0.15), Inches(0.35), size=ts, bold=True, color=tc)
    yoff += Inches(0.38)
    for item in items:
        color = item[1] if isinstance(item, tuple) else bc
        text  = item[0] if isinstance(item, tuple) else item
        T(sl, text, l+Inches(0.18), t+yoff, w-Inches(0.22), Inches(0.32), size=bs, color=color)
        yoff += Inches(0.3)

def img_ph(sl, l, t, w, h, label, sub=""):
    """Image placeholder with dashed border and label."""
    box(sl, l, t, w, h, fill=IMG, lc=CYAN, lw=1.2)
    # Camera icon row
    T(sl, "[ IMAGE ]", l, t+h*0.3, w, Inches(0.4), size=11, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    T(sl, label, l, t+h*0.3+Inches(0.38), w, Inches(0.45), size=10, bold=True, color=WHT, align=PP_ALIGN.CENTER)
    if sub:
        T(sl, sub, l, t+h*0.3+Inches(0.78), w, Inches(0.38), size=9, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

def header(sl, num, title, color=CYAN):
    """Standard slide header bar."""
    grad(sl, 0, 0, prs.slide_width, Inches(0.9), RGBColor(8,8,35), BG, ang=90)
    grad(sl, 0, 0, prs.slide_width, Inches(0.07), color, PURP, ang=0)
    box(sl, Inches(0.25), Inches(0.17), Inches(0.58), Inches(0.58), fill=color)
    num_str = f"{num:02d}" if isinstance(num, int) else str(num)
    T(sl, num_str, Inches(0.25), Inches(0.17), Inches(0.58), Inches(0.58), size=20, bold=True, color=BG, align=PP_ALIGN.CENTER)
    T(sl, title.upper(), Inches(1.0), Inches(0.2), Inches(11.5), Inches(0.55), size=22, bold=True, color=color)

def footer(sl, who, slide_range):
    grad(sl, 0, Inches(7.18), prs.slide_width, Inches(0.32), RGBColor(8,8,30), BG, ang=0)
    T(sl, f"Presenter: {who}   |   Slide {slide_range}   |   Group 3 · Graphics & Multimedia · SMIT",
      Inches(0.3), Inches(7.2), Inches(12.5), Inches(0.28), size=9, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
grad(sl, 0, 0, prs.slide_width, prs.slide_height, RGBColor(4,4,18), RGBColor(10,8,35), ang=135)
# film strip holes top
for i in range(14):
    box(sl, Inches(0.3+i*0.93), Inches(0.1), Inches(0.55), Inches(0.25), fill=RGBColor(15,15,50), lc=RGBColor(30,30,80))
# film strip holes bottom
for i in range(14):
    box(sl, Inches(0.3+i*0.93), Inches(7.15), Inches(0.55), Inches(0.25), fill=RGBColor(15,15,50), lc=RGBColor(30,30,80))

# centre glow
grad(sl, Inches(2.5), Inches(1.2), Inches(8.3), Inches(5.0), RGBColor(0,30,50), BG, ang=90)

# ANIM frames visual (simple)
frame_colors = [PURP, CYAN, GRN, ORG]
frame_labels = ["f1", "f2", "f3", "f4"]
for i, (fc, fl) in enumerate(zip(frame_colors, frame_labels)):
    fx = Inches(0.8 + i*1.05)
    box(sl, fx, Inches(1.5), Inches(0.8), Inches(0.8), fill=RGBColor(8,8,40), lc=fc, lw=1.5)
    T(sl, fl, fx, Inches(1.5), Inches(0.8), Inches(0.8), size=14, bold=True, color=fc, align=PP_ALIGN.CENTER)
    if i < 3:
        T(sl, "▶", Inches(0.8+i*1.05+0.82), Inches(1.7), Inches(0.2), Inches(0.4), size=10, color=MUTED, align=PP_ALIGN.CENTER)

T(sl, "Still frames  →  Illusion of Motion  →", Inches(0.6), Inches(2.5), Inches(4.5), Inches(0.4), size=10, color=MUTED)
T(sl, "ANIMATION", Inches(5.3), Inches(2.45), Inches(7.0), Inches(0.5), size=22, bold=True, color=CYAN)

# Main title
grad(sl, Inches(2.0), Inches(3.0), Inches(9.3), Inches(0.07), CYAN, PURP, ang=0)
T(sl, "ANIMATION", Inches(2.0), Inches(3.1), Inches(9.3), Inches(1.35), size=68, bold=True, color=WHT, align=PP_ALIGN.CENTER)
T(sl, "The Art & Science of Bringing Images to Life", Inches(2.0), Inches(4.5), Inches(9.3), Inches(0.5), size=17, color=CYAN, align=PP_ALIGN.CENTER, italic=True)
grad(sl, Inches(2.0), Inches(5.08), Inches(9.3), Inches(0.05), CYAN, PURP, ang=0)

T(sl, "Graphics & Multimedia  •  SMIT  •  Group 3  •  Topics 1—9", Inches(2.0), Inches(5.2), Inches(9.3), Inches(0.35), size=11, color=MUTED, align=PP_ALIGN.CENTER)

# Names
names = [("Hashir Junaid", "2023F-BCS-358", CYAN), ("Taha Haider", "2023F-BCS-079", PURP), ("Abdul Rahman Baig", "2023F-BCS-100", GRN)]
for i, (name, roll, nc) in enumerate(names):
    nx = Inches(1.1 + i*4.1)
    box(sl, nx, Inches(6.1), Inches(3.6), Inches(1.05), fill=RGBColor(10,10,40), lc=nc, lw=1.2)
    T(sl, name, nx, Inches(6.15), Inches(3.6), Inches(0.42), size=13, bold=True, color=nc, align=PP_ALIGN.CENTER)
    T(sl, roll, nx, Inches(6.57), Inches(3.6), Inches(0.35), size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — TOPICS OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, 0, "What We'll Cover Today", CYAN)
T(sl, "9 Topics  •  Animation: Complete Foundation  •  Group 3", Inches(1.0), Inches(0.75), Inches(10), Inches(0.32), size=11, color=MUTED)

topics = [
    ("01", "Introduction to Animation", CYAN),
    ("02", "Usage of Animation", PURP),
    ("03", "Animation Techniques", GRN),
    ("04", "Types of Animation", ORG),
    ("05", "Computer Animation", PINK),
    ("06", "Levels of Animation", YEL),
    ("07", "12 Basic Principles", CYAN),
    ("08", "Animation Space", PURP),
    ("09", "Animation Process", GRN),
]
cols, rows = 3, 3
for idx, (num, topic, tc) in enumerate(topics):
    c, r = idx % cols, idx // cols
    tx = Inches(0.3 + c*4.33)
    ty = Inches(1.2 + r*1.9)
    box(sl, tx, ty, Inches(4.1), Inches(1.7), fill=CARD, lc=tc, lw=1)
    grad(sl, tx, ty, Inches(0.08), Inches(1.7), tc, CARD, ang=90)
    T(sl, num, tx+Inches(0.15), ty+Inches(0.15), Inches(0.5), Inches(0.5), size=20, bold=True, color=tc)
    T(sl, topic, tx+Inches(0.15), ty+Inches(0.7), Inches(3.8), Inches(0.7), size=13, bold=True, color=WHT)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — INTRODUCTION TO ANIMATION  (Hashir)
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, 1, "Introduction to Animation", CYAN)
footer(sl, "Hashir Junaid", "3 / 17")

# Left column — 3 content cards
CW = Inches(5.2); CH = Inches(1.65)
CX = Inches(0.3); CY_start = Inches(1.05)

# Card 1 — Definition
box(sl, CX, CY_start, CW, CH, fill=CARD, lc=CYAN, lw=0.8)
box(sl, CX, CY_start, Inches(0.07), CH, fill=CYAN)
T(sl, "DEFINITION", CX+Inches(0.15), CY_start+Inches(0.1), CW, Inches(0.35), size=13, bold=True, color=CYAN)
T(sl, "• Still images shown rapidly → illusion of motion", CX+Inches(0.2), CY_start+Inches(0.48), CW-Inches(0.25), Inches(0.3), size=11, color=WHT)
T(sl, "• Brain is 'tricked' by two biological effects", CX+Inches(0.2), CY_start+Inches(0.78), CW-Inches(0.25), Inches(0.3), size=11, color=WHT)
T(sl, "• From Latin 'anima' = to give soul", CX+Inches(0.2), CY_start+Inches(1.08), CW-Inches(0.25), Inches(0.3), size=11, color=MUTED, italic=True)
T(sl, "First animation: Muybridge horse gallop (1878)", CX+Inches(0.2), CY_start+Inches(1.35), CW-Inches(0.25), Inches(0.28), size=10, color=MUTED, italic=True)

# Card 2 — The Two Brain Effects
CY2 = CY_start + CH + Inches(0.12)
box(sl, CX, CY2, CW, Inches(2.25), fill=CARD, lc=PURP, lw=0.8)
box(sl, CX, CY2, Inches(0.07), Inches(2.25), fill=PURP)
T(sl, "THE TWO BRAIN EFFECTS", CX+Inches(0.15), CY2+Inches(0.1), CW, Inches(0.35), size=13, bold=True, color=PURP)

# effect boxes inside
EW = Inches(2.3)
# Persistence of Vision
box(sl, CX+Inches(0.15), CY2+Inches(0.5), EW, Inches(1.55), fill=RGBColor(20,10,50), lc=PURP, lw=0.6)
T(sl, "PERSISTENCE OF VISION", CX+Inches(0.22), CY2+Inches(0.55), EW-Inches(0.1), Inches(0.3), size=10, bold=True, color=PURP)
T(sl, "Physical effect in the eye", CX+Inches(0.22), CY2+Inches(0.85), EW-Inches(0.1), Inches(0.28), size=9.5, color=WHT)
T(sl, "Image stays ~50ms after disappearing", CX+Inches(0.22), CY2+Inches(1.1), EW-Inches(0.1), Inches(0.28), size=9.5, color=WHT)
T(sl, "Afterimages overlap → smooth flow", CX+Inches(0.22), CY2+Inches(1.35), EW-Inches(0.1), Inches(0.28), size=9.5, color=WHT)
T(sl, "Discovered: Peter Roget (1824)", CX+Inches(0.22), CY2+Inches(1.62), EW-Inches(0.1), Inches(0.25), size=8.5, color=MUTED, italic=True)

# Phi Phenomenon
EX2 = CX + Inches(0.15) + EW + Inches(0.2)
box(sl, EX2, CY2+Inches(0.5), EW, Inches(1.55), fill=RGBColor(20,10,50), lc=CYAN, lw=0.6)
T(sl, "PHI PHENOMENON", EX2+Inches(0.07), CY2+Inches(0.55), EW-Inches(0.1), Inches(0.3), size=10, bold=True, color=CYAN)
T(sl, "Psychological brain effect", EX2+Inches(0.07), CY2+Inches(0.85), EW-Inches(0.1), Inches(0.28), size=9.5, color=WHT)
T(sl, "Brain infers motion between objects", EX2+Inches(0.07), CY2+Inches(1.1), EW-Inches(0.1), Inches(0.28), size=9.5, color=WHT)
T(sl, "Examples: flip-books, neon arrows", EX2+Inches(0.07), CY2+Inches(1.35), EW-Inches(0.1), Inches(0.28), size=9.5, color=WHT)
T(sl, "Discovered: Wertheimer (1912)", EX2+Inches(0.07), CY2+Inches(1.62), EW-Inches(0.1), Inches(0.25), size=8.5, color=MUTED, italic=True)

# Card 3 — FPS
CY3 = CY2 + Inches(2.25) + Inches(0.12)
box(sl, CX, CY3, CW, Inches(1.25), fill=CARD, lc=GRN, lw=0.8)
box(sl, CX, CY3, Inches(0.07), Inches(1.25), fill=GRN)
T(sl, "FPS STANDARDS", CX+Inches(0.15), CY3+Inches(0.1), CW, Inches(0.35), size=13, bold=True, color=GRN)
fps_items = [("24fps", "Cinema Gold Standard", CYAN), ("30fps", "TV Broadcast", ORG), ("60fps", "Modern Gaming", GRN), ("< 12fps", "Visible Flicker!", RED)]
for i, (fps, label, fc) in enumerate(fps_items):
    ix = CX + Inches(0.18 + i*1.25)
    box(sl, ix, CY3+Inches(0.5), Inches(1.15), Inches(0.58), fill=RGBColor(15,15,50), lc=fc, lw=0.8)
    T(sl, fps, ix, CY3+Inches(0.52), Inches(1.15), Inches(0.28), size=13, bold=True, color=fc, align=PP_ALIGN.CENTER)
    T(sl, label, ix, CY3+Inches(0.78), Inches(1.15), Inches(0.25), size=8, color=MUTED, align=PP_ALIGN.CENTER)

# Right column — Image placeholder + simple frame diagram
RX = CX + CW + Inches(0.25)
RW = prs.slide_width - RX - Inches(0.25)

img_ph(sl, RX, Inches(1.05), RW, Inches(2.2), "Flip-book / Eye Diagram", "Show how frames create motion")

# Simple drawn diagram — frames → motion
DY = Inches(3.4)
T(sl, "How Animation Works:", RX, DY, RW, Inches(0.32), size=11, bold=True, color=CYAN)
DY += Inches(0.35)
frame_x = RX
for i in range(5):
    fc = [CYAN, PURP, GRN, ORG, PINK][i]
    box(sl, frame_x, DY, Inches(1.2), Inches(1.0), fill=RGBColor(10,10,45), lc=fc, lw=1.2)
    T(sl, f"Frame\n{i+1}", frame_x, DY, Inches(1.2), Inches(1.0), size=11, bold=True, color=fc, align=PP_ALIGN.CENTER)
    frame_x += Inches(1.22)

T(sl, "Shown at 24+ fps  →  BRAIN sees continuous motion!", RX, DY+Inches(1.08), RW, Inches(0.35), size=10, bold=True, color=YEL)
T(sl, "COMBINED = SEAMLESS ANIMATION", RX, DY+Inches(1.45), RW, Inches(0.32), size=11, bold=True, color=WHT)

img_ph(sl, RX, Inches(5.35), RW, Inches(1.7), "Phi Phenomenon Diagram", "Two stationary dots → brain sees movement")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — USAGE OF ANIMATION  (Hashir)
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, 2, "Usage of Animation", PURP)
footer(sl, "Hashir Junaid", "4 / 17")

usage = [
    ("ARTISTIC", CYAN,
     ["• Fine art & experimental films", "• Music video production", "• Title sequence design"],
     ["Disney Fantasia (1940)", "Spider-Verse (Oscar 2018)", "Yellow Submarine"]),
    ("STORYTELLING", PURP,
     ["• Feature films & TV series", "• Video game cutscenes", "• Character-driven dramas"],
     ["Pixar Toy Story", "Studio Ghibli", "Arcane – Netflix (2021)"]),
    ("SCIENTIFIC", GRN,
     ["• Human anatomy explainers", "• Weather & climate models", "• DNA & molecular modeling"],
     ["NASA mission visualizations", "Medical surgical training", "Climate simulations"]),
    ("INSTRUCTIONAL", ORG,
     ["• Educational explainer videos", "• Military & flight simulators", "• Safety procedure demos"],
     ["Khan Academy lessons", "Boeing pilot simulators", "IKEA assembly guides"]),
]

for idx, (title, tc, bullets_list, examples) in enumerate(usage):
    c, r = idx % 2, idx // 2
    bx = Inches(0.3 + c*6.5)
    by = Inches(1.05 + r*3.0)
    BW = Inches(6.2); BH = Inches(2.75)
    box(sl, bx, by, BW, BH, fill=CARD, lc=tc, lw=1.2)
    grad(sl, bx, by, BW, Inches(0.08), tc, CARD, ang=0)
    T(sl, title, bx+Inches(0.18), by+Inches(0.15), BW-Inches(0.2), Inches(0.42), size=16, bold=True, color=tc)
    for i, b in enumerate(bullets_list):
        T(sl, b, bx+Inches(0.18), by+Inches(0.65+i*0.38), BW*0.52, Inches(0.35), size=11, color=WHT)
    # Examples box
    EBX = bx + BW*0.54
    box(sl, EBX, by+Inches(0.55), BW*0.42, Inches(1.7), fill=RGBColor(10,10,40), lc=tc, lw=0.6)
    T(sl, "Examples:", EBX+Inches(0.08), by+Inches(0.6), BW*0.4, Inches(0.28), size=10, bold=True, color=tc)
    for i, ex in enumerate(examples):
        T(sl, f"• {ex}", EBX+Inches(0.08), by+Inches(0.9+i*0.38), BW*0.4, Inches(0.35), size=9.5, color=MUTED)

img_ph(sl, Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.12), "", "")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ANIMATION TECHNIQUES  (Hashir)
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, 3, "Animation Techniques", GRN)
footer(sl, "Hashir Junaid", "5 / 17")

techniques = [
    ("CEL ANIMATION", "1910s – 1990s", CYAN,
     ["Draw characters on transparent", "celluloid (plastic) sheets",
      "• Each sheet = 1 frame", "• Layer over reused backgrounds",
      "• Hand-filmed, frame by frame"],
     "Tom & Jerry · Snow White\nLooney Tunes · Bambi · Akira",
     "Cel layers stacked over BG"),
    ("PATH ANIMATION", "1970s – Present", PURP,
     ["Animator draws the path first", "Object follows it automatically",
      "• Straight lines or Bezier curves", "• Set easing (slow in/out)",
      "• Object can resize along path"],
     "CSS animations · After Effects\nUnity NavMesh · SVG paths",
     "Object on curved path with arrows"),
    ("COMPUTER ANIMATION", "1980s – Present", GRN,
     ["Creates motion entirely in software", "Can simulate both Cel & Path",
      "• 2D digital, 3D CG, or hybrid", "• Physics simulation built-in",
      "• Real-time or pre-rendered"],
     "Pixar · DreamWorks · ILM\nFortnite · Unreal Engine",
     "3D software interface screenshot"),
]

CW3 = Inches(4.2)
for i, (title, era, tc, pts, examples, img_label) in enumerate(techniques):
    cx = Inches(0.28 + i*4.35)
    cy = Inches(1.05)
    CH3 = Inches(6.25)
    box(sl, cx, cy, CW3, CH3, fill=CARD)
    grad(sl, cx, cy, CW3, Inches(0.07), tc, CARD, ang=0)
    # Era badge
    box(sl, cx+Inches(0.1), cy+Inches(0.12), Inches(2.5), Inches(0.32), fill=RGBColor(10,10,40), lc=tc, lw=0.5)
    T(sl, era, cx+Inches(0.15), cy+Inches(0.12), Inches(2.45), Inches(0.32), size=9, color=tc)
    T(sl, title, cx+Inches(0.12), cy+Inches(0.5), CW3-Inches(0.2), Inches(0.45), size=14, bold=True, color=tc)
    # bullets
    for j, pt in enumerate(pts):
        T(sl, pt, cx+Inches(0.18), cy+Inches(1.05+j*0.35), CW3-Inches(0.25), Inches(0.32), size=10.5, color=WHT)
    # Examples
    box(sl, cx+Inches(0.1), cy+Inches(3.0), CW3-Inches(0.2), Inches(0.75), fill=RGBColor(8,8,35), lc=tc, lw=0.6)
    T(sl, "Used In:", cx+Inches(0.18), cy+Inches(3.05), CW3-Inches(0.25), Inches(0.25), size=9, bold=True, color=tc)
    T(sl, examples, cx+Inches(0.18), cy+Inches(3.28), CW3-Inches(0.25), Inches(0.45), size=9, color=MUTED)
    # Image placeholder
    img_ph(sl, cx+Inches(0.1), cy+Inches(3.9), CW3-Inches(0.2), Inches(2.15), img_label)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — CEL ANIMATION DETAILS  (Taha)
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "4A", "Types of Animation — Cel Animation", ORG)
footer(sl, "Taha Haider", "6 / 17")
T(sl, "Named after CELLULOID — transparent plastic sheets. Raoul Barré & John Bray developed the technique ~1914.", Inches(1.0), Inches(0.77), Inches(11), Inches(0.32), size=10, color=MUTED, italic=True)

LW = Inches(5.0)
RX6 = Inches(5.55)
RW6 = Inches(7.5)

cel_cards = [
    ("CEL SHEET", ORG, ["• 1 transparent celluloid sheet = 1 frame", "• Character on cel, background separate", "• Background reused to save cost", "• Colors applied to BACK of sheet", "• 12–24 unique cels per second"]),
    ("KEYFRAMES", YEL, ["• The 'anchor' extreme poses of an action", "• Drawn by SENIOR animators", "• Define timing and structure", "• Example: Frame 1=standing, 12=peak jump, 24=landed"]),
    ("TWEENING (In-betweening)", CYAN, ["• Frames drawn BETWEEN keyframes", "• Done by JUNIOR animators", "• More tweens = smoother motion", "• Modern software auto-tweens between keyframes"]),
    ("PENCIL TEST", GRN, ["• Rough pencil sketches filmed and played back", "• Checks timing BEFORE expensive inking/painting", "• Modern version: 'animatic' preview", "• Every major studio uses this"]),
]

card_h = Inches(1.45)
for i, (title, tc, pts) in enumerate(cel_cards):
    cy = Inches(1.12 + i*1.52)
    box(sl, Inches(0.28), cy, LW, card_h, fill=CARD, lc=tc, lw=0.8)
    box(sl, Inches(0.28), cy, Inches(0.07), card_h, fill=tc)
    T(sl, title, Inches(0.42), cy+Inches(0.1), LW-Inches(0.18), Inches(0.3), size=12, bold=True, color=tc)
    for j, pt in enumerate(pts):
        T(sl, pt, Inches(0.42), cy+Inches(0.42+j*0.27), LW-Inches(0.18), Inches(0.25), size=9.5, color=WHT)

# Right — Cel stack diagram
T(sl, "HOW CELS STACK:", RX6, Inches(1.1), RW6, Inches(0.35), size=13, bold=True, color=ORG)

stack_layers = [("CEL 3 — Top Character Layer", PINK), ("CEL 2 — Mid Character Detail", CYAN), ("CEL 1 — Base Character Outline", PURP), ("BACKGROUND — Static, Reused", GRN)]
for i, (label, lc) in enumerate(stack_layers):
    ly = Inches(1.55 + i*0.7)
    box(sl, RX6, ly, RW6-Inches(0.2), Inches(0.58), fill=RGBColor(10,10,45), lc=lc, lw=1.2)
    T(sl, label, RX6+Inches(0.15), ly+Inches(0.15), RW6-Inches(0.5), Inches(0.3), size=11, bold=True, color=lc)

# Arrow
T(sl, "↓ LAYERED TOGETHER ↓", RX6, Inches(4.45), RW6-Inches(0.2), Inches(0.35), size=12, bold=True, color=ORG, align=PP_ALIGN.CENTER)
box(sl, RX6, Inches(4.85), RW6-Inches(0.2), Inches(0.55), fill=RGBColor(15,15,60), lc=ORG, lw=1.5)
T(sl, "= ONE FRAME of the film", RX6+Inches(0.1), Inches(4.9), RW6-Inches(0.3), Inches(0.42), size=13, bold=True, color=WHT)

img_ph(sl, RX6, Inches(5.55), RW6-Inches(0.2), Inches(1.72), "Tom & Jerry / Snow White clip", "Classic cel animation example")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PATH ANIMATION  (Taha)
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "4B", "Types of Animation — Path Animation", PURP)
footer(sl, "Taha Haider", "7 / 17")
T(sl, "Object follows a predefined trajectory. Animator draws the path; software moves the object automatically.", Inches(1.0), Inches(0.77), Inches(11), Inches(0.32), size=10, color=MUTED, italic=True)

# Left — 6 steps
steps = [
    ("1", "Draw the PATH", "Straight line, Bezier curve, or freehand", CYAN),
    ("2", "Set START & END points", "Define where the object begins and ends", PURP),
    ("3", "Set DURATION", "How long the journey takes (seconds)", GRN),
    ("4", "Apply EASING", "Slow in / slow out for natural feel", ORG),
    ("5", "Set ORIENTATION", "Does object face direction of travel?", PINK),
    ("6", "Preview & Adjust", "Repeat until satisfied, then render", YEL),
]

for i, (num, title, desc, sc) in enumerate(steps):
    sy = Inches(1.1 + i*1.02)
    box(sl, Inches(0.28), sy, Inches(5.5), Inches(0.88), fill=CARD)
    box(sl, Inches(0.28), sy, Inches(0.42), Inches(0.88), fill=sc)
    T(sl, num, Inches(0.28), sy, Inches(0.42), Inches(0.88), size=18, bold=True, color=BG, align=PP_ALIGN.CENTER)
    T(sl, title, Inches(0.78), sy+Inches(0.08), Inches(4.9), Inches(0.32), size=12, bold=True, color=sc)
    T(sl, desc,  Inches(0.78), sy+Inches(0.42), Inches(4.9), Inches(0.38), size=10, color=MUTED)

# Right — Path types
RX7 = Inches(6.0)
RW7 = Inches(7.05)
T(sl, "PATH TYPES:", RX7, Inches(1.05), RW7, Inches(0.35), size=14, bold=True, color=PURP)

path_types = [
    ("STRAIGHT LINE", CYAN, "Mechanical movement, text reveals, UI transitions"),
    ("BEZIER CURVE", PURP, "Natural organic paths — balls, camera, flying objects"),
    ("EASING CURVE", GRN, "Controls speed: slow start → full speed → slow end"),
]
for i, (pt, pc, pdesc) in enumerate(path_types):
    py = Inches(1.5 + i*1.0)
    box(sl, RX7, py, RW7-Inches(0.2), Inches(0.85), fill=CARD, lc=pc, lw=0.8)
    box(sl, RX7, py, Inches(0.06), Inches(0.85), fill=pc)
    T(sl, pt,    RX7+Inches(0.15), py+Inches(0.08), RW7-Inches(0.35), Inches(0.3), size=12, bold=True, color=pc)
    T(sl, pdesc, RX7+Inches(0.15), py+Inches(0.42), RW7-Inches(0.35), Inches(0.35), size=10, color=WHT)

img_ph(sl, RX7, Inches(4.65), RW7-Inches(0.2), Inches(2.65), "Path Animation Diagram", "Ball following Bezier curve path with start→end dots")

T(sl, "Used in: CSS animations  •  After Effects  •  Unity  •  SVG paths  •  Game projectiles",
  RX7, Inches(7.2), RW7, Inches(0.28), size=9.5, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — COMPUTER ANIMATION  (Taha)
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, 5, "Computer Animation", PINK)
footer(sl, "Taha Haider", "8 / 17")
T(sl, "Digital software creates motion — simulating traditional techniques or producing entirely new effects impossible by hand.", Inches(1.0), Inches(0.77), Inches(11), Inches(0.32), size=10, color=MUTED, italic=True)

# Table header
box(sl, Inches(0.3), Inches(1.1), Inches(6.0), Inches(0.5), fill=RGBColor(10,15,50), lc=ORG, lw=1)
box(sl, Inches(6.55), Inches(1.1), Inches(6.5), Inches(0.5), fill=RGBColor(10,15,50), lc=PINK, lw=1)
T(sl, "CEL ANIMATION", Inches(0.3), Inches(1.12), Inches(6.0), Inches(0.45), size=16, bold=True, color=ORG, align=PP_ALIGN.CENTER)
T(sl, "COMPUTER ANIMATION", Inches(6.55), Inches(1.12), Inches(6.5), Inches(0.45), size=16, bold=True, color=PINK, align=PP_ALIGN.CENTER)

# VS badge
box(sl, Inches(6.05), Inches(1.1), Inches(0.45), Inches(0.5), fill=YEL)
T(sl, "VS", Inches(6.05), Inches(1.12), Inches(0.45), Inches(0.44), size=12, bold=True, color=BG, align=PP_ALIGN.CENTER)

comparison = [
    ("Hand-drawn on celluloid sheets", "Created entirely in digital software", ORG, PINK),
    ("Every frame drawn manually (slow!)", "Automation: tweens, physics, rendering", ORG, PINK),
    ("Physical storage — film cans, warehouses", "Digital files — backup & share instantly", ORG, PINK),
    ("Expensive: cels, paint, camera, lab fees", "One-time software cost, no consumables", ORG, PINK),
    ("Months per 1 minute of finished film", "Real-time previews — much faster iterations", ORG, PINK),
    ("Organic, hand-crafted visual warmth", "Any style: cartoon to photorealistic", ORG, PINK),
]

for i, (cel_txt, comp_txt, lc1, lc2) in enumerate(comparison):
    ry = Inches(1.7 + i*0.85)
    fill_c = CARD if i % 2 == 0 else CARD2
    box(sl, Inches(0.3), ry, Inches(6.0), Inches(0.78), fill=fill_c, lc=lc1, lw=0.4)
    box(sl, Inches(6.55), ry, Inches(6.5), Inches(0.78), fill=fill_c, lc=lc2, lw=0.4)
    T(sl, cel_txt,  Inches(0.45), ry+Inches(0.22), Inches(5.7), Inches(0.35), size=11.5, color=WHT)
    T(sl, comp_txt, Inches(6.7),  ry+Inches(0.22), Inches(6.2), Inches(0.35), size=11.5, color=WHT)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — KINEMATICS & MORPHING  (Taha)
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "5B", "Kinematics & Morphing", CYAN)
footer(sl, "Taha Haider", "9 / 17")

KW = Inches(4.1)
kin_data = [
    ("FORWARD KINEMATICS (FK)", CYAN,
     ["Move PARENT joint → children follow", "Shoulder → Elbow → Wrist → Finger", "Like a puppet — control parent, rest follows"],
     ["Robotic / mechanical arms", "Waving & swinging motions", "Tail & tentacle animation"],
     "Arm chain diagram\n(shoulder→elbow→wrist)"),
    ("INVERSE KINEMATICS (IK)", PURP,
     ["Move END POINT → chain auto-adjusts", "Place hand → software calculates elbow & shoulder", "Like reaching for a cup — aim hand, body figures out joints"],
     ["Character walking & running", "Reaching & grabbing actions", "Real-time game characters"],
     "IK reach diagram\n(hand position → full arm)"),
    ("MORPHING", GRN,
     ["Smooth shape/image transformation", "Map matching points on source & target", "Software interpolates all points together"],
     ["MJ 'Black or White' (1991)", "T-1000 — Terminator 2", "Logo reveal animations"],
     "Morphing sequence:\nface A → face B"),
]

for i, (title, tc, how, uses, img_label) in enumerate(kin_data):
    kx = Inches(0.28 + i*4.37)
    ky = Inches(1.05)
    box(sl, kx, ky, KW, Inches(6.3), fill=CARD)
    grad(sl, kx, ky, KW, Inches(0.07), tc, CARD, ang=0)
    T(sl, title, kx+Inches(0.12), ky+Inches(0.15), KW-Inches(0.2), Inches(0.42), size=12, bold=True, color=tc)
    T(sl, "HOW IT WORKS:", kx+Inches(0.12), ky+Inches(0.65), KW-Inches(0.2), Inches(0.28), size=10, bold=True, color=tc)
    for j, h in enumerate(how):
        T(sl, f"• {h}", kx+Inches(0.18), ky+Inches(0.95+j*0.4), KW-Inches(0.25), Inches(0.35), size=9.5, color=WHT)
    T(sl, "BEST FOR:", kx+Inches(0.12), ky+Inches(2.25), KW-Inches(0.2), Inches(0.28), size=10, bold=True, color=tc)
    for j, u in enumerate(uses):
        T(sl, f"→ {u}", kx+Inches(0.18), ky+Inches(2.55+j*0.38), KW-Inches(0.25), Inches(0.35), size=9.5, color=MUTED)
    img_ph(sl, kx+Inches(0.1), ky+Inches(3.8), KW-Inches(0.2), Inches(2.35), img_label)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — LEVELS OF COMPUTER ANIMATION  (Taha)
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, 6, "Levels of Computer Animation", YEL)
footer(sl, "Taha Haider", "10 / 17")

levels = [
    ("BASIC", "★☆☆", GRN,
     ["Animated GIFs & simple sprites", "CSS & JavaScript transitions", "PowerPoint animations", "Basic logo animations"],
     ["MS PowerPoint", "Canva", "GIMP", "Basic CSS/JS"],
     "Days to Weeks", "Beginner",
     "Social media content\nMarketing materials"),
    ("INTERMEDIATE", "★★☆", CYAN,
     ["2D character animation", "Motion graphics & infographics", "Complex path animations", "Basic 3D object animation"],
     ["Adobe After Effects", "Adobe Animate", "Blender (basics)", "Spine"],
     "Several Months", "Intermediate",
     "Motion designer\n2D animator · Game dev"),
    ("ADVANCED", "★★★", PURP,
     ["Full CG feature films", "Realistic physics simulations", "Motion capture integration", "Real-time VFX in games"],
     ["Autodesk Maya", "Houdini", "Unreal Engine", "Cinema 4D"],
     "Multiple Years", "Professional",
     "VFX artist · 3D animator\nTechnical Director"),
]

LVW = Inches(4.1)
for i, (title, stars, tc, what, tools, time_to_learn, skill, career) in enumerate(levels):
    lx = Inches(0.28 + i*4.37)
    ly = Inches(1.05)
    box(sl, lx, ly, LVW, Inches(6.28), fill=CARD)
    grad(sl, lx, ly, LVW, Inches(0.07), tc, CARD, ang=0)
    # Title
    T(sl, stars, lx+Inches(0.12), ly+Inches(0.12), LVW-Inches(0.2), Inches(0.35), size=15, color=tc)
    T(sl, title, lx+Inches(0.12), ly+Inches(0.48), LVW-Inches(0.2), Inches(0.38), size=16, bold=True, color=tc)
    # What you create
    T(sl, "WHAT YOU CREATE:", lx+Inches(0.12), ly+Inches(0.9), LVW-Inches(0.2), Inches(0.28), size=10, bold=True, color=tc)
    for j, w in enumerate(what):
        T(sl, f"• {w}", lx+Inches(0.18), ly+Inches(1.2+j*0.35), LVW-Inches(0.25), Inches(0.32), size=9.5, color=WHT)
    # Tools
    T(sl, "TOOLS:", lx+Inches(0.12), ly+Inches(2.62), LVW-Inches(0.2), Inches(0.28), size=10, bold=True, color=tc)
    for j, tool in enumerate(tools[:3]):
        T(sl, f"→ {tool}", lx+Inches(0.18), ly+Inches(2.92+j*0.33), LVW-Inches(0.25), Inches(0.3), size=9.5, color=MUTED)
    # Stats
    box(sl, lx+Inches(0.1), ly+Inches(4.05), LVW-Inches(0.2), Inches(0.58), fill=RGBColor(8,8,35), lc=tc, lw=0.5)
    T(sl, f"Time: {time_to_learn}  |  Skill: {skill}", lx+Inches(0.15), ly+Inches(4.1), LVW-Inches(0.25), Inches(0.25), size=9, color=tc)
    T(sl, f"Career: {career}", lx+Inches(0.15), ly+Inches(4.35), LVW-Inches(0.25), Inches(0.25), size=9, color=MUTED)
    img_ph(sl, lx+Inches(0.1), ly+Inches(4.72), LVW-Inches(0.2), Inches(1.42), f"{title} level tool screenshot")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — 12 PRINCIPLES Part 1  (Abdul Rahman)
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, 7, "12 Principles of Animation — Part 1 of 3", CYAN)
footer(sl, "Abdul Rahman Baig", "11 / 17")
T(sl, "Developed by Disney masters Ollie Johnston & Frank Thomas (1981). Still the universal GOLD STANDARD.", Inches(1.0), Inches(0.77), Inches(11), Inches(0.32), size=10, color=MUTED, italic=True)

p1 = [
    ("01  SQUASH & STRETCH", CYAN,
     ["Object squashes on IMPACT", "Object stretches when MOVING fast", "GOLDEN RULE: Volume stays constant!", "Shows elasticity, mass, and life"],
     "Most fundamental principle",
     "Rubber ball bounce showing squash on landing, stretch in air"),
    ("02  ANTICIPATION", PURP,
     ["Small opposite action BEFORE main action", "Character crouches before jumping", "Arm swings back before throwing", "Prepares audience mentally"],
     "Without it, actions feel sudden",
     "Bugs Bunny crouching before leap"),
    ("03  STAGING", GRN,
     ["Present ONE clear idea at a time", "Use clear readable silhouettes", "Camera angle serves the story", "Direct the audience's eye"],
     "Think: theater staging, not photography",
     "Character profile with dramatic gesture"),
    ("04  STRAIGHT AHEAD vs POSE-TO-POSE", ORG,
     ["STRAIGHT AHEAD: draw frame by frame → fluid & spontaneous",
      "POSE-TO-POSE: draw key poses first → controlled & consistent",
      "Best work COMBINES both methods"],
     "Straight ahead = fire/water | P-to-P = character acting",
     "Side-by-side comparison of both methods"),
]

for i, (title, tc, pts, note, img_label) in enumerate(p1):
    c, r = i % 2, i // 2
    px = Inches(0.28 + c*6.52)
    py = Inches(1.1 + r*2.85)
    PW = Inches(6.28); PH = Inches(2.65)
    box(sl, px, py, PW, PH, fill=CARD, lc=tc, lw=1)
    grad(sl, px, py, PW, Inches(0.06), tc, CARD, ang=0)
    T(sl, title, px+Inches(0.12), py+Inches(0.12), PW*0.55, Inches(0.38), size=12, bold=True, color=tc)
    for j, pt in enumerate(pts):
        T(sl, f"• {pt}", px+Inches(0.15), py+Inches(0.55+j*0.37), PW*0.52, Inches(0.34), size=9.5, color=WHT)
    T(sl, note, px+Inches(0.15), py+PH-Inches(0.32), PW*0.52, Inches(0.28), size=9, color=MUTED, italic=True)
    # Image placeholder right side of each card
    img_ph(sl, px+PW*0.56, py+Inches(0.12), PW*0.41, PH-Inches(0.2), img_label)

# Principle counter bottom
T(sl, "Principles 1–4 of 12  |  Continued next slide →", Inches(0.3), Inches(7.22), Inches(12), Inches(0.25), size=9, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — 12 PRINCIPLES Part 2  (Abdul Rahman)
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, 7, "12 Principles of Animation — Part 2 of 3", PURP)
footer(sl, "Abdul Rahman Baig", "12 / 17")
T(sl, "These principles govern how objects move AFTER the main action — the details that separate amateur from professional.", Inches(1.0), Inches(0.77), Inches(11), Inches(0.32), size=10, color=MUTED, italic=True)

p2 = [
    ("05  FOLLOW THROUGH & OVERLAP", CYAN,
     ["Parts keep moving AFTER main action stops", "Hair swings after head stops", "OVERLAP: different parts move at different rates", "Makes motion feel organic & believable"],
     "Coat settles after character stops",
     "Character stopping, coat still swinging"),
    ("06  EASE IN & EASE OUT", PURP,
     ["Slow start + slow end for every movement", "Nothing starts at full speed or stops instantly", "A car: slow → full speed → slow stop", "Linear = robotic | Ease = natural"],
     "Most impactful for 'weight' feeling",
     "Bezier easing curve graph"),
    ("07  ARCS", GRN,
     ["All natural motion follows ARC paths", "Arms, heads, eyes — all move in arcs", "Thrown objects follow a parabola", "EXCEPTION: robots move in straight lines"],
     "Breaking arcs = instant 'robot' look",
     "Arm swing arc diagram"),
    ("08  SECONDARY ACTION", ORG,
     ["Supporting actions that ENRICH main movement", "Walking (main) + arms swing (secondary)", "Running (main) + hair flows (secondary)", "MUST NOT overpower the main action"],
     "It supports — never steals focus",
     "Character walking with secondary elements"),
]

for i, (title, tc, pts, note, img_label) in enumerate(p2):
    c, r = i % 2, i // 2
    px = Inches(0.28 + c*6.52)
    py = Inches(1.1 + r*2.85)
    PW = Inches(6.28); PH = Inches(2.65)
    box(sl, px, py, PW, PH, fill=CARD, lc=tc, lw=1)
    grad(sl, px, py, PW, Inches(0.06), tc, CARD, ang=0)
    T(sl, title, px+Inches(0.12), py+Inches(0.12), PW*0.55, Inches(0.38), size=12, bold=True, color=tc)
    for j, pt in enumerate(pts):
        T(sl, f"• {pt}", px+Inches(0.15), py+Inches(0.55+j*0.37), PW*0.52, Inches(0.34), size=9.5, color=WHT)
    T(sl, note, px+Inches(0.15), py+PH-Inches(0.32), PW*0.52, Inches(0.28), size=9, color=MUTED, italic=True)
    img_ph(sl, px+PW*0.56, py+Inches(0.12), PW*0.41, PH-Inches(0.2), img_label)

T(sl, "Principles 5–8 of 12  |  Continued next slide →", Inches(0.3), Inches(7.22), Inches(12), Inches(0.25), size=9, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — 12 PRINCIPLES Part 3  (Abdul Rahman)
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, 7, "12 Principles of Animation — Part 3 of 3", GRN)
footer(sl, "Abdul Rahman Baig", "13 / 17")
T(sl, "Walt Disney: 'The secret of animation is the exaggeration of the truth.' These final 4 deliver exactly that truth.", Inches(1.0), Inches(0.77), Inches(11), Inches(0.32), size=10, color=MUTED, italic=True)

p3 = [
    ("09  TIMING", CYAN,
     ["Number of frames = speed & personality", "MORE frames = slower, heavier movement", "FEWER frames = faster, lighter movement",
      "Eye blink: 6–8f | Head turn: 10–15f | Heavy punch: 3–4f"],
     "Wrong timing = nothing else matters",
     "Frame count chart for different actions"),
    ("10  EXAGGERATION", PURP,
     ["Push actions BEYOND realistic limits", "Giant cartoon eyes, impossible stretching", "Over-the-top emotional reactions", "NOT unrealistic — making TRUTH more visible"],
     "Exaggerate the ESSENCE, not randomness",
     "Exaggerated cartoon expressions"),
    ("11  SOLID DRAWING", GRN,
     ["Characters must feel like 3D objects", "Understand basic 3D shapes (cube, sphere, cylinder)", "Avoid 'twinning' — mirror-image poses", "Strong line of action in every pose"],
     "Even 2D must FEEL three-dimensional",
     "Line of action pose examples"),
    ("12  APPEAL", ORG,
     ["Characters must be interesting to WATCH", "NOT just cuteness — villains have appeal too!",
      "Clear, readable design & strong silhouette", "Dynamic, asymmetric poses with personality"],
     "TEST: Would audience want to keep watching?",
     "Strong character silhouettes"),
]

for i, (title, tc, pts, note, img_label) in enumerate(p3):
    c, r = i % 2, i // 2
    px = Inches(0.28 + c*6.52)
    py = Inches(1.1 + r*2.85)
    PW = Inches(6.28); PH = Inches(2.65)
    box(sl, px, py, PW, PH, fill=CARD, lc=tc, lw=1)
    grad(sl, px, py, PW, Inches(0.06), tc, CARD, ang=0)
    T(sl, title, px+Inches(0.12), py+Inches(0.12), PW*0.55, Inches(0.38), size=12, bold=True, color=tc)
    for j, pt in enumerate(pts):
        T(sl, f"• {pt}", px+Inches(0.15), py+Inches(0.55+j*0.37), PW*0.52, Inches(0.34), size=9.5, color=WHT)
    T(sl, note, px+Inches(0.15), py+PH-Inches(0.32), PW*0.52, Inches(0.28), size=9, color=MUTED, italic=True)
    img_ph(sl, px+PW*0.56, py+Inches(0.12), PW*0.41, PH-Inches(0.2), img_label)

T(sl, "All 12 Principles covered! Developed by Disney (1981) — still used by EVERY professional animator today.", Inches(0.3), Inches(7.22), Inches(12), Inches(0.25), size=9, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — ANIMATION SPACE  (Abdul Rahman)
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, 8, "Animation Space", YEL)
footer(sl, "Abdul Rahman Baig", "14 / 17")
T(sl, "Animation can exist in 3 types of spatial dimensions — each with different tools, techniques, and visual results.", Inches(1.0), Inches(0.77), Inches(11), Inches(0.32), size=10, color=MUTED, italic=True)

space_data = [
    ("2D SPACE", "X + Y axes only", CYAN,
     ["Flat — like drawing on paper", "No real depth perception", "Simple to create & render", "Low hardware requirements", "Classic cartoon aesthetic"],
     ["Tom & Jerry, classic Disney", "Original Mario & Pac-Man", "Mobile game animations", "Flash web animations"],
     ["Adobe Animate", "Toon Boom", "Spine", "TVPaint"],
     "2D flat diagram\n(X and Y axes only)"),
    ("2½D SPACE", "2D art + depth illusion", PURP,
     ["2D sprites with illusion of depth", "Parallax scrolling creates depth effect", "Technically still 2D sprites", "'Pseudo-3D' or 'Layered 2D'", "Great performance-to-visual ratio"],
     ["Paper Mario", "Ori and the Blind Forest", "South Park", "Disney multiplane camera"],
     ["After Effects (Z-depth)", "Unity 2D with parallax", "Spine Pro", ""],
     "Parallax layers diagram\n(foreground/mid/background)"),
    ("3D SPACE", "X + Y + Z axes", GRN,
     ["True depth, volume, perspective", "Viewable from ANY camera angle", "Realistic lighting & shadows", "Maximum visual versatility", "Requires powerful hardware"],
     ["Toy Story, Frozen, Moana", "Fortnite, Call of Duty", "Architectural visualizations", "Medical & surgical simulations"],
     ["Blender", "Autodesk Maya", "Houdini", "Unreal Engine"],
     "3D coordinate axes\n(X, Y, Z cube)"),
]

SPW = Inches(4.1)
for i, (title, sub, tc, chars, examples, tools, img_label) in enumerate(space_data):
    sx = Inches(0.28 + i*4.37)
    sy = Inches(1.1)
    box(sl, sx, sy, SPW, Inches(6.25), fill=CARD)
    grad(sl, sx, sy, SPW, Inches(0.07), tc, CARD, ang=0)
    T(sl, title, sx+Inches(0.12), sy+Inches(0.12), SPW-Inches(0.2), Inches(0.42), size=16, bold=True, color=tc)
    T(sl, sub, sx+Inches(0.12), sy+Inches(0.55), SPW-Inches(0.2), Inches(0.28), size=11, color=MUTED, italic=True)
    for j, c in enumerate(chars):
        T(sl, f"• {c}", sx+Inches(0.18), sy+Inches(0.88+j*0.37), SPW-Inches(0.25), Inches(0.34), size=9.5, color=WHT)
    T(sl, "Examples:", sx+Inches(0.12), sy+Inches(2.82), SPW-Inches(0.2), Inches(0.28), size=10, bold=True, color=tc)
    for j, ex in enumerate(examples[:3]):
        T(sl, f"• {ex}", sx+Inches(0.18), sy+Inches(3.1+j*0.33), SPW-Inches(0.25), Inches(0.3), size=9, color=MUTED)
    T(sl, "Tools:", sx+Inches(0.12), sy+Inches(4.1), SPW-Inches(0.2), Inches(0.28), size=10, bold=True, color=tc)
    T(sl, " · ".join(t for t in tools if t), sx+Inches(0.18), sy+Inches(4.4), SPW-Inches(0.25), Inches(0.28), size=9, color=MUTED)
    img_ph(sl, sx+Inches(0.1), sy+Inches(4.75), SPW-Inches(0.2), Inches(1.38), img_label)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — ANIMATION PROCESS  (Abdul Rahman)
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, 9, "Animation Process", GRN)
footer(sl, "Abdul Rahman Baig", "15 / 17")
T(sl, "Every animation project — from a 10-second ad to a 2-hour feature film — follows this same 4-stage pipeline.", Inches(1.0), Inches(0.77), Inches(11), Inches(0.32), size=10, color=MUTED, italic=True)

steps_data = [
    ("1", "ORGANIZE\nEXECUTION", CYAN,
     ["Define project goals & scope", "Assign team roles & deadlines", "Create style guide", "Budget & risk planning"],
     "OUTPUT: Schedule, moodboard,\nstoryboard, project bible",
     "Project planning board\nor Gantt chart"),
    ("2", "CHOOSE\nANIMATION TOOL", PURP,
     ["Match tool to project type", "Consider team skill level", "2D: Adobe Animate, Toon Boom",
      "3D: Maya, Blender | MG: After Effects"],
     "Games: Unity / Unreal Engine",
     "Software logos:\nAfter Effects, Blender, Maya"),
    ("3", "BUILD & TWEAK\nSEQUENCES", GRN,
     ["Create rough blocking first", "Refine with keyframes", "Apply the 12 Principles", "Team review cycles — repeat!"],
     "ITERATIVE STEP — repeat many\ntimes before moving to Step 4",
     "Animation timeline with\nkeyframes and tweens"),
    ("4", "POST-PROCESS\nANIMATION", ORG,
     ["Render final frames", "Color grading & VFX overlays", "Sound design & music mix", "QA testing & format for delivery"],
     "OUTPUT: Final MP4, MOV, WebM\nArchive project files",
     "Render farm or\npost-production setup"),
]

SW = Inches(3.1)
for i, (num, title, tc, pts, output, img_label) in enumerate(steps_data):
    sx = Inches(0.25 + i*3.27)
    sy = Inches(1.1)
    SH = Inches(6.25)
    box(sl, sx, sy, SW, SH, fill=CARD)
    grad(sl, sx, sy, SW, Inches(0.07), tc, CARD, ang=0)
    # Step number
    box(sl, sx+Inches(0.12), sy+Inches(0.1), Inches(0.55), Inches(0.55), fill=tc)
    T(sl, num, sx+Inches(0.12), sy+Inches(0.1), Inches(0.55), Inches(0.55), size=22, bold=True, color=BG, align=PP_ALIGN.CENTER)
    T(sl, title, sx+Inches(0.78), sy+Inches(0.1), SW-Inches(0.85), Inches(0.62), size=11, bold=True, color=tc)
    for j, pt in enumerate(pts):
        T(sl, f"• {pt}", sx+Inches(0.15), sy+Inches(0.82+j*0.4), SW-Inches(0.22), Inches(0.36), size=9.5, color=WHT)
    # Output box
    box(sl, sx+Inches(0.1), sy+Inches(2.55), SW-Inches(0.2), Inches(0.7), fill=RGBColor(8,8,35), lc=tc, lw=0.5)
    T(sl, output, sx+Inches(0.15), sy+Inches(2.6), SW-Inches(0.25), Inches(0.6), size=8.5, color=tc)
    img_ph(sl, sx+Inches(0.1), sy+Inches(3.38), SW-Inches(0.2), Inches(2.72), img_label)
    # Arrow between steps
    if i < 3:
        T(sl, "▶", sx+SW, sy+Inches(2.8), Inches(0.22), Inches(0.5), size=14, bold=True, color=tc, align=PP_ALIGN.CENTER)

# Iterative note
box(sl, Inches(0.25), Inches(7.0), Inches(12.82), Inches(0.22), fill=RGBColor(10,30,10), lc=GRN, lw=0.5)
T(sl, "KEY INSIGHT: Step 3 (Build & Tweak) is ITERATIVE — you cycle through it many times before Post-Processing.",
  Inches(0.35), Inches(7.02), Inches(12.5), Inches(0.2), size=9.5, bold=True, color=GRN)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — SUMMARY  (Abdul Rahman)
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
header(sl, 0, "Summary — Key Takeaways", CYAN)
footer(sl, "Abdul Rahman Baig", "16 / 17")

summary = [
    ("01", "Introduction",  CYAN,  "Still frames + Persistence of Vision + Phi Phenomenon = seamless motion illusion"),
    ("02", "Usage",         PURP,  "Art, storytelling, science, education — animation powers every major industry"),
    ("03", "Techniques",    GRN,   "3 core methods: Cel (hand-drawn), Path (trajectory), Computer (digital — most versatile)"),
    ("04", "Cel Animation", ORG,   "Keyframes + tweening = smooth motion. Pencil test validates timing before final production"),
    ("05", "Path Anim.",    PINK,  "Object follows preset path (straight or Bezier) with easing for natural weighted motion"),
    ("06", "Computer Anim", YEL,   "FK/IK chains for rigs. Morphing for shape transforms. Faster & more flexible than traditional"),
    ("07", "12 Principles", CYAN,  "Disney's Squash & Stretch, Anticipation, Timing, Ease, Arcs and 7 more — the universal standard"),
    ("08", "Anim Space",    PURP,  "2D (flat) → 2½D (pseudo-depth via parallax) → 3D (full volumetric, any camera angle)"),
    ("09", "Process",       GRN,   "Organize → Choose Tool → Build & Tweak (iterative!) → Post-Process — every studio uses this"),
]

SH = Inches(0.62)
for i, (num, title, tc, text) in enumerate(summary):
    r, c = i // 2, i % 2
    if i == 8: # last one full width
        sx, sy = Inches(0.28), Inches(1.1 + 4*SH + Inches(0.1)*4 + Inches(0.08))
        SW2 = Inches(12.75)
    else:
        sx = Inches(0.28 + c*6.52)
        sy = Inches(1.1 + r*(SH+Inches(0.08)))
        SW2 = Inches(6.28)
    box(sl, sx, sy, SW2, SH, fill=CARD, lc=tc, lw=0.6)
    box(sl, sx, sy, Inches(0.07), SH, fill=tc)
    box(sl, sx+Inches(0.1), sy+Inches(0.12), Inches(0.42), Inches(0.38), fill=tc)
    T(sl, num, sx+Inches(0.1), sy+Inches(0.12), Inches(0.42), Inches(0.38), size=11, bold=True, color=BG, align=PP_ALIGN.CENTER)
    T(sl, title, sx+Inches(0.6), sy+Inches(0.07), Inches(1.2), Inches(0.3), size=10, bold=True, color=tc)
    T(sl, text, sx+Inches(0.6), sy+Inches(0.35), SW2-Inches(0.65), Inches(0.24), size=9, color=WHT)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — THANK YOU
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide()
grad(sl, 0, 0, prs.slide_width, prs.slide_height, RGBColor(4,4,18), RGBColor(10,8,35), ang=135)
for i in range(14):
    box(sl, Inches(0.3+i*0.93), Inches(0.1), Inches(0.55), Inches(0.25), fill=RGBColor(15,15,50), lc=RGBColor(30,30,80))
    box(sl, Inches(0.3+i*0.93), Inches(7.15), Inches(0.55), Inches(0.25), fill=RGBColor(15,15,50), lc=RGBColor(30,30,80))

grad(sl, Inches(2.5), Inches(1.5), Inches(8.3), Inches(4.0), RGBColor(0,30,50), BG, ang=90)
grad(sl, Inches(2.0), Inches(2.35), Inches(9.3), Inches(0.07), CYAN, PURP, ang=0)
T(sl, "THANK YOU", Inches(2.0), Inches(2.45), Inches(9.3), Inches(1.2), size=60, bold=True, color=WHT, align=PP_ALIGN.CENTER)
T(sl, "Animation: Where Imagination Meets Technology", Inches(2.0), Inches(3.7), Inches(9.3), Inches(0.5), size=16, color=CYAN, align=PP_ALIGN.CENTER, italic=True)
grad(sl, Inches(2.0), Inches(4.28), Inches(9.3), Inches(0.06), CYAN, PURP, ang=0)
T(sl, "Group 3  •  Graphics & Multimedia  •  SMIT", Inches(2.0), Inches(4.42), Inches(9.3), Inches(0.35), size=12, color=MUTED, align=PP_ALIGN.CENTER)

for i, (name, roll, nc) in enumerate(names):
    nx = Inches(1.1 + i*4.1)
    box(sl, nx, Inches(5.1), Inches(3.6), Inches(1.05), fill=RGBColor(10,10,40), lc=nc, lw=1.2)
    T(sl, name, nx, Inches(5.15), Inches(3.6), Inches(0.42), size=13, bold=True, color=nc, align=PP_ALIGN.CENTER)
    T(sl, roll, nx, Inches(5.57), Inches(3.6), Inches(0.35), size=10, color=MUTED, align=PP_ALIGN.CENTER)

T(sl, "Topics: Introduction • Usage • Techniques • Types • Computer Animation • Levels • 12 Principles • Space • Process",
  Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.3), size=9, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════
output_path = "Animation_Group3_CLEAN.pptx"
prs.save(output_path)
print(f"Saved: {output_path}  ({len(prs.slides)} slides)")
