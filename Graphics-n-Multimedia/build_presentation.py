"""
Animation Presentation Builder
Group 3 - Graphics & Multimedia - SMIT
Topics 1-9
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── COLOUR PALETTE ────────────────────────────────────────────
BG       = RGBColor(8,   8,  35)   # Deep navy
CARD     = RGBColor(18, 18,  62)   # Card background
STRIPE   = RGBColor(12, 12,  48)   # Alt row stripe
CYAN     = RGBColor(0,  210, 255)
PURPLE   = RGBColor(150, 80, 220)
ORANGE   = RGBColor(255, 145,  0)
GREEN    = RGBColor(46,  213, 115)
PINK     = RGBColor(255,  90, 160)
YELLOW   = RGBColor(255, 215,   0)
WHITE    = RGBColor(255, 255, 255)
LIGHT    = RGBColor(195, 205, 235)
DARK     = RGBColor(8,    8,  35)

# ── HELPERS ───────────────────────────────────────────────────
def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide

def rect(slide, l, t, w, h, fill, line_color=None):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(1)
    else:
        sh.line.color.rgb = fill
    return sh

def txt(slide, text, l, t, w, h,
        size=14, color=WHITE, bold=False, italic=False,
        align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text        = text
    run.font.size   = Pt(size)
    run.font.color.rgb = color
    run.font.bold   = bold
    run.font.italic = italic
    return box

def header(slide, section_num, title, accent):
    """Standard slide header."""
    rect(slide, Inches(0), Inches(0), Inches(0.32), Inches(7.5), accent)
    rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.12), RGBColor(5,5,40))
    txt(slide, section_num,
        Inches(0.42), Inches(0.08), Inches(1.1), Inches(1.0),
        size=38, color=accent, bold=True)
    txt(slide, title,
        Inches(1.6), Inches(0.18), Inches(11.3), Inches(0.9),
        size=28, color=WHITE, bold=True)

def darken(color, factor=0.35):
    return RGBColor(int(color[0]*factor), int(color[1]*factor), int(color[2]*factor))

def info_bar(slide, text, color=None, y=None):
    if color is None:
        color = CYAN
    if y is None:
        y = Inches(1.18)
    rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.8), darken(color))
    txt(slide, text, Inches(0.65), y+Inches(0.07), Inches(12.0), Inches(0.72),
        size=12, color=WHITE, italic=True)

def card(slide, x, y, w, h, accent, title, body, title_size=13, body_size=10.5):
    rect(slide, x, y, w, h, CARD)
    rect(slide, x, y, w, Inches(0.62), accent)
    txt(slide, title, x+Inches(0.06), y+Inches(0.04), w-Inches(0.12), Inches(0.58),
        size=title_size, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, body, x+Inches(0.12), y+Inches(0.72), w-Inches(0.24), h-Inches(0.82),
        size=body_size, color=LIGHT)

# ════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ╔══════════════════════════════════════════╗
# ║  SLIDE 1 — TITLE                         ║
# ╚══════════════════════════════════════════╝
s1 = new_slide(prs)
rect(s1, Inches(0),     Inches(0),    Inches(13.33), Inches(0.18), CYAN)
rect(s1, Inches(0),     Inches(7.32), Inches(13.33), Inches(0.18), PURPLE)
rect(s1, Inches(0),     Inches(0),    Inches(13.33), Inches(7.5),
     RGBColor(8,8,35))    # re-apply so stripes render on top
rect(s1, Inches(0),     Inches(0),    Inches(13.33), Inches(0.18), CYAN)
rect(s1, Inches(0),     Inches(7.32), Inches(13.33), Inches(0.18), PURPLE)

# Glow card
rect(s1, Inches(1.5), Inches(0.5), Inches(10.33), Inches(6.5), RGBColor(13,13,52))

txt(s1, "ANIMATION",
    Inches(1), Inches(0.7), Inches(11.33), Inches(2.2),
    size=90, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

txt(s1, "The Art of Bringing Life to Motion",
    Inches(1), Inches(2.9), Inches(11.33), Inches(0.8),
    size=24, color=LIGHT, italic=True, align=PP_ALIGN.CENTER)

rect(s1, Inches(3.8), Inches(3.85), Inches(5.73), Inches(0.06), ORANGE)

txt(s1, "Graphics & Multimedia  |  SMIT",
    Inches(1), Inches(4.0), Inches(11.33), Inches(0.55),
    size=16, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

rect(s1, Inches(3.5), Inches(4.65), Inches(6.33), Inches(0.9), RGBColor(0,140,190))
txt(s1, "Topics 1 — 9  |  Group 3",
    Inches(3.5), Inches(4.65), Inches(6.33), Inches(0.9),
    size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(s1,
    "Hashir Junaid (2023F-BCS-358)   •   Taha Haider (2023F-BCS-079)   •   Abdul Rahman Baig (2023F-BCS-100)",
    Inches(0.5), Inches(5.75), Inches(12.33), Inches(0.6),
    size=12, color=LIGHT, align=PP_ALIGN.CENTER)

# ╔══════════════════════════════════════════╗
# ║  SLIDE 2 — AGENDA                        ║
# ╚══════════════════════════════════════════╝
s2 = new_slide(prs)
rect(s2, Inches(0), Inches(0), Inches(13.33), Inches(1.2), RGBColor(0,140,190))
txt(s2, "WHAT WE'LL COVER",
    Inches(0.4), Inches(0.1), Inches(12), Inches(1.05),
    size=34, color=WHITE, bold=True)

agenda = [
    ("01", "Introduction to Animation",    CYAN),
    ("02", "Usage of Animation",           PURPLE),
    ("03", "Animation Techniques",         ORANGE),
    ("04", "Types of Animation",           GREEN),
    ("05", "Computer Animation",           PINK),
    ("06", "Levels of Computer Animation", YELLOW),
    ("07", "12 Basic Principles",          CYAN),
    ("08", "Animation Space",              PURPLE),
    ("09", "Animation Process",            ORANGE),
]

col1, col2 = agenda[:5], agenda[5:]

for i, (n, t, c) in enumerate(col1):
    y = Inches(1.35) + i * Inches(1.18)
    rect(s2, Inches(0.4), y, Inches(0.7), Inches(0.9), c)
    txt(s2, n, Inches(0.4), y, Inches(0.7), Inches(0.9),
        size=19, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    rect(s2, Inches(1.15), y, Inches(5.2), Inches(0.9), CARD)
    txt(s2, t, Inches(1.25), y+Inches(0.12), Inches(5.0), Inches(0.75),
        size=16, color=WHITE)

for i, (n, t, c) in enumerate(col2):
    y = Inches(1.35) + i * Inches(1.18)
    rect(s2, Inches(7.0), y, Inches(0.7), Inches(0.9), c)
    txt(s2, n, Inches(7.0), y, Inches(0.7), Inches(0.9),
        size=19, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    rect(s2, Inches(7.75), y, Inches(5.2), Inches(0.9), CARD)
    txt(s2, t, Inches(7.85), y+Inches(0.12), Inches(5.0), Inches(0.75),
        size=16, color=WHITE)

# ╔══════════════════════════════════════════╗
# ║  SLIDE 3 — INTRODUCTION (4 concepts)     ║
# ╚══════════════════════════════════════════╝
s3 = new_slide(prs)
header(s3, "01", "INTRODUCTION TO ANIMATION", CYAN)

rect(s3, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.05), RGBColor(0,90,130))
txt(s3,
    '"Animation is the technique of sequentially displaying STATIC images rapidly enough to\n'
    ' create the ILLUSION OF MOVEMENT — typically at 24 or more frames per second."',
    Inches(0.65), Inches(1.28), Inches(12.0), Inches(1.0),
    size=14, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

concepts = [
    (CYAN,   "DEFINITION",
     "• From Latin 'anima' — to give life\n"
     "• Series of still images shown rapidly\n"
     "• Standard: 24 fps (film), 30 fps (TV)\n"
     "• Below 12 fps → choppy, flickery\n"
     "• First animation: Eadweard Muybridge\n  (1878) horse gallop sequence"),
    (ORANGE, "PERSISTENCE\nOF VISION",
     "• Retina retains image ~50ms after\n  stimulus disappears\n"
     "• Afterimage overlaps the next frame\n"
     "• Creates a continuous visual flow\n"
     "• Discovered: Peter Mark Roget (1824)\n"
     "• WHY films don't flicker at 24fps"),
    (GREEN,  "PHI PHENOMENON",
     "• PSYCHOLOGICAL effect (not retinal)\n"
     "• Brain infers motion between two\n  stationary objects shown in sequence\n"
     "• Discovered: Max Wertheimer (1912)\n"
     "• Classic example: flip-book\n"
     "• Traffic light arrows 'moving'"),
    (PURPLE, "ILLUSION OF\nMOVEMENT",
     "• Result of POV + Phi Phenomenon\n"
     "• Neither alone is enough\n"
     "• Combined = seamless motion\n"
     "• The brain WANTS to see motion\n"
     "• Animators exploit this 'glitch'\n  in human perception"),
]

for i, (c, t, b) in enumerate(concepts):
    x = Inches(0.45) + i * Inches(3.22)
    card(s3, x, Inches(2.35), Inches(3.0), Inches(4.85), c, t, b, title_size=13)

# ╔══════════════════════════════════════════╗
# ║  SLIDE 4 — USAGE OF ANIMATION            ║
# ╚══════════════════════════════════════════╝
s4 = new_slide(prs)
header(s4, "02", "USAGE OF ANIMATION", PURPLE)

usages = [
    (CYAN,   "ARTISTIC PURPOSES",
     "• Fine art & experimental film\n"
     "• Abstract visual compositions\n"
     "• Music video production\n"
     "• Title sequence design\n\n"
     "EXAMPLES:\n"
     "  Disney's Fantasia (1940)\n"
     "  Yellow Submarine (1968)\n"
     "  Spider-Man: Into the\n  Spider-Verse (2018)"),
    (ORANGE, "STORYTELLING",
     "• Narrative films & TV series\n"
     "• Video game cutscenes\n"
     "• Interactive web comics\n"
     "• Character-driven dramas\n\n"
     "EXAMPLES:\n"
     "  Pixar's Toy Story series\n"
     "  Studio Ghibli classics\n"
     "  Arcane (Netflix, 2021)\n"
     "  Attack on Titan"),
    (GREEN,  "SCIENTIFIC\nVISUALIZATION",
     "• Medical anatomy animations\n"
     "• Weather system simulations\n"
     "• Physics & fluid dynamics\n"
     "• DNA/molecular modeling\n\n"
     "EXAMPLES:\n"
     "  NASA space animations\n"
     "  DNA replication videos\n"
     "  Climate change models\n"
     "  Surgical training sims"),
    (PINK,   "INSTRUCTIONAL\nPURPOSES",
     "• Educational explainer videos\n"
     "• Flight & military simulators\n"
     "• Safety training demos\n"
     "• Product assembly guides\n\n"
     "EXAMPLES:\n"
     "  Khan Academy lessons\n"
     "  Pilot training software\n"
     "  IKEA assembly guides\n"
     "  YouTube tutorials"),
]

for i, (c, t, b) in enumerate(usages):
    x = Inches(0.4) + i * Inches(3.22)
    rect(s4, x, Inches(1.22), Inches(3.0), Inches(6.05), CARD)
    rect(s4, x, Inches(1.22), Inches(3.0), Inches(0.18), c)
    txt(s4, t, x, Inches(1.45), Inches(3.0), Inches(0.8),
        size=13, color=c, bold=True, align=PP_ALIGN.CENTER)
    txt(s4, b, x+Inches(0.12), Inches(2.35), Inches(2.76), Inches(4.8),
        size=11.5, color=LIGHT)

# ╔══════════════════════════════════════════╗
# ║  SLIDE 5 — ANIMATION TECHNIQUES (3-way) ║
# ╚══════════════════════════════════════════╝
s5 = new_slide(prs)
header(s5, "03", "ANIMATION TECHNIQUES", ORANGE)

techniques = [
    (CYAN,   "CEL ANIMATION",
     "Traditional hand-drawn technique\n\n"
     "• Drawings on transparent celluloid sheets\n"
     "• Cels layered over static backgrounds\n"
     "• Each cel = one individual frame\n"
     "• Photographed frame-by-frame\n\n"
     "WORKFLOW:\n"
     "  Story → Storyboard → Layout →\n"
     "  Key animation → In-betweening →\n"
     "  Inking → Painting → Photography\n\n"
     "ERA: 1910s – 1990s\n"
     "STUDIOS: Disney, Warner Bros,\n"
     "          Hanna-Barbera"),
    (ORANGE, "PATH ANIMATION",
     "Object follows a predefined trajectory\n\n"
     "• Animator defines the path shape\n"
     "• Object moves along it automatically\n"
     "• Supports straight lines & Bezier curves\n"
     "• Object can resize/reshape along path\n\n"
     "USE CASES:\n"
     "  Logo reveal animations\n"
     "  Camera fly-through paths\n"
     "  UI element transitions\n"
     "  Game projectile trajectories\n\n"
     "TOOLS: After Effects motion paths,\n"
     "        CSS animations, SVG paths"),
    (GREEN,  "COMPUTER ANIMATION",
     "Digital creation of moving images\n\n"
     "• Can simulate cel & path techniques\n"
     "• Plus physics, particles, simulation\n"
     "• 2D digital, 3D CG, or hybrid\n"
     "• Real-time or pre-rendered\n\n"
     "ADVANTAGES:\n"
     "  Unlimited camera angles\n"
     "  Easy to iterate & modify\n"
     "  Physics simulation\n"
     "  Procedural animation\n\n"
     "ERA: 1980s – present\n"
     "STUDIOS: Pixar, DreamWorks, ILM"),
]

for i, (c, t, b) in enumerate(techniques):
    x = Inches(0.55) + i * Inches(4.26)
    rect(s5, x, Inches(1.22), Inches(4.0), Inches(6.0), CARD)
    rect(s5, x, Inches(1.22), Inches(4.0), Inches(0.65), c)
    rect(s5, x, Inches(1.22), Inches(0.55), Inches(0.65), DARK)
    txt(s5, str(i+1), x, Inches(1.22), Inches(0.55), Inches(0.65),
        size=20, color=c, bold=True, align=PP_ALIGN.CENTER)
    txt(s5, t, x+Inches(0.6), Inches(1.28), Inches(3.3), Inches(0.58),
        size=14, color=DARK, bold=True)
    txt(s5, b, x+Inches(0.14), Inches(1.95), Inches(3.72), Inches(5.1),
        size=11, color=LIGHT)
    if i < 2:
        txt(s5, "▶", Inches(4.4) + i*Inches(4.26), Inches(4.1),
            Inches(0.5), Inches(0.8), size=20, color=ORANGE, bold=True)

# ╔══════════════════════════════════════════╗
# ║  SLIDE 6 — CEL ANIMATION (4 concepts)   ║
# ╚══════════════════════════════════════════╝
s6 = new_slide(prs)
header(s6, "04", "TYPES OF ANIMATION — CEL ANIMATION", CYAN)
info_bar(s6,
    "Named after CELLULOID — the transparent plastic sheets used to draw frames."
    "  Developed ~1914 by Raoul Barré & John Bray. Dominated animation for 80 years.")

cel_cards = [
    (CYAN,   "CEL ANIMATION",
     "• Transparent sheet per frame\n"
     "• Cels stacked: character over BG\n"
     "• Allows background reuse\n"
     "• Reduced drawing workload\n"
     "• Colors applied to back of cel\n"
     "• ~12-24 cels per second of film\n\n"
     "Used in: Tom & Jerry,\n"
     "Looney Tunes, Snow White"),
    (ORANGE, "KEYFRAMES",
     "• The 'anchor' frames of action\n"
     "• Show extreme or critical poses\n"
     "• Drawn by senior animators\n"
     "• Define timing & structure\n"
     "• Example: frame 1 (standing)\n"
     "  and frame 12 (peak of jump)\n\n"
     "Modern: keyframes in software\n"
     "timelines work identically"),
    (GREEN,  "TWEENING\n(In-Betweening)",
     "• Frames drawn BETWEEN keyframes\n"
     "• Creates smooth transition\n"
     "• Done by junior animators\n"
     "• More tweens = smoother motion\n"
     "• At 24fps: ~22 tweens per action\n\n"
     "Modern software AUTO-tweens:\n"
     "  Adobe Animate, Flash,\n"
     "  Toon Boom Harmony"),
    (PURPLE, "PENCIL TEST",
     "• Quick pre-production test\n"
     "• Rough pencil sketches filmed\n"
     "• Played back to check timing\n"
     "• Catches errors BEFORE costly\n"
     "  inking & painting stages\n"
     "• Saves significant time & cost\n\n"
     "Modern equivalent:\n"
     "  Animatic / rough cut preview"),
]

for i, (c, t, b) in enumerate(cel_cards):
    x = Inches(0.45) + i * Inches(3.22)
    card(s6, x, Inches(2.1), Inches(3.0), Inches(5.1), c, t, b)

# ╔══════════════════════════════════════════╗
# ║  SLIDE 7 — PATH ANIMATION                ║
# ╚══════════════════════════════════════════╝
s7 = new_slide(prs)
header(s7, "04", "TYPES OF ANIMATION — PATH ANIMATION", ORANGE)
info_bar(s7,
    "PATH ANIMATION: An object moves along a predefined trajectory through space."
    "  The animator draws the path; software handles movement along it.")

# Left: concepts panel
rect(s7, Inches(0.45), Inches(2.1), Inches(6.0), Inches(5.15), CARD)
txt(s7, "KEY CONCEPTS", Inches(0.55), Inches(2.18), Inches(5.8), Inches(0.5),
    size=14, color=ORANGE, bold=True)

path_pts = [
    (ORANGE, "PREDETERMINED PATH",
     "Route through space defined BEFORE animation starts.\nLike a roller-coaster track — the object follows exactly."),
    (CYAN,   "STRAIGHT LINES",
     "Linear movement between two points. Used for mechanical objects, UI reveals, and text animations."),
    (GREEN,  "CURVES (Bezier)",
     "Organic, natural paths defined by control points. Used for bouncing balls, birds, camera movements."),
    (PINK,   "RESIZING & RESHAPING",
     "Object scales or deforms as it follows the path — adds dynamism (e.g. ball squashing on landing)."),
]

y_off = Inches(2.75)
for c, title, desc in path_pts:
    rect(s7, Inches(0.55), y_off, Inches(0.1), Inches(0.85), c)
    txt(s7, title, Inches(0.75), y_off, Inches(5.3), Inches(0.4),
        size=12, color=c, bold=True)
    txt(s7, desc, Inches(0.75), y_off+Inches(0.38), Inches(5.5), Inches(0.5),
        size=10.5, color=LIGHT)
    y_off += Inches(1.1)

# Right: workflow steps panel
rect(s7, Inches(6.85), Inches(2.1), Inches(6.0), Inches(5.15), CARD)
txt(s7, "HOW PATH ANIMATION WORKS", Inches(6.95), Inches(2.18), Inches(5.8), Inches(0.5),
    size=14, color=ORANGE, bold=True)

steps = [
    (CYAN,   "STEP 1", "Define the PATH shape (line, curve, or custom)"),
    (ORANGE, "STEP 2", "Set START and END points along the path"),
    (GREEN,  "STEP 3", "Define TIMING — duration of the movement"),
    (PURPLE, "STEP 4", "Apply EASING — speed variation along the path"),
    (PINK,   "STEP 5", "Add OBJECT ORIENTATION — does it face the direction of travel?"),
    (YELLOW, "STEP 6", "Preview, adjust, and render final animation"),
]

y_off = Inches(2.78)
for c, step, desc in steps:
    rect(s7, Inches(6.95), y_off, Inches(1.0), Inches(0.65), c)
    txt(s7, step, Inches(6.95), y_off, Inches(1.0), Inches(0.65),
        size=11, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    txt(s7, desc, Inches(8.05), y_off+Inches(0.08), Inches(4.7), Inches(0.6),
        size=11, color=LIGHT)
    y_off += Inches(0.78)

rect(s7, Inches(0.45), Inches(7.05), Inches(12.4), Inches(0.38),
     RGBColor(0, 80, 0))
txt(s7, "Real-World Examples:  CSS @keyframes  •  After Effects motion paths  •  Unity NavMesh agent paths  •  Sprite trajectories in games",
    Inches(0.6), Inches(7.06), Inches(12.1), Inches(0.36),
    size=11, color=GREEN)

# ╔══════════════════════════════════════════╗
# ║  SLIDE 8 — COMPUTER ANIMATION INTRO     ║
# ╚══════════════════════════════════════════╝
s8 = new_slide(prs)
header(s8, "05", "COMPUTER ANIMATION", GREEN)
info_bar(s8,
    "Computer animation uses digital software to create moving images — simulating traditional"
    " techniques or producing entirely new effects impossible by hand.", GREEN)

# VS comparison
txt(s8, "CEL ANIMATION  vs  COMPUTER ANIMATION",
    Inches(0.5), Inches(2.1), Inches(12.3), Inches(0.55),
    size=17, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)

rect(s8, Inches(0.5), Inches(2.72), Inches(5.5), Inches(0.52), CYAN)
txt(s8, "CEL ANIMATION", Inches(0.5), Inches(2.72), Inches(5.5), Inches(0.52),
    size=14, color=DARK, bold=True, align=PP_ALIGN.CENTER)

rect(s8, Inches(7.3), Inches(2.72), Inches(5.5), Inches(0.52), GREEN)
txt(s8, "COMPUTER ANIMATION", Inches(7.3), Inches(2.72), Inches(5.5), Inches(0.52),
    size=14, color=DARK, bold=True, align=PP_ALIGN.CENTER)

rect(s8, Inches(6.05), Inches(2.6), Inches(1.2), Inches(1.05), ORANGE)
txt(s8, "VS", Inches(6.05), Inches(2.6), Inches(1.2), Inches(1.05),
    size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

rows = [
    ("Hand-drawn on physical celluloid sheets", "Created entirely with digital software"),
    ("Highly time-consuming — every frame by hand", "Automation handles tweens, physics, rendering"),
    ("Physical storage: film cans, shelving", "Digital files — easy to backup and share"),
    ("Expensive materials: cels, paint, camera", "One-time software cost; no materials"),
    ("Difficult to reuse or modify frames", "Copy, paste, modify assets instantly"),
    ("Organic, hand-crafted visual warmth", "Precise, photorealistic, or any style"),
    ("Production time: months per minute", "Faster pipeline; real-time previews"),
]

for i, (cel, comp) in enumerate(rows):
    y = Inches(3.3) + i * Inches(0.56)
    bg = CARD if i % 2 == 0 else STRIPE
    rect(s8, Inches(0.5), y, Inches(5.5), Inches(0.52), bg)
    rect(s8, Inches(7.3), y, Inches(5.5), Inches(0.52), bg)
    txt(s8, cel,  Inches(0.62), y+Inches(0.06), Inches(5.25), Inches(0.44),
        size=11, color=LIGHT)
    txt(s8, comp, Inches(7.42), y+Inches(0.06), Inches(5.25), Inches(0.44),
        size=11, color=GREEN)

# ╔══════════════════════════════════════════╗
# ║  SLIDE 9 — IK, FK & MORPHING            ║
# ╚══════════════════════════════════════════╝
s9 = new_slide(prs)
header(s9, "05", "KINEMATICS, INVERSE KINEMATICS & MORPHING", GREEN)

panels = [
    (CYAN,   "FORWARD KINEMATICS (FK)",
     "Animating from PARENT to CHILD joints.\n\n"
     "HOW IT WORKS:\n"
     "  Move the SHOULDER first\n"
     "  → Elbow follows shoulder\n"
     "  → Wrist follows elbow\n"
     "  → Finger follows wrist\n\n"
     "ANALOGY:\n"
     "Like a puppet — you control\n"
     "the parent; limbs follow.\n\n"
     "USED FOR:\n"
     "  Mechanical / robotic arms\n"
     "  Waving motions\n"
     "  Tail or tentacle animation\n"
     "  Flag ripple effects"),
    (ORANGE, "INVERSE KINEMATICS (IK)",
     "Reverse of FK — move the END\nPOINT; chain adjusts automatically.\n\n"
     "HOW IT WORKS:\n"
     "  Move the HAND position\n"
     "  → Software calculates elbow\n"
     "  → And shoulder positions\n"
     "  → Works backwards up chain\n\n"
     "ANALOGY:\n"
     "Reach for a cup — you aim\n"
     "your hand; body figures\n"
     "out the rest automatically.\n\n"
     "USED FOR:\n"
     "  Walking & running cycles\n"
     "  Character reaching/grabbing\n"
     "  Real-time game characters"),
    (PURPLE, "MORPHING",
     "Smooth transformation of one\nshape/image into a DIFFERENT one.\n\n"
     "HOW IT WORKS:\n"
     "  Map points on source\n"
     "  Map corresponding target points\n"
     "  Software interpolates all\n"
     "  points simultaneously\n\n"
     "TYPES:\n"
     "  2D Morphing — image to image\n"
     "  3D Morphing — mesh deformation\n"
     "  Shape morphing — vector paths\n\n"
     "FAMOUS EXAMPLES:\n"
     "  MJ 'Black or White' video\n"
     "  T-1000 in Terminator 2\n"
     "  Werewolf transformations\n"
     "  Logo reveal animations"),
]

for i, (c, t, b) in enumerate(panels):
    x = Inches(0.45) + i * Inches(4.3)
    rect(s9, x, Inches(1.22), Inches(4.0), Inches(6.05), CARD)
    rect(s9, x, Inches(1.22), Inches(4.0), Inches(0.62), c)
    txt(s9, t, x+Inches(0.08), Inches(1.27), Inches(3.85), Inches(0.56),
        size=13, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    txt(s9, b, x+Inches(0.14), Inches(1.92), Inches(3.72), Inches(5.2),
        size=11, color=LIGHT)

# ╔══════════════════════════════════════════╗
# ║  SLIDE 10 — LEVELS OF COMPUTER ANIM    ║
# ╚══════════════════════════════════════════╝
s10 = new_slide(prs)
header(s10, "06", "LEVELS OF COMPUTER ANIMATION", YELLOW)

levels = [
    (CYAN,   "BASIC", "★☆☆",
     "Entry-level animation using simple tools\n\n"
     "WHAT YOU CREATE:\n"
     "  • Animated GIFs\n"
     "  • CSS / PowerPoint transitions\n"
     "  • Simple 2D sprite movement\n"
     "  • Basic logo animations\n"
     "  • Slideshow animations\n\n"
     "TOOLS:\n"
     "  MS PowerPoint, Canva,\n"
     "  GIMP, basic CSS\n\n"
     "SKILL LEVEL: Beginner\n"
     "TIME TO LEARN: Days – weeks"),
    (ORANGE, "INTERMEDIATE", "★★☆",
     "Deeper control over animation systems\n\n"
     "WHAT YOU CREATE:\n"
     "  • 2D character animation\n"
     "  • Motion graphics & infographics\n"
     "  • Complex path animations\n"
     "  • Basic 3D object animation\n"
     "  • Game character movement\n\n"
     "TOOLS:\n"
     "  Adobe After Effects,\n"
     "  Adobe Animate, Blender (basics),\n"
     "  Spine, Toon Boom\n\n"
     "SKILL LEVEL: Intermediate\n"
     "TIME TO LEARN: Months"),
    (PURPLE, "ADVANCED", "★★★",
     "Production-quality, studio-level work\n\n"
     "WHAT YOU CREATE:\n"
     "  • Full CG feature films\n"
     "  • Realistic physics simulations\n"
     "  • Procedural & generative anim\n"
     "  • Motion capture integration\n"
     "  • Real-time VFX in games\n"
     "  • AI-assisted animation\n\n"
     "TOOLS:\n"
     "  Autodesk Maya, Houdini,\n"
     "  Unreal Engine, Cinema 4D\n\n"
     "SKILL LEVEL: Professional\n"
     "TIME TO LEARN: Years"),
]

for i, (c, t, stars, b) in enumerate(levels):
    x = Inches(0.5) + i * Inches(4.28)
    rect(s10, x, Inches(1.22), Inches(4.0), Inches(6.05), CARD)
    rect(s10, x, Inches(1.22), Inches(4.0), Inches(0.95), c)
    txt(s10, stars, x, Inches(1.24), Inches(4.0), Inches(0.5),
        size=20, color=DARK, align=PP_ALIGN.CENTER)
    txt(s10, t, x, Inches(1.68), Inches(4.0), Inches(0.52),
        size=16, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    txt(s10, b, x+Inches(0.14), Inches(2.25), Inches(3.72), Inches(4.9),
        size=11, color=LIGHT)

# ╔══════════════════════════════════════════╗
# ║  SLIDE 11 — 12 PRINCIPLES (1 of 3)      ║
# ╚══════════════════════════════════════════╝
s11 = new_slide(prs)
header(s11, "07", "12 BASIC PRINCIPLES OF ANIMATION  —  Part 1 of 3", CYAN)
info_bar(s11,
    "Developed by Disney animators Ollie Johnston & Frank Thomas (book: The Illusion of Life, 1981)."
    "  Still the GOLD STANDARD for all animation worldwide.")

p1 = [
    (CYAN,   "01  SQUASH & STRETCH",
     "Exaggerate deformation for life & weight.\n\n"
     "• Ball SQUASHES on impact\n"
     "• Ball STRETCHES when in fast motion\n"
     "• Total volume stays CONSTANT\n"
     "• Shows elasticity & mass\n\n"
     "EXAMPLE:\n"
     "  Rubber ball bounce —\n"
     "  squash on landing,\n"
     "  stretch in mid-air"),
    (ORANGE, "02  ANTICIPATION",
     "Small opposite action BEFORE main action.\n\n"
     "• Character crouches before jumping\n"
     "• Arm swings back before throwing\n"
     "• Pitcher winds up before pitching\n"
     "• Prepares audience mentally\n\n"
     "EXAMPLE:\n"
     "  Bugs Bunny ducking\n"
     "  down before a huge\n"
     "  comedic leap"),
    (GREEN,  "03  STAGING",
     "Clear presentation of one idea at a time.\n\n"
     "• Use clear silhouettes\n"
     "• Camera angle tells the story\n"
     "• Avoid busy, distracting BGs\n"
     "• Direct the audience's eye\n\n"
     "EXAMPLE:\n"
     "  Character in profile view\n"
     "  making dramatic gesture —\n"
     "  instantly readable"),
    (PURPLE, "04  STRAIGHT AHEAD\n& POSE-TO-POSE",
     "Two animation approaches:\n\n"
     "STRAIGHT AHEAD:\n"
     "  Draw frame by frame start\n"
     "  to finish — spontaneous,\n"
     "  fluid, unpredictable\n\n"
     "POSE-TO-POSE:\n"
     "  Draw key poses first, fill\n"
     "  in betweens — controlled,\n"
     "  planned, consistent\n\n"
     "Best work combines BOTH"),
]

for i, (c, t, b) in enumerate(p1):
    x = Inches(0.45) + i * Inches(3.22)
    card(s11, x, Inches(2.05), Inches(3.0), Inches(5.2), c, t, b)

# ╔══════════════════════════════════════════╗
# ║  SLIDE 12 — 12 PRINCIPLES (2 of 3)      ║
# ╚══════════════════════════════════════════╝
s12 = new_slide(prs)
header(s12, "07", "12 BASIC PRINCIPLES OF ANIMATION  —  Part 2 of 3", ORANGE)
info_bar(s12,
    "These principles are NOT rules — they are tools. Master them, then bend them deliberately for artistic effect.",
    ORANGE)

p2 = [
    (CYAN,   "05  FOLLOW THROUGH\n& OVERLAPPING ACTION",
     "FOLLOW THROUGH:\n"
     "  Parts keep moving after\n"
     "  main action stops.\n"
     "  • Hair swings after head stops\n"
     "  • Coat settles after character\n\n"
     "OVERLAPPING:\n"
     "  Different parts move at\n"
     "  different rates/delays.\n"
     "  • Arms lag behind torso\n"
     "  • Makes motion feel organic"),
    (ORANGE, "06  EASE IN & EASE OUT\n(Slow In / Slow Out)",
     "Objects ACCELERATE at start and\nDECELERATE at end of movement.\n\n"
     "• Nothing in nature starts/stops\n"
     "  at full speed instantly\n"
     "• A car: slow start → full speed\n"
     "  → gradual brake\n"
     "• Gives weight & natural feel\n\n"
     "In software:\n"
     "  Bézier easing curves\n"
     "  Linear = robotic/artificial\n"
     "  Ease = natural/believable"),
    (GREEN,  "07  ARCS",
     "Natural movement follows ARC-\nshape paths, not straight lines.\n\n"
     "• Arms swing in arcs\n"
     "• Heads nod in arcs\n"
     "• Thrown objects: parabola\n"
     "• Even eyes dart in arcs\n\n"
     "EXCEPTION:\n"
     "  Mechanical/robot movement\n"
     "  IS straight — use straight\n"
     "  lines intentionally for\n"
     "  mechanical characters"),
    (PURPLE, "08  SECONDARY ACTION",
     "Additional actions that SUPPORT\nand enrich the main action.\n\n"
     "EXAMPLES:\n"
     "  Walking (main)\n"
     "  + Arms swinging (secondary)\n\n"
     "  Running (main)\n"
     "  + Heavy breathing (secondary)\n\n"
     "RULE:\n"
     "  Secondary should NEVER\n"
     "  overpower the main action\n"
     "  — it supports, not steals"),
]

for i, (c, t, b) in enumerate(p2):
    x = Inches(0.45) + i * Inches(3.22)
    card(s12, x, Inches(2.05), Inches(3.0), Inches(5.2), c, t, b)

# ╔══════════════════════════════════════════╗
# ║  SLIDE 13 — 12 PRINCIPLES (3 of 3)      ║
# ╚══════════════════════════════════════════╝
s13 = new_slide(prs)
header(s13, "07", "12 BASIC PRINCIPLES OF ANIMATION  —  Part 3 of 3", PURPLE)
info_bar(s13,
    "Walt Disney: 'The secret of animation is the exaggeration of the truth.' — These final 4 principles deliver that truth.",
    PURPLE)

p3 = [
    (CYAN,   "09  TIMING",
     "Number of frames = perceived\nspeed, weight, and feel.\n\n"
     "MORE frames = slower, heavier\n"
     "FEWER frames = faster, lighter\n\n"
     "EXAMPLES AT 24fps:\n"
     "  Eye blink: 6-8 frames\n"
     "  Head turn: 10-15 frames\n"
     "  Heavy punch: 4 frames\n"
     "  Gentle wave: 30+ frames\n\n"
     "TIMING is the HEARTBEAT\n"
     "of animation — get it wrong\n"
     "and nothing else matters"),
    (ORANGE, "10  EXAGGERATION",
     "Push actions beyond realistic\nlimits for emotional impact.\n\n"
     "• Giant cartoon eyes\n"
     "• Impossible stretching\n"
     "• Huge emotional expressions\n"
     "• Over-the-top reactions\n\n"
     "NOT about being unrealistic —\nit's making the TRUTH more\nVISIBLE and ENTERTAINING.\n\n"
     "RULE:\n"
     "  Exaggerate the ESSENCE,\n"
     "  not random features.\n"
     "  A sad character cries\n"
     "  buckets — not leaks"),
    (GREEN,  "11  SOLID DRAWING",
     "Characters must appear as 3D\nobjects with weight and form.\n\n"
     "PRINCIPLES:\n"
     "  • Understand 3D basic shapes\n"
     "  • 'Draw through' forms\n"
     "  • Consistent proportions\n"
     "  • Avoid 'twinning' (mirror poses)\n"
     "  • Strong line of action\n\n"
     "Even in 2D, the audience\nshould FEEL the 3D volume.\n\n"
     "Modern: 3D rigs must deform\nbelievably under movement"),
    (PURPLE, "12  APPEAL",
     "Characters must be visually\ninteresting to WATCH.\n\n"
     "APPEAL ≠ CUTENESS:\n"
     "  • A powerful villain has appeal\n"
     "  • A complex anti-hero has it\n"
     "  • Even monsters can have it\n\n"
     "HOW TO ACHIEVE:\n"
     "  • Clear, readable design\n"
     "  • Charisma & personality\n"
     "  • Avoid stiff, symmetric poses\n"
     "  • Dynamic silhouette\n\n"
     "TEST: Would an audience\nwant to KEEP watching\nthis character?"),
]

for i, (c, t, b) in enumerate(p3):
    x = Inches(0.45) + i * Inches(3.22)
    card(s13, x, Inches(2.05), Inches(3.0), Inches(5.2), c, t, b)

# ╔══════════════════════════════════════════╗
# ║  SLIDE 14 — ANIMATION SPACE             ║
# ╚══════════════════════════════════════════╝
s14 = new_slide(prs)
header(s14, "08", "ANIMATION SPACE", PINK)

spaces = [
    (CYAN,   "2D SPACE",
     "Objects exist in TWO dimensions\n(X-axis width + Y-axis height)\n\n"
     "CHARACTERISTICS:\n"
     "  • Flat, like drawing on paper\n"
     "  • No real depth perception\n"
     "  • Simple to create & render\n"
     "  • Low hardware requirements\n"
     "  • Classic cartoon aesthetic\n"
     "  • Fast production pipeline\n\n"
     "EXAMPLES:\n"
     "  Tom & Jerry\n"
     "  Classic Disney (pre-1990)\n"
     "  Original Mario & Pac-Man\n"
     "  Most mobile game animations\n"
     "  Flash web animations\n\n"
     "TOOLS:\n"
     "  Adobe Animate, Toon Boom,\n"
     "  Spine, TVPaint"),
    (ORANGE, "2½D SPACE\n(Two-and-a-Half D)",
     "2D art creating the ILLUSION\nof depth via layers & parallax.\n\n"
     "CHARACTERISTICS:\n"
     "  • Parallax scrolling = depth\n"
     "  • Objects appear on Z-axis\n"
     "    but are still 2D sprites\n"
     "  • 'Pseudo-3D' / 'Layered 2D'\n"
     "  • Richer than 2D, simpler than 3D\n"
     "  • Great performance/visual trade-off\n\n"
     "EXAMPLES:\n"
     "  Paper Mario games\n"
     "  South Park animation\n"
     "  Ori and the Blind Forest\n"
     "  Disney multiplane camera\n"
     "  (Sleeping Beauty, 1959)\n\n"
     "TOOLS:\n"
     "  After Effects (Z-depth layers)"),
    (GREEN,  "3D SPACE",
     "Full THREE dimensions:\nX-axis + Y-axis + Z-axis (depth)\n\n"
     "CHARACTERISTICS:\n"
     "  • True depth, volume, perspective\n"
     "  • Viewable from ANY angle\n"
     "  • Realistic lighting & shadows\n"
     "  • Requires powerful hardware\n"
     "  • Longer render times\n"
     "  • Maximum visual versatility\n\n"
     "EXAMPLES:\n"
     "  Toy Story, Frozen, Moana\n"
     "  Fortnite, Call of Duty\n"
     "  Architectural visualization\n"
     "  Medical simulations\n"
     "  VFX in live-action films\n\n"
     "TOOLS:\n"
     "  Blender, Autodesk Maya,\n"
     "  Cinema 4D, Houdini"),
]

for i, (c, t, b) in enumerate(spaces):
    x = Inches(0.5) + i * Inches(4.28)
    rect(s14, x, Inches(1.22), Inches(4.0), Inches(6.05), CARD)
    rect(s14, x, Inches(1.22), Inches(4.0), Inches(0.65), c)
    txt(s14, t, x+Inches(0.08), Inches(1.27), Inches(3.85), Inches(0.6),
        size=14, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    txt(s14, b, x+Inches(0.14), Inches(1.95), Inches(3.72), Inches(5.2),
        size=10.5, color=LIGHT)

# ╔══════════════════════════════════════════╗
# ║  SLIDE 15 — ANIMATION PROCESS           ║
# ╚══════════════════════════════════════════╝
s15 = new_slide(prs)
header(s15, "09", "ANIMATION PROCESS", ORANGE)
info_bar(s15,
    "All professional animation — from a 10-second ad to a 2-hour film — follows this 4-stage pipeline.",
    ORANGE)

process = [
    (CYAN,   "STEP 1",  "ORGANIZE\nEXECUTION",
     "• Define project goals & scope\n"
     "• Identify team roles\n"
     "• Set milestones & deadlines\n"
     "• Gather reference material\n"
     "• Create style guide\n"
     "• Budget & resource planning\n"
     "• Risk assessment\n\n"
     "OUTPUT:\n"
     "  Project bible, schedule,\n"
     "  mood boards, storyboard"),
    (ORANGE, "STEP 2",  "CHOOSE\nANIMATION TOOL",
     "• Match tool to project type\n"
     "• Consider team skill level\n"
     "• Account for budget\n\n"
     "2D ANIMATION:\n"
     "  Adobe Animate, Toon Boom\n\n"
     "3D ANIMATION:\n"
     "  Maya, Blender, Cinema 4D\n\n"
     "MOTION GRAPHICS:\n"
     "  After Effects, Motion\n\n"
     "GAMES (real-time):\n"
     "  Unity, Unreal Engine"),
    (GREEN,  "STEP 3",  "BUILD & TWEAK\nSEQUENCES",
     "• Create rough blocking first\n"
     "• Refine with keyframes\n"
     "• Add secondary animation\n"
     "• Apply the 12 Principles\n"
     "• Sync with audio / music\n"
     "• Team review cycles\n"
     "• Iterate until approved\n\n"
     "NOTE:\n"
     "  This step is ITERATIVE —\n"
     "  you cycle through it\n"
     "  many times before moving on"),
    (PURPLE, "STEP 4",  "POST-PROCESS\nANIMATION",
     "• Render final frames\n"
     "• Color grading & correction\n"
     "• Add visual effects / overlays\n"
     "• Sound design & music mix\n"
     "• Quality assurance testing\n"
     "• Format for delivery target\n"
     "  (web, broadcast, cinema)\n"
     "• Archive project files\n\n"
     "OUTPUT:\n"
     "  Final deliverable file\n"
     "  (MP4, MOV, WebM, etc.)"),
]

for i, (c, snum, stitle, b) in enumerate(process):
    x = Inches(0.45) + i * Inches(3.22)
    rect(s15, x, Inches(2.05), Inches(3.0), Inches(5.2), CARD)
    rect(s15, x, Inches(2.05), Inches(3.0), Inches(0.45), c)
    txt(s15, snum, x, Inches(2.05), Inches(3.0), Inches(0.45),
        size=13, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    txt(s15, stitle, x, Inches(2.5), Inches(3.0), Inches(0.75),
        size=13, color=c, bold=True, align=PP_ALIGN.CENTER)
    txt(s15, b, x+Inches(0.12), Inches(3.3), Inches(2.76), Inches(3.85),
        size=10.5, color=LIGHT)
    if i < 3:
        txt(s15, "→", Inches(3.35)+i*Inches(3.22), Inches(4.5),
            Inches(0.4), Inches(0.7), size=22, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

rect(s15, Inches(0.45), Inches(7.05), Inches(12.4), Inches(0.38), RGBColor(0, 90, 0))
txt(s15,
    "KEY INSIGHT: Animation is ITERATIVE — Step 3 is repeated multiple times before you ever reach Step 4.",
    Inches(0.6), Inches(7.07), Inches(12.1), Inches(0.34),
    size=12, color=GREEN, bold=True)

# ╔══════════════════════════════════════════╗
# ║  SLIDE 16 — SUMMARY                     ║
# ╚══════════════════════════════════════════╝
s16 = new_slide(prs)
rect(s16, Inches(0), Inches(0), Inches(13.33), Inches(1.28), RGBColor(0,140,190))
txt(s16, "SUMMARY — KEY TAKEAWAYS",
    Inches(0.4), Inches(0.1), Inches(12.5), Inches(1.1),
    size=34, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

left_pts = [
    ("01", CYAN,   "Introduction",    "Animation = sequential images exploiting Persistence of Vision + Phi Phenomenon to create seamless motion."),
    ("02", PURPLE, "Usage",           "Used in art, storytelling, scientific visualization, and instructional training across all industries."),
    ("03", ORANGE, "Techniques",      "3 core techniques: Cel (hand-drawn), Path (trajectory-based), Computer (digital, most versatile)."),
    ("04", GREEN,  "Cel Animation",   "Relies on keyframes + tweening to create smooth motion efficiently. Pencil test validates timing early."),
    ("05", PINK,   "Path Animation",  "Object follows a predefined path — straight line or Bezier curve — with easing for natural feel."),
]
right_pts = [
    ("06", CYAN,   "Computer Anim",   "Digital animation with FK/IK chains and morphing — far faster, more flexible than traditional cel."),
    ("07", PURPLE, "Levels",          "Basic (GIFs/CSS) → Intermediate (motion graphics/2D chars) → Advanced (feature films, real-time VFX)."),
    ("08", ORANGE, "12 Principles",   "Disney's foundational rules — Squash & Stretch, Anticipation, Timing, Ease, Arcs, and 8 more. Still universal."),
    ("09", GREEN,  "Animation Space", "2D (flat), 2½D (pseudo-depth via parallax), 3D (full volumetric space with any camera angle)."),
    ("  ", YELLOW, "Process",         "Organize → Choose Tool → Build & Tweak Sequences → Post-Process. Iterative, not linear."),
]

for i, (n, c, title, desc) in enumerate(left_pts):
    y = Inches(1.45) + i * Inches(1.18)
    rect(s16, Inches(0.35), y, Inches(0.65), Inches(0.95), c)
    txt(s16, n, Inches(0.35), y, Inches(0.65), Inches(0.95),
        size=16, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    rect(s16, Inches(1.05), y, Inches(5.7), Inches(0.95), CARD)
    txt(s16, title, Inches(1.12), y+Inches(0.04), Inches(5.55), Inches(0.4),
        size=12, color=c, bold=True)
    txt(s16, desc, Inches(1.12), y+Inches(0.42), Inches(5.55), Inches(0.55),
        size=10.5, color=LIGHT)

for i, (n, c, title, desc) in enumerate(right_pts):
    y = Inches(1.45) + i * Inches(1.18)
    rect(s16, Inches(7.0), y, Inches(0.65), Inches(0.95), c)
    txt(s16, n, Inches(7.0), y, Inches(0.65), Inches(0.95),
        size=16, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    rect(s16, Inches(7.7), y, Inches(5.25), Inches(0.95), CARD)
    txt(s16, title, Inches(7.78), y+Inches(0.04), Inches(5.1), Inches(0.4),
        size=12, color=c, bold=True)
    txt(s16, desc, Inches(7.78), y+Inches(0.42), Inches(5.1), Inches(0.55),
        size=10.5, color=LIGHT)

rect(s16, Inches(6.6), Inches(1.3), Inches(0.07), Inches(6.0), RGBColor(40,40,90))

# ╔══════════════════════════════════════════╗
# ║  SLIDE 17 — THANK YOU                   ║
# ╚══════════════════════════════════════════╝
s17 = new_slide(prs)
rect(s17, Inches(0), Inches(0),    Inches(13.33), Inches(0.2), CYAN)
rect(s17, Inches(0), Inches(7.3),  Inches(13.33), Inches(0.2), PURPLE)
rect(s17, Inches(1.5), Inches(0.4), Inches(10.33), Inches(6.7), RGBColor(13,13,52))

txt(s17, "THANK YOU",
    Inches(0.8), Inches(0.8), Inches(11.73), Inches(2.5),
    size=88, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

txt(s17, "Animation: Where Imagination Meets Technology",
    Inches(0.8), Inches(3.2), Inches(11.73), Inches(0.75),
    size=21, color=LIGHT, italic=True, align=PP_ALIGN.CENTER)

rect(s17, Inches(3.5), Inches(4.08), Inches(6.33), Inches(0.07), ORANGE)

txt(s17, "Group 3  —  Graphics & Multimedia  —  SMIT",
    Inches(0.8), Inches(4.25), Inches(11.73), Inches(0.55),
    size=15, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

members = [
    "Hashir Junaid           2023F-BCS-358",
    "Taha Haider             2023F-BCS-079",
    "Abdul Rahman Baig       2023F-BCS-100",
]
for i, m in enumerate(members):
    txt(s17, m, Inches(0.8), Inches(4.95)+i*Inches(0.5),
        Inches(11.73), Inches(0.48),
        size=15, color=LIGHT, align=PP_ALIGN.CENTER)

txt(s17, "Topics Covered: Introduction • Usage • Techniques • Types • Computer Animation • Levels • 12 Principles • Space • Process",
    Inches(0.8), Inches(6.7), Inches(11.73), Inches(0.5),
    size=11, color=RGBColor(100,100,160), align=PP_ALIGN.CENTER)

# ── SAVE ─────────────────────────────────────────────────────
out = r"C:\Users\Noman Traders\Desktop\SMIT -AI-CLAUDE\Graphics n Multimedia\Animation_Group3.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
